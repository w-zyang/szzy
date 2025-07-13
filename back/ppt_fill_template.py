#!/usr/bin/env python
"""
PPT模板填充工具
使用JSON元数据智能填充PPT模板
"""

import os
import sys
import json
import logging
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT  # 修正了PP_ALIGN的导入问题
from io import BytesIO
import requests
from image_service import get_image_for_slide

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("ppt_template_fill.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("ppt_template_fill")

class PPTTemplateFiller:
    """PPT模板填充器"""
    
    def __init__(self, template_path, metadata_path=None):
        """
        初始化PPT模板填充器
        
        Args:
            template_path: 模板文件路径
            metadata_path: 元数据文件路径，如果为None则自动查找
        """
        self.template_path = template_path
        self.prs = None
        self.metadata = None
        self.metadata_path = metadata_path
        
        # 加载模板
        self.load_template()
        
        # 加载元数据
        self.load_metadata()
    
    def load_template(self):
        """加载PPT模板"""
        try:
            if not os.path.exists(self.template_path):
                logger.error(f"模板文件不存在: {self.template_path}")
                raise FileNotFoundError(f"模板文件不存在: {self.template_path}")
                
            logger.info(f"加载PPT模板: {self.template_path}")
            self.prs = Presentation(self.template_path)
            logger.info(f"模板加载成功，包含 {len(self.prs.slides)} 张幻灯片")
        except Exception as e:
            logger.error(f"加载模板失败: {str(e)}")
            logger.error(traceback.format_exc())
            raise RuntimeError(f"加载模板失败: {str(e)}")
    
    def load_metadata(self):
        """加载模板元数据"""
        try:
            # 如果未指定元数据路径，则自动查找
            if not self.metadata_path:
                template_name = os.path.basename(self.template_path)
                base_name, _ = os.path.splitext(template_name)
                template_dir = os.path.dirname(self.template_path)
                self.metadata_path = os.path.join(template_dir, f"{base_name}.json")
            
            # 检查元数据文件是否存在
            if not os.path.exists(self.metadata_path):
                logger.warning(f"元数据文件不存在: {self.metadata_path}")
                # 尝试自动生成元数据
                try:
                    # 导入模板分析器
                    sys.path.append(os.path.dirname(os.path.abspath(__file__)))
                    from ppt_template_analyzer import analyze_template
                    
                    # 分析模板
                    self.metadata_path = analyze_template(self.template_path)
                    logger.info(f"自动生成元数据成功: {self.metadata_path}")
                    
                    # 加载新生成的元数据
                    with open(self.metadata_path, 'r', encoding='utf-8') as f:
                        self.metadata = json.load(f)
                    
                    logger.info(f"元数据加载成功，包含 {len(self.metadata['slides'])} 张幻灯片信息")
                    return
                except Exception as e:
                    logger.error(f"自动生成元数据失败: {str(e)}")
                    logger.error(traceback.format_exc())
                    logger.warning("将使用基本填充方法")
                    return
            
            # 加载元数据
            logger.info(f"加载元数据: {self.metadata_path}")
            with open(self.metadata_path, 'r', encoding='utf-8') as f:
                self.metadata = json.load(f)
            
            logger.info(f"元数据加载成功，包含 {len(self.metadata['slides'])} 张幻灯片信息")
        except Exception as e:
            logger.error(f"加载元数据失败: {str(e)}")
            logger.error(traceback.format_exc())
            logger.warning("将使用基本填充方法")
    
    def find_best_slide_for_content(self, content_type, slide_index=None):
        """
        根据内容类型找到最合适的幻灯片
        
        Args:
            content_type: 内容类型（cover, content, keypoints, image, table, chart, summary）
            slide_index: 指定的幻灯片索引，如果为None则自动查找
            
        Returns:
            最合适的幻灯片索引
        """
        if not self.metadata:
            # 如果没有元数据，使用基本规则
            logger.info(f"没有元数据，使用基本规则选择幻灯片")
            
            # 检查self.prs是否为None
            if self.prs is None or not hasattr(self.prs, 'slides') or len(self.prs.slides) == 0:
                logger.warning("模板对象为空或无幻灯片，返回默认索引0")
                return 0
                
            if content_type == 'cover':
                return 0
            elif content_type == 'summary':
                return len(self.prs.slides) - 1
            elif slide_index is not None and 0 <= slide_index < len(self.prs.slides):
                return slide_index
            else:
                return min(1, len(self.prs.slides) - 1)
        
        # 检查self.prs是否为None
        if self.prs is None or not hasattr(self.prs, 'slides') or len(self.prs.slides) == 0:
            logger.warning("模板对象为空或无幻灯片，返回默认索引0")
            return 0
            
        # 如果指定了幻灯片索引且在有效范围内，直接返回
        if slide_index is not None and 0 <= slide_index < len(self.prs.slides):
            logger.info(f"使用指定的幻灯片索引: {slide_index}")
            return slide_index
        
        # 根据内容类型查找合适的幻灯片
        suitable_slides = []
        for i, slide_info in enumerate(self.metadata['slides']):
            if 'suitable_for' in slide_info and content_type in slide_info['suitable_for']:
                suitable_slides.append(i)
        
        # 如果找到合适的幻灯片，返回第一个
        if suitable_slides:
            logger.info(f"找到{len(suitable_slides)}个适合{content_type}的幻灯片: {suitable_slides}")
            return suitable_slides[0]
        
        # 如果没有找到，使用基本规则
        logger.info(f"未找到适合{content_type}的幻灯片，使用基本规则")
        if content_type == 'cover':
            return 0
        elif content_type == 'summary':
            return len(self.prs.slides) - 1
        else:
            return min(1, len(self.prs.slides) - 1)
    
    def fill_slide(self, slide, content):
        """
        填充幻灯片内容，保留原始设计和背景，智能调整模板
        
        Args:
            slide: 幻灯片对象
            content: 内容数据
            
        Returns:
            slide: 填充后的幻灯片对象
        """
        try:
            # 获取幻灯片尺寸信息
            slide_width = getattr(slide, 'slide_width', 9144000)  # 默认宽度值
            slide_height = getattr(slide, 'slide_height', 5143500)  # 默认高度值
            
            # 记录所有形状，以便后续引用
            all_shapes = list(slide.shapes)
            text_shapes = [shape for shape in all_shapes if shape.has_text_frame]
            placeholder_shapes = [shape for shape in all_shapes if hasattr(shape, 'is_placeholder') and shape.is_placeholder]
            picture_shapes = [shape for shape in all_shapes if hasattr(shape, 'image') or (hasattr(shape, 'shape_type') and shape.shape_type == 13)]  # 13 = PICTURE
            
            logger.info(f"幻灯片包含 {len(all_shapes)} 个形状，其中 {len(text_shapes)} 个文本框，{len(placeholder_shapes)} 个占位符，{len(picture_shapes)} 个图片")
            
            # 获取主题关键词，用于替换图片和相关性判断
            topic_keywords = self._extract_topic_keywords(content)
            logger.info(f"提取的主题关键词: {topic_keywords}")
            
            # 分析文本框的位置和大小，推断其用途
            title_shapes = []
            subtitle_shapes = []
            content_shapes = []
            other_shapes = []
            
            # 分析文本框的位置和大小
            for shape in text_shapes:
                # 获取形状的位置和大小
                top = getattr(shape, 'top', 0)
                left = getattr(shape, 'left', 0)
                width = getattr(shape, 'width', 0)
                height = getattr(shape, 'height', 0)
                
                # 计算相对位置（百分比）
                rel_top = top / slide_height
                rel_left = left / slide_width
                rel_width = width / slide_width
                rel_height = height / slide_height
                
                # 判断文本框的用途
                is_title = False
                is_subtitle = False
                is_content = False
                
                # 检查是否是标准占位符
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    if hasattr(shape, 'placeholder_format') and hasattr(shape.placeholder_format, 'type'):
                        if shape.placeholder_format.type == 1:  # 1 = TITLE
                            is_title = True
                        elif shape.placeholder_format.type == 2:  # 2 = SUBTITLE/BODY
                            if rel_top < 0.3:  # 在上部，可能是副标题
                                is_subtitle = True
                            else:  # 在下部，可能是正文
                                is_content = True
                        elif shape.placeholder_format.type in [3, 4, 5, 6, 7]:  # 其他类型的文本占位符
                            is_content = True
                
                # 如果不是标准占位符，根据位置和大小推断
                if not (is_title or is_subtitle or is_content):
                    # 标题通常在顶部，宽度大，高度小
                    if rel_top < 0.2 and rel_width > 0.5:
                        is_title = True
                    # 副标题通常在标题下方，宽度大，高度小
                    elif rel_top < 0.3 and rel_width > 0.4:
                        is_subtitle = True
                    # 内容通常占据幻灯片的主要部分
                    elif rel_width > 0.4 and rel_height > 0.3:
                        is_content = True
                    # 其他情况，可能是注释、页脚等
                    else:
                        pass
                
                # 根据判断结果分类
                if is_title:
                    title_shapes.append(shape)
                elif is_subtitle:
                    subtitle_shapes.append(shape)
                elif is_content:
                    content_shapes.append(shape)
                else:
                    other_shapes.append(shape)
            
            # 按面积排序（从大到小）
            title_shapes.sort(key=lambda s: getattr(s, 'width', 0) * getattr(s, 'height', 0), reverse=True)
            subtitle_shapes.sort(key=lambda s: getattr(s, 'width', 0) * getattr(s, 'height', 0), reverse=True)
            content_shapes.sort(key=lambda s: getattr(s, 'width', 0) * getattr(s, 'height', 0), reverse=True)
            
            # 填充标题
            if 'title' in content and content['title']:
                if title_shapes:
                    # 使用最大的标题形状
                    title_shape = title_shapes[0]
                    title_shape.text_frame.text = content['title']
                    logger.info(f"填充标题: {content['title']}")
                elif text_shapes:
                    # 如果没有找到标题形状，使用第一个文本形状
                    title_shape = text_shapes[0]
                    title_shape.text_frame.text = content['title']
                    logger.info(f"使用第一个文本框填充标题: {content['title']}")
            
            # 填充副标题
            if 'subtitle' in content and content['subtitle']:
                if subtitle_shapes:
                    # 使用最大的副标题形状
                    subtitle_shape = subtitle_shapes[0]
                    subtitle_shape.text_frame.text = content['subtitle']
                    logger.info(f"填充副标题: {content['subtitle']}")
                elif len(text_shapes) > 1 and not title_shapes:
                    # 如果没有找到副标题形状，但有多个文本形状，使用第二个
                    subtitle_shape = text_shapes[1]
                    subtitle_shape.text_frame.text = content['subtitle']
                    logger.info(f"使用第二个文本框填充副标题: {content['subtitle']}")
            
            # 填充主要内容
            if 'content' in content and content['content']:
                if content_shapes:
                    # 使用最大的内容形状
                    content_shape = content_shapes[0]
                    content_shape.text_frame.text = content['content']
                    logger.info(f"填充内容: {content['content'][:30]}...")
                elif len(text_shapes) > len(title_shapes) + len(subtitle_shapes):
                    # 如果没有找到内容形状，但有额外的文本形状，使用剩余最大的
                    used_shape_ids = [id(s) for s in title_shapes + subtitle_shapes]
                    available_shapes = [s for s in text_shapes if id(s) not in used_shape_ids]
                    if available_shapes:
                        # 按面积排序，选择最大的
                        available_shapes.sort(key=lambda s: getattr(s, 'width', 0) * getattr(s, 'height', 0), reverse=True)
                        content_shape = available_shapes[0]
                        content_shape.text_frame.text = content['content']
                        logger.info(f"使用额外文本框填充内容: {content['content'][:30]}...")
            
            # 填充关键点
            if 'keypoints' in content and content['keypoints']:
                # 尝试找到合适的形状填充关键点
                keypoints_shape = None
                
                # 如果有内容形状但没有填充内容，可以用来填充关键点
                if content_shapes and 'content' not in content:
                    keypoints_shape = content_shapes[0]
                # 或者使用额外的内容形状
                elif len(content_shapes) > 1:
                    keypoints_shape = content_shapes[1]
                # 或者使用其他未使用的文本形状
                else:
                    used_shape_ids = []
                    if 'title' in content and title_shapes:
                        used_shape_ids.append(id(title_shapes[0]))
                    if 'subtitle' in content and subtitle_shapes:
                        used_shape_ids.append(id(subtitle_shapes[0]))
                    if 'content' in content and content_shapes:
                        used_shape_ids.append(id(content_shapes[0]))
                    
                    available_shapes = [s for s in text_shapes if id(s) not in used_shape_ids]
                    if available_shapes:
                        # 按面积排序，选择最大的
                        available_shapes.sort(key=lambda s: getattr(s, 'width', 0) * getattr(s, 'height', 0), reverse=True)
                        keypoints_shape = available_shapes[0]
                
                # 如果找到了形状，填充要点
                if keypoints_shape:
                    # 清除现有内容
                    text_frame = keypoints_shape.text_frame
                    text_frame.clear()
                    
                    # 添加要点
                    for i, point in enumerate(content['keypoints']):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        # 处理不同类型的要点
                        if isinstance(point, dict):
                            p.text = point.get('text', '')
                            p.level = point.get('level', 0)
                        elif isinstance(point, list):
                            p.text = ' - '.join([str(item) for item in point])
                        else:
                            p.text = str(point)
                        
                        # 设置项目符号
                        p.bullet = True
                    
                    logger.info(f"填充了 {len(content['keypoints'])} 个要点")
            
            # 处理装饰性图片和内容图片
            # 1. 首先处理用户指定的图片
            user_image_data = None
            if 'image' in content and content['image']:
                user_image_data = get_image_for_slide(content)
            
            # 2. 获取主题相关图片（如果没有用户指定图片）
            topic_image_data = None
            if not user_image_data and topic_keywords:
                # 尝试获取主题相关图片
                try:
                    from image_service import ImageService
                    image_service = ImageService()
                    
                    # 构建查询
                    query = " ".join(topic_keywords[:3])  # 使用前3个关键词
                    image_urls = image_service.search_image(query, max_results=3)
                    
                    if image_urls and len(image_urls) > 0:
                        # 尝试下载多个图片，如果第一个失败可以使用备选
                        for url in image_urls:
                            try:
                                import requests
                                response = requests.get(url, timeout=10)
                                if response.status_code == 200:
                                    topic_image_data = response.content
                                    logger.info(f"获取到主题相关图片: {url}")
                                    break
                            except Exception as e:
                                logger.error(f"下载图片失败: {url}, {str(e)}")
                                continue
                except Exception as e:
                    logger.error(f"获取主题相关图片失败: {str(e)}")
            
            # 3. 如果没有图片，使用默认图片
            if not user_image_data and not topic_image_data:
                try:
                    # 根据主题或内容选择合适的默认图片
                    default_image_path = None
                    subject_lower = content.get('subject', '').lower()
                    title_lower = content.get('title', '').lower() if content.get('title') else ''
                    content_lower = str(content.get('content', '')).lower() if content.get('content') else ''
                    
                    # 尝试基于学科和内容匹配默认图片
                    if '生物' in subject_lower or '细胞' in title_lower or '细胞' in content_lower:
                        if '植物' in title_lower or '植物' in content_lower:
                            default_image_path = os.path.join('default_images', 'plant_cell.jpg')
                        else:
                            default_image_path = os.path.join('default_images', 'animal_cell.jpg')
                    elif '生物' in subject_lower:
                        default_image_path = os.path.join('default_images', 'biology.jpg')
                    elif '总结' in title_lower or '结论' in title_lower:
                        default_image_path = os.path.join('default_images', 'conclusion.jpg')
                    elif '封面' in title_lower or '标题' in title_lower or content.get('type') == 'cover':
                        default_image_path = os.path.join('default_images', 'cover.jpg')
                    else:
                        default_image_path = os.path.join('default_images', 'default.jpg')
                    
                    # 加载默认图片
                    if default_image_path and os.path.exists(default_image_path):
                        with open(default_image_path, 'rb') as f:
                            topic_image_data = f.read()
                            logger.info(f"使用默认图片: {default_image_path}")
                except Exception as e:
                    logger.error(f"加载默认图片失败: {str(e)}")
            
            # 4. 处理图片占位符和装饰图片
            image_data = user_image_data or topic_image_data
            if image_data:
                # 查找图片占位符
                image_placeholder = None
                for shape in placeholder_shapes:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 18:  # 18 = PICTURE
                        image_placeholder = shape
                        break
                
                # 如果找到图片占位符，填充图片
                if image_placeholder:
                    # 记录原始位置和大小
                    left, top, width, height = image_placeholder.left, image_placeholder.top, image_placeholder.width, image_placeholder.height
                    
                    # 添加图片到相同位置
                    try:
                        slide.shapes.add_picture(BytesIO(image_data), left, top, width, height)
                        logger.info(f"添加图片到占位符位置")
                    except Exception as e:
                        logger.error(f"添加图片到占位符失败: {str(e)}")
                
                # 如果有装饰性图片，替换它们
                decorative_images = []
                for shape in picture_shapes:
                    # 获取形状的位置和大小
                    left = getattr(shape, 'left', 0)
                    top = getattr(shape, 'top', 0)
                    width = getattr(shape, 'width', 0)
                    height = getattr(shape, 'height', 0)
                    area = width * height
                    
                    # 计算相对位置（百分比）
                    rel_top = top / slide_height
                    rel_left = left / slide_width
                    rel_width = width / slide_width
                    rel_height = height / slide_height
                    
                    # 判断是否是装饰性图片
                    is_decorative = (rel_width < 0.3 or rel_height < 0.3)  # 小图片
                    is_decorative = is_decorative or (rel_left < 0.1 or rel_left > 0.9 or rel_top < 0.1 or rel_top > 0.9)  # 边缘图片
                    
                    if is_decorative:
                        decorative_images.append((shape, left, top, width, height, area))
                
                # 如果找到装饰性图片，替换第一个
                if decorative_images:
                    # 按面积排序，从小到大
                    decorative_images.sort(key=lambda x: x[5])
                    
                    # 替换最小的装饰性图片
                    shape, left, top, width, height, _ = decorative_images[0]
                    try:
                        # 添加新图片到相同位置
                        slide.shapes.add_picture(BytesIO(image_data), left, top, width, height)
                        logger.info(f"替换装饰性图片")
                    except Exception as e:
                        logger.error(f"替换装饰性图片失败: {str(e)}")
                
                # 如果没有图片占位符和装饰性图片，但需要添加图片，在合适位置添加
                if not image_placeholder and not decorative_images and (content.get('layout', '').lower() == 'image' or 'image' in content):
                    try:
                        # 计算合适的位置和大小
                        left = int(slide_width * 0.5)  # 右半部分
                        top = int(slide_height * 0.3)
                        width = int(slide_width * 0.4)
                        height = int(slide_height * 0.4)
                        
                        # 添加图片
                        slide.shapes.add_picture(BytesIO(image_data), left, top, width, height)
                        logger.info(f"添加图片到自定义位置")
                    except Exception as e:
                        logger.error(f"添加图片到自定义位置失败: {str(e)}")
            
            # 填充表格
            if 'table' in content and content['table']:
                table_data = content['table']
                
                # 查找表格占位符
                table_shape = None
                for shape in placeholder_shapes:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 12:  # 12 = TABLE
                        table_shape = shape
                        break
                
                # 如果找到表格占位符，填充表格
                if table_shape:
                    # 记录原始位置和大小
                    left, top, width, height = table_shape.left, table_shape.top, table_shape.width, table_shape.height
                    
                    # 创建表格
                    rows = len(table_data)
                    cols = max(len(row) for row in table_data) if rows > 0 else 0
                    
                    if rows > 0 and cols > 0:
                        try:
                            # 添加表格
                            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                            
                            # 填充表格数据
                            for i, row_data in enumerate(table_data):
                                for j, cell_data in enumerate(row_data):
                                    if j < cols:
                                        cell = table.cell(i, j)
                                        cell.text = str(cell_data)
                            
                            logger.info(f"添加了 {rows}x{cols} 的表格")
                        except Exception as e:
                            logger.error(f"添加表格失败: {str(e)}")
                else:
                    # 如果没有找到表格占位符，在合适的位置添加表格
                    try:
                        # 计算合适的位置和大小
                        left = int(slide_width * 0.1)
                        top = int(slide_height * 0.4)
                        width = int(slide_width * 0.8)
                        height = int(slide_height * 0.4)
                        
                        # 创建表格
                        rows = len(table_data)
                        cols = max(len(row) for row in table_data) if rows > 0 else 0
                        
                        if rows > 0 and cols > 0:
                            # 添加表格
                            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                            
                            # 填充表格数据
                            for i, row_data in enumerate(table_data):
                                for j, cell_data in enumerate(row_data):
                                    if j < cols:
                                        cell = table.cell(i, j)
                                        cell.text = str(cell_data)
                            
                            logger.info(f"添加了 {rows}x{cols} 的表格到自定义位置")
                    except Exception as e:
                        logger.error(f"添加表格到自定义位置失败: {str(e)}")
            
            # 处理未使用的文本框（可以添加页码、注释等）
            used_shape_ids = []
            if 'title' in content and title_shapes:
                used_shape_ids.extend([id(s) for s in title_shapes])
            if 'subtitle' in content and subtitle_shapes:
                used_shape_ids.extend([id(s) for s in subtitle_shapes])
            if 'content' in content and content_shapes:
                used_shape_ids.append(id(content_shapes[0])) # 只标记第一个内容形状为已使用
            
            # 查找未使用的文本框
            unused_text_shapes = [s for s in text_shapes if id(s) not in used_shape_ids]
            
                            # 如果有页码或页脚信息，可以添加到未使用的文本框
            if 'footer' in content and content['footer'] and unused_text_shapes:
                # 找到最底部的文本框作为页脚
                footer_shape = min(unused_text_shapes, key=lambda s: -s.top)
                footer_shape.text_frame.text = content['footer']
                logger.info(f"添加页脚: {content['footer']}")
                unused_text_shapes.remove(footer_shape)
            
            # 最后一步：优化文本布局，处理重复内容和排版问题
            self._optimize_text_layout(slide, content)
            
            return slide
        except Exception as e:
            logger.error(f"填充幻灯片失败: {str(e)}")
            logger.error(traceback.format_exc())
            return slide
            
    def _extract_topic_keywords(self, content):
        """
        从内容中提取主题关键词
        
        Args:
            content: 内容数据
            
        Returns:
            关键词列表
        """
        keywords = []
        
        # 从标题中提取关键词
        if 'title' in content and content['title']:
            title_words = content['title'].split()
            keywords.extend([word for word in title_words if len(word) > 1])
        
        # 从内容中提取关键词
        if 'content' in content and content['content']:
            # 提取前10个词作为关键词
            content_words = content['content'].split()[:10]
            keywords.extend([word for word in content_words if len(word) > 1 and word not in keywords])
        
        # 从关键点中提取关键词
        if 'keypoints' in content and content['keypoints']:
            for point in content['keypoints'][:3]:  # 只使用前3个关键点
                if isinstance(point, str):
                    point_words = point.split()[:3]  # 每个关键点取前3个词
                    keywords.extend([word for word in point_words if len(word) > 1 and word not in keywords])
        
        # 如果有学科信息，添加学科关键词
        if 'subject' in content:
            keywords.append(content['subject'])
        
        return keywords
            
    def _replace_placeholder_texts(self, slide, content, topic_keywords):
        """
        替换幻灯片中的占位符文本，删除无关文本和重复内容
        
        Args:
            slide: 幻灯片对象
            content: 内容数据
            topic_keywords: 主题关键词
        """
        try:
            if slide is None:
                logger.warning("slide对象为None，无法处理占位符文本")
                return
                
            # 常见的占位符文本模式
            title_placeholders = ["点击此处添加标题", "添加标题", "Click to add title", "Title", "请输入您的标题", "您的标题", 
                                "YOUR TITLE", "CLICK TO EDIT", "CLICK TO ADD TITLE", "TITLE HERE", "ADD TITLE"]
            content_placeholders = ["点击添加正文", "添加正文", "点击此处添加正文", "Click to add text", "Content", "请输入内容", "您的内容",
                                  "YOUR CONTENT", "CLICK TO ADD TEXT", "ADD CONTENT", "TEXT HERE"]
            subtitle_placeholders = ["点击此处添加副标题", "添加副标题", "Click to add subtitle", "Subtitle", "请输入您的副标题", "您的副标题",
                                   "YOUR SUBTITLE", "SUBTITLE HERE", "ADD SUBTITLE"]
            footer_placeholders = ["页脚", "Footer", "FOOTER", "PAGE", "页码"]
            
            # 通用的占位符关键词(不区分大小写)
            generic_placeholders = ["请输入", "您的", "点击", "添加", "输入", "Click", "Enter", "Your", "Type", "Add", "Edit", "Text"]
            
            # 需要完全移除的文本模式(这些常见于模板中但在最终PPT中应该被删除)
            remove_patterns = ["YOUR LOGO", "LOGO HERE", "Company Name", "Company Logo", "Department", "Date", "Author Name", 
                             "Your Name", "ABCDE", "abcde", "Lorem ipsum", "lorem ipsum", "Example", "EXAMPLE", 
                             "Sample", "SAMPLE", "placeholder", "PLACEHOLDER", "Click to edit"]
            
            # 获取幻灯片尺寸信息
            slide_height = getattr(slide, 'slide_height', 7200000)  # 默认值，避免属性不存在
            
            # 跟踪已填充的内容，避免重复
            filled_content = set()
            filled_title = False
            filled_subtitle = False
            
            # 收集所有文本形状，以便后续处理
            text_shapes = []
            title_shapes = []
            subtitle_shapes = []
            content_shapes = []
            other_shapes = []
            
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                    
                text_frame = shape.text_frame
                if not text_frame.text:
                    continue
                
                original_text = text_frame.text.strip()
                shape_top = getattr(shape, 'top', 0)
                shape_left = getattr(shape, 'left', 0)
                width = getattr(shape, 'width', 0)
                height = getattr(shape, 'height', 0)
                
                # 计算相对位置（百分比）
                slide_width = getattr(slide, 'slide_width', 9144000)  # 默认宽度值
                slide_height = getattr(slide, 'slide_height', 5143500)  # 默认高度值
                rel_top = shape_top / slide_height
                rel_left = shape_left / slide_width
                rel_width = width / slide_width
                rel_height = height / slide_height
                
                # 判断文本框的用途
                is_title = False
                is_subtitle = False
                is_content = False
                
                # 检查是否是标准占位符
                if hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    if hasattr(shape, 'placeholder_format') and hasattr(shape.placeholder_format, 'type'):
                        if shape.placeholder_format.type == 1:  # 1 = TITLE
                            is_title = True
                        elif shape.placeholder_format.type == 2:  # 2 = SUBTITLE/BODY
                            if rel_top < 0.3:  # 在上部，可能是副标题
                                is_subtitle = True
                            else:  # 在下部，可能是正文
                                is_content = True
                        elif shape.placeholder_format.type in [3, 4, 5, 6, 7]:  # 其他类型的文本占位符
                            is_content = True
                
                # 如果不是标准占位符，根据位置和大小推断
                if not (is_title or is_subtitle or is_content):
                    # 标题通常在顶部，宽度大，高度小
                    if rel_top < 0.2 and rel_width > 0.5:
                        is_title = True
                    # 副标题通常在标题下方，宽度大，高度小
                    elif rel_top < 0.3 and rel_width > 0.4:
                        is_subtitle = True
                    # 内容通常占据幻灯片的主要部分
                    elif rel_width > 0.4 and rel_height > 0.3:
                        is_content = True
                    # 其他情况，可能是注释、页脚等
                    else:
                        pass
                
                # 根据判断结果分类
                if is_title:
                    title_shapes.append(shape)
                elif is_subtitle:
                    subtitle_shapes.append(shape)
                elif is_content:
                    content_shapes.append(shape)
                else:
                    other_shapes.append(shape)
                
                # 添加到所有文本形状列表
                text_shapes.append(shape)
                
            # 按面积排序（从大到小）
            title_shapes.sort(key=lambda s: getattr(s, 'width', 0) * getattr(s, 'height', 0), reverse=True)
            subtitle_shapes.sort(key=lambda s: getattr(s, 'width', 0) * getattr(s, 'height', 0), reverse=True)
            content_shapes.sort(key=lambda s: getattr(s, 'width', 0) * getattr(s, 'height', 0), reverse=True)
            
            # 不再需要按垂直位置排序，因为我们已经按类型分类了文本框
            
            # 填充内容，使用之前分类好的文本框
            
            # 删除不需要的部分，我们已经重写了文本框处理逻辑
            
            # 填充标题
            if 'title' in content and content['title']:
                if title_shapes:
                    # 使用最大的标题形状
                    title_shape = title_shapes[0]
                    title_shape.text_frame.text = content['title']
                    logger.info(f"填充标题: {content['title']}")
                elif text_shapes:
                    # 如果没有找到标题形状，使用第一个文本形状
                    title_shape = text_shapes[0]
                    title_shape.text_frame.text = content['title']
                    logger.info(f"使用第一个文本框填充标题: {content['title']}")
            
            # 填充副标题
            if 'subtitle' in content and content['subtitle']:
                if subtitle_shapes:
                    # 使用最大的副标题形状
                    subtitle_shape = subtitle_shapes[0]
                    subtitle_shape.text_frame.text = content['subtitle']
                    logger.info(f"填充副标题: {content['subtitle']}")
                elif len(text_shapes) > 1 and not title_shapes:
                    # 如果没有找到副标题形状，但有多个文本形状，使用第二个
                    subtitle_shape = text_shapes[1]
                    subtitle_shape.text_frame.text = content['subtitle']
                    logger.info(f"使用第二个文本框填充副标题: {content['subtitle']}")
            
            # 填充主要内容
            if 'content' in content and content['content']:
                if content_shapes:
                    # 使用最大的内容形状
                    content_shape = content_shapes[0]
                    content_shape.text_frame.text = content['content']
                    logger.info(f"填充内容: {content['content'][:30]}...")
                elif len(text_shapes) > len(title_shapes) + len(subtitle_shapes):
                    # 如果没有找到内容形状，但有额外的文本形状，使用剩余最大的
                    used_shape_ids = [id(s) for s in title_shapes + subtitle_shapes]
                    available_shapes = [s for s in text_shapes if id(s) not in used_shape_ids]
                    if available_shapes:
                        # 按面积排序，选择最大的
                        available_shapes.sort(key=lambda s: getattr(s, 'width', 0) * getattr(s, 'height', 0), reverse=True)
                        content_shape = available_shapes[0]
                        content_shape.text_frame.text = content['content']
                        logger.info(f"使用额外文本框填充内容: {content['content'][:30]}...")
                
        except Exception as e:
            logger.error(f"替换占位符文本失败: {str(e)}")
            logger.error(traceback.format_exc())
            
    def _text_similarity(self, text1, text2):
        """
        计算两段文本的相似度
        
        Args:
            text1: 第一段文本
            text2: 第二段文本
            
        Returns:
            相似度分数，范围0-1
        """
        # 简单实现，基于共同单词的比例
        if not text1 or not text2:
            return 0
            
        # 转换为小写并分词
        words1 = set(text1.lower().split())
        words2 = set(text2.lower().split())
        
        # 计算交集大小
        common_words = words1.intersection(words2)
        
        # 计算相似度分数
        similarity = len(common_words) / max(len(words1), len(words2))
        
        return similarity
    
    def _process_images(self, slide, content, topic_keywords):
        """
        处理幻灯片中的图片，替换不相关的装饰性图片为主题相关图片
        
        Args:
            slide: 幻灯片对象
            content: 内容数据
            topic_keywords: 主题关键词
        """
        try:
            if slide is None:
                logger.warning("slide对象为None，无法处理图片")
                return
                
            # 获取幻灯片尺寸信息，避免属性不存在
            slide_width = getattr(slide, 'slide_width', 9144000)  # 默认宽度值
            slide_height = getattr(slide, 'slide_height', 5143500)  # 默认高度值
                
            # 如果用户指定了图片，优先使用用户提供的图片
            if 'image' in content and content['image']:
                user_image_data = get_image_for_slide(content)
                if user_image_data:
                    # 查找图片占位符或位置适合的现有图片
                    image_shapes = []
                    for shape in slide.shapes:
                        # 检查是否是图片
                        is_picture = (hasattr(shape, 'image') or 
                                    (hasattr(shape, 'shape_type') and shape.shape_type == 13))  # 13 = PICTURE
                        if is_picture:
                            image_shapes.append((shape, shape.left, shape.top, shape.width, shape.height))
                    
                    if image_shapes:
                        # 优先使用中央或较大的图片位置
                        slide_center_x = slide_width / 2
                        slide_center_y = slide_height / 2
                        
                        # 计算每个图片到中心的距离以及图片大小
                        for i, (shape, left, top, width, height) in enumerate(image_shapes):
                            center_x = left + width/2
                            center_y = top + height/2
                            distance = ((center_x - slide_center_x) ** 2 + (center_y - slide_center_y) ** 2) ** 0.5
                            area = width * height
                            image_shapes[i] = (shape, left, top, width, height, distance, area)
                        
                        # 根据距离和大小排序，优先选择距离近且较大的图片
                        image_shapes.sort(key=lambda x: (x[5], -x[6]))
                        
                        # 使用最合适的位置
                        best_shape = image_shapes[0]
                        shape, left, top, width, height = best_shape[:5]
                        
                        # 添加新图片
                        slide.shapes.add_picture(BytesIO(user_image_data), left, top, width, height)
                        logger.info(f"替换图片为用户提供的图片")
                        return
                    else:
                        # 如果没有找到合适的位置，添加到默认位置
                        left = int(slide_width * 0.5)
                        top = int(slide_height * 0.3)
                        width = int(slide_width * 0.4)
                        height = int(slide_height * 0.4)
                        
                        slide.shapes.add_picture(BytesIO(user_image_data), left, top, width, height)
                        logger.info(f"添加用户提供的图片到默认位置")
                        return
            
            # 如果没有用户提供的图片，尝试获取主题相关图片
            if not topic_keywords:
                logger.info("没有主题关键词，跳过图片处理")
                return
                
            # 获取主题相关的图片数据
            topic_image_data = None
            
            # 构建优化的搜索查询
            # 从主题和内容标题中提取关键词
            query_keywords = []
            
            # 优先使用主题的关键词
            if topic_keywords:
                # 过滤掉长度小于2的关键词
                filtered_keywords = [k for k in topic_keywords if len(k) >= 2]
                # 优先使用前3个关键词
                query_keywords.extend(filtered_keywords[:3])
            
            # 如果关键词不足3个，从标题和内容中补充
            if len(query_keywords) < 3:
                if 'title' in content and content['title']:
                    title_words = content['title'].split()
                    title_keywords = [w for w in title_words if len(w) >= 3 and w not in query_keywords]
                    query_keywords.extend(title_keywords[:2])
            
            if len(query_keywords) < 3:
                if 'subject' in content and content['subject']:
                    query_keywords.append(content['subject'])
            
            # 构建最终查询
            if query_keywords:
                query = " ".join(query_keywords)
                
                # 尝试获取主题相关图片
                try:
                    from image_service import ImageService
                    image_service = ImageService()
                    image_urls = image_service.search_image(query, max_results=3)
                    
                    if image_urls and len(image_urls) > 0:
                        # 尝试下载多个图片，如果第一个失败可以使用备选
                        for url in image_urls:
                            try:
                                import requests
                                response = requests.get(url, timeout=10)
                                if response.status_code == 200:
                                    topic_image_data = response.content
                                    logger.info(f"获取到主题相关图片: {url}")
                                    break
                            except Exception as e:
                                logger.error(f"下载图片失败: {url}, {str(e)}")
                                continue
                except Exception as e:
                    logger.error(f"获取主题相关图片失败: {str(e)}")
            
            # 如果没有成功获取主题相关图片，尝试使用默认图库
            if not topic_image_data:
                try:
                    # 根据主题或内容选择合适的默认图片
                    default_image_path = None
                    subject_lower = content.get('subject', '').lower()
                    title_lower = content.get('title', '').lower()
                    content_lower = str(content.get('content', '')).lower()
                    
                    # 尝试基于学科和内容匹配默认图片
                    if '生物' in subject_lower or '细胞' in title_lower or '细胞' in content_lower:
                        if '植物' in title_lower or '植物' in content_lower:
                            default_image_path = os.path.join('default_images', 'plant_cell.jpg')
                        else:
                            default_image_path = os.path.join('default_images', 'animal_cell.jpg')
                    elif '生物' in subject_lower:
                        default_image_path = os.path.join('default_images', 'biology.jpg')
                    elif '总结' in title_lower or '结论' in title_lower:
                        default_image_path = os.path.join('default_images', 'conclusion.jpg')
                    elif '封面' in title_lower or '标题' in title_lower or content.get('type') == 'cover':
                        default_image_path = os.path.join('default_images', 'cover.jpg')
                    else:
                        default_image_path = os.path.join('default_images', 'default.jpg')
                    
                    # 加载默认图片
                    if default_image_path and os.path.exists(default_image_path):
                        with open(default_image_path, 'rb') as f:
                            topic_image_data = f.read()
                            logger.info(f"使用默认图片: {default_image_path}")
                except Exception as e:
                    logger.error(f"加载默认图片失败: {str(e)}")
            
            # 如果仍然没有图片，直接返回
            if not topic_image_data:
                logger.warning("无法获取合适的主题图片，跳过图片替换")
                return
            
            # 识别和分类幻灯片中的图片
            decorative_images = []  # 装饰性图片
            content_images = []     # 内容相关图片
            
            for shape in slide.shapes:
                # 检查是否是图片
                is_picture = (hasattr(shape, 'image') or 
                             (hasattr(shape, 'shape_type') and shape.shape_type == 13))  # 13 = PICTURE
                
                if not is_picture:
                    continue
                    
                # 记录图片位置和大小
                left = getattr(shape, 'left', 0)
                top = getattr(shape, 'top', 0)
                width = getattr(shape, 'width', 0) 
                height = getattr(shape, 'height', 0)
                area = width * height
                
                # 尝试判断图片是否为装饰性的
                # 1. 位于边缘的小图片更可能是装饰性的
                is_at_edge = (left < slide_width * 0.1 or 
                            left + width > slide_width * 0.9 or
                            top < slide_height * 0.1 or
                            top + height > slide_height * 0.9)
                            
                # 2. 面积较小的图片更可能是装饰性的
                is_small = area < (slide_width * slide_height * 0.1)
                
                # 3. 宽高比异常的图片可能是装饰性的(如Logo、图标等)
                aspect_ratio = width / max(1, height)
                has_unusual_ratio = aspect_ratio < 0.3 or aspect_ratio > 3
                
                # 根据特征判断图片类型
                if (is_at_edge and is_small) or (is_small and has_unusual_ratio):
                    # 可能是装饰性图片
                    decorative_images.append((shape, left, top, width, height, area))
                else:
                    # 可能是内容相关图片
                    content_images.append((shape, left, top, width, height, area))
            
            # 决定替换哪些图片
            images_to_replace = []
            
            # 优先替换装饰性图片
            if decorative_images:
                # 按面积从小到大排序
                decorative_images.sort(key=lambda x: x[5])
                # 添加装饰性图片，从最小的开始
                images_to_replace.extend(decorative_images)
            
            # 如果没有装饰性图片，考虑替换内容图片
            # 注意：仅当内容图片很可能是模板中的示例图片时才替换
            elif content_images and len(content_images) == 1:
                images_to_replace.extend(content_images)
            
            # 如果找到了需要替换的图片，替换第一个
            if images_to_replace:
                shape, left, top, width, height, _ = images_to_replace[0]
                
                try:
                    # 添加新图片到相同位置
                    slide.shapes.add_picture(BytesIO(topic_image_data), left, top, width, height)
                    logger.info(f"替换图片为主题相关图片")
                except Exception as e:
                    logger.error(f"替换图片失败: {str(e)}")
                    logger.error(traceback.format_exc())
            else:
                logger.info("未找到合适的图片位置进行替换")
        except Exception as e:
            logger.error(f"处理图片失败: {str(e)}")
            logger.error(traceback.format_exc())
    
    def _optimize_text_layout(self, slide, content):
        """
        优化幻灯片中文本框的排版，处理重叠文本和排版问题
        
        Args:
            slide: 幻灯片对象
            content: 内容数据
        """
        try:
            if slide is None:
                logger.warning("slide对象为None，无法优化文本布局")
                return
                
            # 获取幻灯片尺寸信息
            slide_width = getattr(slide, 'slide_width', 9144000)  # 默认宽度值
            slide_height = getattr(slide, 'slide_height', 5143500)  # 默认高度值
            
            # 收集所有文本形状
            text_shapes = []
            for shape in slide.shapes:
                if not shape.has_text_frame or not shape.text_frame.text:
                    continue
                
                # 记录文本形状及其位置信息
                text_shapes.append({
                    'shape': shape,
                    'text': shape.text_frame.text.strip(),
                    'top': getattr(shape, 'top', 0),
                    'left': getattr(shape, 'left', 0),
                    'width': getattr(shape, 'width', 0),
                    'height': getattr(shape, 'height', 0),
                    'area': getattr(shape, 'width', 0) * getattr(shape, 'height', 0),
                    'is_title': getattr(shape, 'top', 0) < slide_height / 4,  # 顶部1/4区域视为标题区
                    'is_right_side': getattr(shape, 'left', 0) > slide_width / 2  # 右半部分
                })
            
            # 如果文本框太少，不需要优化
            if len(text_shapes) <= 1:
                return
            
            # 特殊处理：检查右侧区域是否有重复的标题文本
            if 'title' in content and content['title']:
                title_text = content['title']
                right_side_titles = []
                
                # 查找右侧区域中包含标题文本的形状
                for shape_info in text_shapes:
                    if shape_info['is_right_side'] and title_text in shape_info['text']:
                        right_side_titles.append(shape_info)
                
                # 如果右侧区域有标题文本，可能是重复的，需要清理
                if right_side_titles:
                    # 保留最大的一个（可能是主要内容区的标题）
                    if len(right_side_titles) > 1:
                        right_side_titles.sort(key=lambda x: x['area'], reverse=True)
                        for shape_info in right_side_titles[1:]:
                            shape_info['shape'].text_frame.text = ""
                            logger.info(f"删除右侧区域重复的标题文本: {shape_info['text'][:30]}...")
                    
                    # 如果右侧只有一个标题文本，检查是否有内容，如果没有则删除
                    elif 'content' in content and content['content'] and len(right_side_titles[0]['text']) < len(content['content']) * 0.5:
                        # 如果标题文本明显短于内容，可能是重复的标题，应该替换为内容
                        right_side_titles[0]['shape'].text_frame.text = content['content']
                        logger.info(f"将右侧区域的标题文本替换为内容: {right_side_titles[0]['text'][:30]} -> {content['content'][:30]}...")
                
            # 按垂直位置排序
            text_shapes.sort(key=lambda x: x['top'])
            
            # 检测重叠的文本框
            overlapping_shapes = []
            for i in range(len(text_shapes)):
                for j in range(i+1, len(text_shapes)):
                    shape1 = text_shapes[i]
                    shape2 = text_shapes[j]
                    
                    # 检查两个文本框是否重叠
                    if (shape1['left'] < shape2['left'] + shape2['width'] and
                        shape1['left'] + shape1['width'] > shape2['left'] and
                        shape1['top'] < shape2['top'] + shape2['height'] and
                        shape1['top'] + shape1['height'] > shape2['top']):
                        
                        # 计算重叠区域
                        overlap_width = min(shape1['left'] + shape1['width'], shape2['left'] + shape2['width']) - max(shape1['left'], shape2['left'])
                        overlap_height = min(shape1['top'] + shape1['height'], shape2['top'] + shape2['height']) - max(shape1['top'], shape2['top'])
                        overlap_area = overlap_width * overlap_height
                        
                        # 如果重叠面积超过较小文本框的30%，认为是显著重叠
                        smaller_area = min(shape1['area'], shape2['area'])
                        if overlap_area > smaller_area * 0.3:
                            # 记录重叠对
                            overlapping_shapes.append((shape1, shape2))
            
            # 处理重叠的文本框
            for shape1, shape2 in overlapping_shapes:
                # 如果内容相似，保留面积较大的，删除较小的
                similarity = self._text_similarity(shape1['text'], shape2['text'])
                if similarity > 0.5:  # 内容相似
                    if shape1['area'] >= shape2['area']:
                        # 保留shape1，删除shape2
                        shape2['shape'].text_frame.text = ""
                        logger.info(f"删除重叠的较小文本框: {shape2['text'][:30]}...")
                    else:
                        # 保留shape2，删除shape1
                        shape1['shape'].text_frame.text = ""
                        logger.info(f"删除重叠的较小文本框: {shape1['text'][:30]}...")
            
            # 处理右侧区域的多个文本框（常见于模板）
            right_side_shapes = [s for s in text_shapes if s['is_right_side'] and s['shape'].text_frame.text]
            
            # 如果右侧有多个内容文本框，可能需要合并或选择最佳的一个
            if len(right_side_shapes) > 1:
                # 按面积排序，找出最大的文本框
                right_side_shapes.sort(key=lambda x: x['area'], reverse=True)
                main_shape = right_side_shapes[0]
                
                # 检查主要内容是否已填充到最大的文本框
                if 'content' in content and content['content']:
                    has_content = False
                    for shape in right_side_shapes:
                        if content['content'] in shape['shape'].text_frame.text:
                            has_content = True
                            main_shape = shape  # 更新主要文本框为包含内容的文本框
                            break
                    
                    # 如果没有找到包含内容的文本框，填充到最大的文本框
                    if not has_content:
                        main_shape['shape'].text_frame.text = content['content']
                        logger.info(f"填充内容到右侧最大文本框")
                
                # 清除其他较小的右侧文本框
                for shape in right_side_shapes:
                    if shape['shape'] is not main_shape['shape']:
                        # 检查是否与主文本框内容重复
                        if self._text_similarity(shape['text'], main_shape['text']) > 0.3:
                            shape['shape'].text_frame.text = ""
                            logger.info(f"删除右侧重复文本框: {shape['text'][:30]}...")
            
            # 优化标题区域
            title_shapes = [s for s in text_shapes if s['is_title'] and s['shape'].text_frame.text]
            if len(title_shapes) > 1 and 'title' in content:
                # 找出最适合作为标题的文本框（通常是最上面的）
                title_shapes.sort(key=lambda x: x['top'])
                main_title = title_shapes[0]
                
                # 确保主标题包含正确的标题文本
                if content['title'] not in main_title['shape'].text_frame.text:
                    main_title['shape'].text_frame.text = content['title']
                
                # 清除其他标题区域的重复文本
                for shape in title_shapes[1:]:
                    if self._text_similarity(shape['text'], content['title']) > 0.5:
                        shape['shape'].text_frame.text = ""
                        logger.info(f"删除重复的标题文本: {shape['text']}")
            
            # 最后检查：确保没有重复的标题文本出现在内容区域
            if 'title' in content and content['title']:
                for shape_info in text_shapes:
                    # 如果不是主标题区域但包含标题文本
                    if not shape_info['is_title'] and self._text_similarity(shape_info['text'], content['title']) > 0.7:
                        # 如果文本框只包含标题文本，删除它
                        if len(shape_info['text']) <= len(content['title']) * 1.2:
                            shape_info['shape'].text_frame.text = ""
                            logger.info(f"删除内容区域的重复标题: {shape_info['text'][:30]}...")
        
        except Exception as e:
            logger.error(f"优化文本布局失败: {str(e)}")
            logger.error(traceback.format_exc())
    
    def _clear_template_content(self, slide):
        """
        清空模板中的所有文本内容，保留结构和设计元素
        
        Args:
            slide: 幻灯片对象
        """
        try:
            if slide is None:
                logger.warning("slide对象为None，无法清空内容")
                return
                
            # 收集需要保留的形状类型
            shapes_to_keep = []
            shapes_to_clear = []
            
            for shape in slide.shapes:
                # 检查是否是文本框
                if shape.has_text_frame:
                    # 保存原始形状，但清空文本
                    shapes_to_clear.append(shape)
                else:
                    # 非文本形状（如图片、形状等）保留
                    shapes_to_keep.append(shape)
            
            # 清空所有文本框的内容
            for shape in shapes_to_clear:
                try:
                    # 保留文本框但清空内容
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        # 记录原始文本（用于调试）
                        original_text = shape.text_frame.text.strip()
                        if original_text:
                            logger.info(f"清空文本框内容: {original_text[:30]}...")
                        
                        # 清空文本框
                        shape.text_frame.clear()
                        
                        # 确保文本框仍然可见
                        if hasattr(shape.text_frame, 'paragraphs') and len(shape.text_frame.paragraphs) == 0:
                            shape.text_frame.text = ""
                except Exception as e:
                    logger.error(f"清空文本框失败: {str(e)}")
            
            logger.info(f"已清空模板内容: {len(shapes_to_clear)}个文本框已清空，保留{len(shapes_to_keep)}个非文本元素")
            
        except Exception as e:
            logger.error(f"清空模板内容失败: {str(e)}")
            logger.error(traceback.format_exc())

    def create_ppt(self, slides_data, output_path):
        """
        创建PPT
        
        Args:
            slides_data: 幻灯片数据列表
            output_path: 输出文件路径
            
        Returns:
            bool: 成功返回True，失败返回False
        """
        try:
            logger.info(f"开始创建PPT: {output_path}")
            logger.info(f"使用模板: {self.template_path}")
            logger.info(f"使用元数据: {self.metadata_path}")
            logger.info(f"幻灯片数量: {len(slides_data)}")
            
            # 确保模板文件存在
            if not os.path.exists(self.template_path):
                logger.error(f"模板文件不存在: {self.template_path}")
                return False
                
            # 确保prs已加载
            if self.prs is None:
                logger.error("模板未成功加载，无法创建PPT")
                return False
                
            # 直接使用模板文件作为基础，而不是创建新的演示文稿
            try:
                output_prs = Presentation(self.template_path)
            except Exception as e:
                logger.error(f"加载模板失败: {str(e)}")
                return False
            
            if output_prs is None:
                logger.error("无法创建演示文稿对象")
                return False
                
            # 记录模板信息
            if self.metadata:
                logger.info(f"模板元数据: {len(self.metadata['slides'])}张幻灯片")
                for i, slide_info in enumerate(self.metadata['slides']):
                    suitable_for = slide_info.get('suitable_for', [])
                    logger.info(f"  幻灯片{i}: 适合类型={suitable_for}")
            else:
                logger.warning("没有模板元数据，将使用基本填充方法")
            
            # 获取模板中的幻灯片数量
            try:
                template_slides_count = len(output_prs.slides)
                logger.info(f"模板中包含 {template_slides_count} 张幻灯片")
            except Exception as e:
                logger.error(f"获取幻灯片数量失败: {str(e)}")
                return False
            
            # 如果内容页数超过模板页数，需要复制模板页来扩展
            if len(slides_data) > template_slides_count:
                logger.info(f"内容页数({len(slides_data)})超过模板页数({template_slides_count})，需要复制模板页")
                try:
                    # 保留原始幻灯片，以便后续复制
                    original_slides = list(output_prs.slides)
                    
                    # 需要添加的幻灯片数量
                    slides_to_add = len(slides_data) - template_slides_count
                    
                    # 复制模板中的幻灯片来扩展
                    for i in range(slides_to_add):
                        # 选择一个合适的模板页进行复制（这里简单地循环使用模板页）
                        template_index = (i % (template_slides_count - 2)) + 1  # 跳过首尾页
                        if template_index < 0 or template_index >= len(original_slides):
                            template_index = 0  # 默认使用第一页
                            
                        template_slide = original_slides[template_index]
                        
                        # 复制幻灯片
                        slide_layout = template_slide.slide_layout
                        output_prs.slides.add_slide(slide_layout)
                        
                        logger.info(f"复制了模板幻灯片 {template_index} 作为新的幻灯片")
                except Exception as e:
                    logger.error(f"复制幻灯片失败: {str(e)}")
                    # 继续处理，使用已有的幻灯片
            
            # 首先清空所有幻灯片的内容
            logger.info("开始清空模板中的所有文本内容")
            for i, slide in enumerate(output_prs.slides):
                self._clear_template_content(slide)
            
            # 填充内容到幻灯片
            try:
                slides_count = len(output_prs.slides)
                for i, slide_data in enumerate(slides_data):
                    if i >= slides_count:
                        logger.warning(f"内容页数超出，跳过第 {i+1} 页")
                        continue
                        
                    logger.info(f"处理第 {i+1} 张幻灯片")
                    
                    # 获取幻灯片类型和布局
                    slide_type = slide_data.get('type', '').lower()
                    slide_layout = slide_data.get('layout', '').lower()
                    
                    # 确定内容类型
                    content_type = 'content'  # 默认类型
                    if i == 0 or slide_type == 'cover' or slide_layout == 'cover':
                        content_type = 'cover'
                    elif i == len(slides_data) - 1 or slide_type == 'summary' or slide_type == 'conclusion' or slide_layout == 'summary':
                        content_type = 'summary'
                    elif 'keypoints' in slide_data and slide_data['keypoints']:
                        content_type = 'keypoints'
                    elif 'table' in slide_data and slide_data['table']:
                        content_type = 'table'
                    elif 'image' in slide_data and slide_data['image'] or slide_layout == 'image':
                        content_type = 'image'
                    
                    logger.info(f"  内容类型: {content_type}")
                    
                    # 获取当前幻灯片
                    slide = output_prs.slides[i]
                    
                    # 填充内容到幻灯片，保留原有背景和设计元素
                    self.fill_slide(slide, slide_data)
            except Exception as e:
                logger.error(f"填充幻灯片内容失败: {str(e)}")
                logger.error(traceback.format_exc())
                # 尝试继续保存
            
            # 保存演示文稿
            try:
                logger.info(f"保存PPT到: {output_path}")
                output_prs.save(output_path)
                logger.info("PPT保存成功")
                
                # 验证文件是否已保存
                if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
                    logger.info(f"PPT文件保存成功，大小: {os.path.getsize(output_path)} 字节")
                    return True
                else:
                    logger.error("保存PPT失败，文件不存在或大小为0")
                    return False
            except Exception as e:
                logger.error(f"保存PPT失败: {str(e)}")
                logger.error(traceback.format_exc())
                return False
        except Exception as e:
            logger.error(f"创建PPT失败: {str(e)}")
            logger.error(traceback.format_exc())
            return False

def fill_ppt_template(template_path, slides_data, output_path, metadata_path=None):
    """
    填充PPT模板的便捷函数
    
    Args:
        template_path: 模板文件路径
        slides_data: 幻灯片数据列表
        output_path: 输出文件路径
        metadata_path: 元数据文件路径，如果为None则自动查找
        
    Returns:
        bool: 成功返回True，失败返回False
    """
    try:
        filler = PPTTemplateFiller(template_path, metadata_path)
        return filler.create_ppt(slides_data, output_path)
    except Exception as e:
        logger.error(f"填充PPT模板失败: {str(e)}")
        logger.error(traceback.format_exc())
        return False

if __name__ == "__main__":
    # 命令行使用方式
    if len(sys.argv) < 4:
        print("用法: python ppt_fill_template.py <template.pptx> <slides_data.json> <output.pptx> [metadata.json]")
        sys.exit(1)
        
    template_path = sys.argv[1]
    slides_data_path = sys.argv[2]
    output_path = sys.argv[3]
    metadata_path = sys.argv[4] if len(sys.argv) > 4 else None
    
    try:
        # 加载幻灯片数据
        with open(slides_data_path, 'r', encoding='utf-8') as f:
                slides_data = json.load(f)
        
        # 填充模板
        success = fill_ppt_template(template_path, slides_data, output_path, metadata_path)
        
        if success:
            print(f"PPT生成成功: {output_path}")
        else:
            print("PPT生成失败")
            sys.exit(1)
    except Exception as e:
        print(f"错误: {e}")
        sys.exit(1) 
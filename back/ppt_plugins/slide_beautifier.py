"""
幻灯片美化插件
提供增强视觉效果的组件
"""

import logging
import traceback
import os
from io import BytesIO
import random
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from improved_ppt_generator import SlideComponentBase

# 配置日志
logger = logging.getLogger("slide_beautifier_plugin")

# 检查PIL可用性
try:
    from PIL import Image, ImageEnhance, ImageFilter, ImageDraw
    HAS_PIL = True
except ImportError:
    HAS_PIL = False
    logger.warning("PIL未安装，高级图片处理功能不可用")

class BackgroundStyleComponent(SlideComponentBase):
    """设置幻灯片背景样式的组件"""
    
    def apply(self, slide, data, context=None):
        """应用背景样式到幻灯片"""
        # 如果数据中没有指定背景样式，则不做处理
        if not data.get('background_style'):
            return slide
            
        try:
            bg_style = data['background_style']
            style_type = bg_style.get('type', 'solid')
            
            # 获取幻灯片的背景对象
            background = slide.background
            
            # 根据样式类型设置不同的背景
            if style_type == 'solid':
                self._apply_solid_background(slide, bg_style, context)
            elif style_type == 'gradient':
                self._apply_gradient_background(slide, bg_style, context)
            elif style_type == 'pattern':
                self._apply_pattern_background(slide, bg_style, context)
            elif style_type == 'image':
                self._apply_image_background(slide, bg_style, context)
                
        except Exception as e:
            logger.error(f"设置背景样式失败: {str(e)}")
            logger.error(traceback.format_exc())
            
        return slide
        
    def _apply_solid_background(self, slide, bg_style, context):
        """应用纯色背景"""
        # 获取颜色
        color_str = bg_style.get('color', '#FFFFFF')
        
        # 将颜色字符串转换为RGB
        if color_str.startswith('#'):
            r = int(color_str[1:3], 16)
            g = int(color_str[3:5], 16)
            b = int(color_str[5:7], 16)
        else:
            # 默认为白色
            r, g, b = 255, 255, 255
            
        # 设置背景填充
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
        
    def _apply_gradient_background(self, slide, bg_style, context):
        """应用渐变背景"""
        # 在没有原生渐变支持的情况下，创建一个渐变图像背景
        if HAS_PIL:
            width = slide.slide_width
            height = slide.slide_height
            
            # 创建一个新图像
            img = Image.new('RGB', (int(width / 9525), int(height / 9525)))  # 9525为EMU到像素的转换比例
            draw = ImageDraw.Draw(img)
            
            # 获取渐变方向
            direction = bg_style.get('direction', 'horizontal')
            
            # 获取颜色
            start_color_str = bg_style.get('start_color', '#FFFFFF')
            end_color_str = bg_style.get('end_color', '#000000')
            
            # 解析颜色
            if start_color_str.startswith('#'):
                start_r = int(start_color_str[1:3], 16)
                start_g = int(start_color_str[3:5], 16)
                start_b = int(start_color_str[5:7], 16)
            else:
                start_r, start_g, start_b = 255, 255, 255
                
            if end_color_str.startswith('#'):
                end_r = int(end_color_str[1:3], 16)
                end_g = int(end_color_str[3:5], 16)
                end_b = int(end_color_str[5:7], 16)
            else:
                end_r, end_g, end_b = 0, 0, 0
                
            # 创建渐变
            img_width = img.width
            img_height = img.height
            
            if direction == 'horizontal':
                for x in range(img_width):
                    # 计算插值比例
                    ratio = x / img_width
                    
                    # 插值计算当前颜色
                    r = int(start_r + (end_r - start_r) * ratio)
                    g = int(start_g + (end_g - start_g) * ratio)
                    b = int(start_b + (end_b - start_b) * ratio)
                    
                    # 绘制垂直线
                    draw.line([(x, 0), (x, img_height)], fill=(r, g, b))
            else:  # vertical
                for y in range(img_height):
                    # 计算插值比例
                    ratio = y / img_height
                    
                    # 插值计算当前颜色
                    r = int(start_r + (end_r - start_r) * ratio)
                    g = int(start_g + (end_g - start_g) * ratio)
                    b = int(start_b + (end_b - start_b) * ratio)
                    
                    # 绘制水平线
                    draw.line([(0, y), (img_width, y)], fill=(r, g, b))
            
            # 将图像保存到内存中
            img_byte_arr = BytesIO()
            img.save(img_byte_arr, format='PNG')
            img_byte_arr.seek(0)
            
            # 设置背景图片
            slide.shapes.add_picture(img_byte_arr, 0, 0, width, height)
            
    def _apply_pattern_background(self, slide, bg_style, context):
        """应用图案背景"""
        # 创建一个基本的图案背景
        if HAS_PIL:
            width = slide.slide_width
            height = slide.slide_height
            
            # 获取图案类型和颜色
            pattern_type = bg_style.get('pattern', 'dots')
            bg_color_str = bg_style.get('bg_color', '#FFFFFF')
            fg_color_str = bg_style.get('fg_color', '#000000')
            
            # 解析颜色
            if bg_color_str.startswith('#'):
                bg_r = int(bg_color_str[1:3], 16)
                bg_g = int(bg_color_str[3:5], 16)
                bg_b = int(bg_color_str[5:7], 16)
            else:
                bg_r, bg_g, bg_b = 255, 255, 255
                
            if fg_color_str.startswith('#'):
                fg_r = int(fg_color_str[1:3], 16)
                fg_g = int(fg_color_str[3:5], 16)
                fg_b = int(fg_color_str[5:7], 16)
            else:
                fg_r, fg_g, fg_b = 0, 0, 0
                
            # 创建基础图像
            img = Image.new('RGB', (int(width / 9525), int(height / 9525)), (bg_r, bg_g, bg_b))
            draw = ImageDraw.Draw(img)
            
            # 创建图案
            img_width = img.width
            img_height = img.height
            
            if pattern_type == 'dots':
                # 创建点状图案
                dot_size = 2
                spacing = 20
                for x in range(0, img_width, spacing):
                    for y in range(0, img_height, spacing):
                        draw.ellipse([(x, y), (x + dot_size, y + dot_size)], 
                                     fill=(fg_r, fg_g, fg_b))
            elif pattern_type == 'grid':
                # 创建网格图案
                spacing = 30
                for x in range(0, img_width, spacing):
                    draw.line([(x, 0), (x, img_height)], fill=(fg_r, fg_g, fg_b), width=1)
                for y in range(0, img_height, spacing):
                    draw.line([(0, y), (img_width, y)], fill=(fg_r, fg_g, fg_b), width=1)
            elif pattern_type == 'stripes':
                # 创建条纹图案
                spacing = 20
                for y in range(0, img_height, spacing * 2):
                    draw.rectangle([(0, y), (img_width, y + spacing)], 
                                   fill=(fg_r, fg_g, fg_b))
            
            # 将图像保存到内存中
            img_byte_arr = BytesIO()
            img.save(img_byte_arr, format='PNG')
            img_byte_arr.seek(0)
            
            # 设置背景图片
            slide.shapes.add_picture(img_byte_arr, 0, 0, width, height)
            
    def _apply_image_background(self, slide, bg_style, context):
        """应用图片背景"""
        # 获取图片路径或URL
        image_path = bg_style.get('image_path')
        
        if not image_path:
            return
            
        try:
            # 读取图片数据
            image_data = None
            
            if os.path.exists(image_path):
                # 从本地文件读取
                with open(image_path, 'rb') as f:
                    image_data = f.read()
            elif image_path.startswith('http'):
                # 从URL下载
                import requests
                response = requests.get(image_path)
                response.raise_for_status()
                image_data = response.content
                
            if image_data:
                # 如果需要处理图片
                if HAS_PIL and bg_style.get('process_image', False):
                    # 加载图片
                    img = Image.open(BytesIO(image_data))
                    
                    # 应用各种图片处理
                    opacity = bg_style.get('opacity', 1.0)
                    brightness = bg_style.get('brightness', 1.0)
                    contrast = bg_style.get('contrast', 1.0)
                    blur = bg_style.get('blur', 0)
                    
                    # 调整亮度和对比度
                    enhancer = ImageEnhance.Brightness(img)
                    img = enhancer.enhance(brightness)
                    
                    enhancer = ImageEnhance.Contrast(img)
                    img = enhancer.enhance(contrast)
                    
                    # 模糊处理
                    if blur > 0:
                        img = img.filter(ImageFilter.GaussianBlur(radius=blur))
                    
                    # 调整透明度
                    if opacity < 1.0 and img.mode != 'RGBA':
                        img = img.convert('RGBA')
                        
                        # 创建透明度遮罩
                        alpha = Image.new('L', img.size, int(opacity * 255))
                        img.putalpha(alpha)
                    
                    # 将处理后的图片保存到内存
                    img_byte_arr = BytesIO()
                    img.save(img_byte_arr, format='PNG')
                    img_byte_arr.seek(0)
                    image_data = img_byte_arr.read()
                
                # 设置背景图片
                # 先创建一个全幻灯片大小的形状
                left = 0
                top = 0
                width = slide.slide_width
                height = slide.slide_height
                
                # 添加图片
                slide.shapes.add_picture(BytesIO(image_data), left, top, width, height)
                
        except Exception as e:
            logger.error(f"设置图片背景失败: {str(e)}")
            logger.error(traceback.format_exc())

class ShapeStyleComponent(SlideComponentBase):
    """美化幻灯片形状的组件"""
    
    def apply(self, slide, data, context=None):
        """应用形状美化"""
        if not data.get('shape_style'):
            return slide
            
        try:
            # 获取形状样式配置
            shape_style = data['shape_style']
            
            # 对幻灯片中的每个形状应用样式
            for shape in slide.shapes:
                self._apply_shape_style(shape, shape_style, context)
                
        except Exception as e:
            logger.error(f"应用形状样式失败: {str(e)}")
            logger.error(traceback.format_exc())
            
        return slide
        
    def _apply_shape_style(self, shape, shape_style, context):
        """应用样式到单个形状"""
        # 跳过图片形状
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return
            
        # 应用轮廓样式
        if hasattr(shape, 'line') and shape_style.get('outline', False):
            outline = shape_style.get('outline', {})
            
            # 轮廓颜色
            color_str = outline.get('color')
            if color_str and color_str.startswith('#'):
                r = int(color_str[1:3], 16)
                g = int(color_str[3:5], 16)
                b = int(color_str[5:7], 16)
                shape.line.color.rgb = RGBColor(r, g, b)
                
            # 轮廓宽度
            width = outline.get('width')
            if width:
                shape.line.width = Pt(width)
        
        # 应用填充样式
        if hasattr(shape, 'fill') and shape_style.get('fill', False):
            fill = shape_style.get('fill', {})
            
            # 填充颜色
            color_str = fill.get('color')
            if color_str and color_str.startswith('#'):
                r = int(color_str[1:3], 16)
                g = int(color_str[3:5], 16)
                b = int(color_str[5:7], 16)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(r, g, b)
                
        # 应用阴影
        if hasattr(shape, 'shadow') and shape_style.get('shadow', False):
            pass  # python-pptx目前不支持阴影属性设置
            
        # 应用文字格式
        if shape.has_text_frame and shape_style.get('text', False):
            text_style = shape_style.get('text', {})
            
            for paragraph in shape.text_frame.paragraphs:
                # 段落对齐方式
                align = text_style.get('align')
                if align:
                    if align == 'left':
                        paragraph.alignment = PP_ALIGN.LEFT
                    elif align == 'center':
                        paragraph.alignment = PP_ALIGN.CENTER
                    elif align == 'right':
                        paragraph.alignment = PP_ALIGN.RIGHT
                
                # 字体样式
                for run in paragraph.runs:
                    # 字体大小
                    size = text_style.get('size')
                    if size:
                        run.font.size = Pt(size)
                    
                    # 字体颜色
                    color_str = text_style.get('color')
                    if color_str and color_str.startswith('#'):
                        r = int(color_str[1:3], 16)
                        g = int(color_str[3:5], 16)
                        b = int(color_str[5:7], 16)
                        run.font.color.rgb = RGBColor(r, g, b)
                    
                    # 粗体
                    if 'bold' in text_style:
                        run.font.bold = text_style['bold']
                    
                    # 斜体
                    if 'italic' in text_style:
                        run.font.italic = text_style['italic']
                    
                    # 下划线
                    if 'underline' in text_style:
                        run.font.underline = text_style['underline']

class FooterComponent(SlideComponentBase):
    """添加页脚的组件"""
    
    def apply(self, slide, data, context=None):
        """添加页脚到幻灯片"""
        if not data.get('footer'):
            return slide
            
        try:
            footer_config = data['footer']
            
            # 获取页脚文本
            footer_text = footer_config.get('text', '')
            show_slide_number = footer_config.get('show_slide_number', True)
            show_date = footer_config.get('show_date', False)
            
            # 在幻灯片底部添加文本框作为页脚
            left = Inches(0.5)
            top = slide.slide_height - Inches(0.6)
            width = slide.slide_width - Inches(1)
            height = Inches(0.5)
            
            # 创建文本框
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            
            # 构建页脚文本
            full_text = footer_text
            
            # 添加日期
            if show_date:
                import datetime
                today = datetime.date.today().strftime("%Y-%m-%d")
                if full_text:
                    full_text += " | "
                full_text += today
            
            # 添加幻灯片编号
            if show_slide_number and context and 'slide_number' in context:
                if full_text:
                    full_text += " | "
                full_text += f"第 {context['slide_number']} 页"
            
            # 设置页脚文本
            text_frame.text = full_text
            
            # 设置文本格式
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.RIGHT
                for run in paragraph.runs:
                    run.font.size = Pt(8)
                    run.font.color.rgb = RGBColor(128, 128, 128)
                    
        except Exception as e:
            logger.error(f"添加页脚失败: {str(e)}")
            logger.error(traceback.format_exc())
            
        return slide

class WatermarkComponent(SlideComponentBase):
    """添加水印的组件"""
    
    def apply(self, slide, data, context=None):
        """添加水印到幻灯片"""
        # 如果全局上下文中有水印设置或幻灯片数据中指定了水印
        watermark_config = None
        
        if context and context.get('watermark'):
            watermark_config = context['watermark']
        elif data.get('watermark'):
            watermark_config = data['watermark']
            
        if not watermark_config:
            return slide
            
        try:
            # 获取水印文本和样式
            text = watermark_config.get('text', '')
            if not text:
                return slide
                
            # 水印位置和大小
            left = Inches(1)
            top = Inches(1.5)
            width = slide.slide_width - Inches(2)
            height = slide.slide_height - Inches(3)
            
            # 创建文本框
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            
            # 设置水印文本
            text_frame.text = text
            
            # 设置文本格式
            for paragraph in text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    # 设置字体大小
                    run.font.size = Pt(watermark_config.get('size', 48))
                    
                    # 设置字体颜色和透明度
                    color_str = watermark_config.get('color', '#CCCCCC')
                    if color_str.startswith('#'):
                        r = int(color_str[1:3], 16)
                        g = int(color_str[3:5], 16)
                        b = int(color_str[5:7], 16)
                    else:
                        r, g, b = 204, 204, 204  # 默认浅灰色
                        
                    run.font.color.rgb = RGBColor(r, g, b)
                    
                    # 设置旋转角度
                    # 注意: python-pptx目前不支持设置文本框的旋转角度
                    # 需要通过其他方式实现
                    
        except Exception as e:
            logger.error(f"添加水印失败: {str(e)}")
            logger.error(traceback.format_exc())
            
        return slide 
import os
import sys
import logging
import json
import shutil
from pathlib import Path
from pptx import Presentation

# 配置日志
logger = logging.getLogger("ppt_engine.template_manager")

class TemplateManager:
    """模板管理系统，负责模板的读取、解析和转换"""
    
    def __init__(self):
        """初始化模板管理器"""
        self.template_cache = {}
        self.base_path = Path(os.path.dirname(os.path.abspath(__file__)))
        self.template_directory = os.path.join(self.base_path.parent, "ppt_templates")
        self.html_template_directory = os.path.join(self.base_path, "html_templates")
        
        # 确保目录存在
        os.makedirs(self.template_directory, exist_ok=True)
        os.makedirs(self.html_template_directory, exist_ok=True)
        
    def get_template(self, template_name=None):
        """
        获取指定模板，如未指定则返回默认模板
        
        Args:
            template_name: 模板名称或路径
            
        Returns:
            template_data: 包含模板信息的字典
        """
        if not template_name:
            template_name = self._get_default_template()
            
        # 检查缓存
        if template_name in self.template_cache:
            return self.template_cache[template_name]
            
        # 确定模板路径
        template_path = self._get_template_path(template_name)
        if not os.path.exists(template_path):
            raise ValueError(f"模板不存在: {template_path}")
            
        # 转换模板
        template_data = self._convert_template(template_path)
        
        # 缓存模板
        self.template_cache[template_name] = template_data
        
        return template_data
        
    def _get_default_template(self):
        """获取默认模板名称"""
        # 首选绿色圆点模板，其次是任何可用模板
        default_template = "绿色圆点.pptx"
        if os.path.exists(os.path.join(self.template_directory, default_template)):
            return default_template
            
        # 查找任何PPTX文件
        for file in os.listdir(self.template_directory):
            if file.endswith(".pptx"):
                return file
                
        raise ValueError("找不到默认模板，请确保ppt_templates目录中有.pptx文件")
    
    def _get_template_path(self, template_name):
        """获取模板的完整路径"""
        # 如果提供的是完整路径
        if os.path.exists(template_name):
            return template_name
            
        # 如果提供的是模板名称
        if not template_name.endswith(".pptx"):
            template_name += ".pptx"
            
        return os.path.join(self.template_directory, template_name)
        
    def _convert_template(self, template_path):
        """
        将PPTX模板转换为HTML模板
        
        Args:
            template_path: PPTX模板路径
            
        Returns:
            template_data: 包含模板信息的字典
        """
        logger.info(f"开始转换模板: {template_path}")
        
        # 获取模板名称
        template_name = os.path.splitext(os.path.basename(template_path))[0]
        
        # 创建HTML模板目录
        html_template_dir = os.path.join(self.html_template_directory, template_name)
        os.makedirs(html_template_dir, exist_ok=True)
        
        # 加载PPTX文件
        prs = Presentation(template_path)
        
        # 提取模板信息
        template_info = self._extract_template_info(prs, template_name)
        
        # 提取布局信息
        layouts = self._extract_layouts(prs)
        
        # 提取主题样式
        theme = self._extract_theme(prs)
        
        # 保存模板元数据
        metadata_path = os.path.join(html_template_dir, "template_info.json")
        with open(metadata_path, 'w', encoding='utf-8') as f:
            json.dump({
                "name": template_name,
                "info": template_info,
                "layouts": layouts,
                "theme": theme
            }, f, ensure_ascii=False, indent=2)
            
        # 创建HTML模板文件
        self._create_html_templates(html_template_dir, layouts, theme)
        
        # 返回模板数据
        return {
            "name": template_name,
            "info": template_info,
            "layouts": layouts,
            "theme": theme,
            "dir": html_template_dir
        }
        
    def _extract_template_info(self, prs, template_name):
        """提取模板基本信息"""
        # 获取幻灯片尺寸
        width = prs.slide_width
        height = prs.slide_height
        
        # 计算幻灯片数量
        slide_count = len(prs.slides)
        
        # 计算布局数量
        layout_count = len(prs.slide_layouts)
        
        # 读取模板描述（如果有）
        description = ""
        json_path = os.path.join(self.template_directory, f"{template_name}.json")
        if os.path.exists(json_path):
            try:
                with open(json_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    description = data.get('desc', '')
            except Exception as e:
                logger.warning(f"读取模板描述失败: {str(e)}")
        
        return {
            "width": width,
            "height": height,
            "slide_count": slide_count,
            "layout_count": layout_count,
            "description": description
        }
        
    def _extract_layouts(self, prs):
        """提取幻灯片布局信息"""
        layouts = {}
        
        for i, layout in enumerate(prs.slide_layouts):
            # 获取布局名称
            name = layout.name or f"Layout {i+1}"
            
            # 获取占位符信息
            placeholders = []
            for ph in layout.placeholders:
                placeholders.append({
                    "idx": ph.placeholder_format.idx,
                    "type": ph.placeholder_format.type,
                    "name": ph.name,
                    "position": {
                        "left": ph.left,
                        "top": ph.top,
                        "width": ph.width,
                        "height": ph.height
                    }
                })
                
            # 确定布局类型
            layout_type = self._determine_layout_type(layout, placeholders)
            
            layouts[name] = {
                "id": i,
                "type": layout_type,
                "placeholders": placeholders
            }
            
        return layouts
        
    def _determine_layout_type(self, layout, placeholders):
        """确定布局类型"""
        name = layout.name.lower() if layout.name else ""
        
        # 根据名称判断
        if any(x in name for x in ["title", "标题", "封面"]):
            return "title"
        elif any(x in name for x in ["section", "章节"]):
            return "section"
        elif any(x in name for x in ["比较", "compare"]):
            return "comparison"
        elif any(x in name for x in ["图文", "image"]):
            return "image_text"
        elif any(x in name for x in ["表格", "table"]):
            return "table"
        elif any(x in name for x in ["图表", "chart"]):
            return "chart"
        elif any(x in name for x in ["列表", "bullet"]):
            return "bullet_points"
        
        # 根据占位符判断
        has_title = False
        has_content = False
        has_image = False
        
        for ph in placeholders:
            ph_type = ph.get("type")
            if ph_type == 1:  # 标题占位符
                has_title = True
            elif ph_type == 2:  # 内容占位符
                has_content = True
            elif ph_type == 18:  # 图片占位符
                has_image = True
                
        if has_title and has_content and has_image:
            return "image_text"
        elif has_title and has_content:
            return "content"
        elif has_title:
            return "title_only"
        
        # 默认类型
        return "blank"
        
    def _extract_theme(self, prs):
        """提取模板主题样式"""
        theme = {
            "colors": {
                "primary": "#2E7D32",  # 默认主色
                "secondary": "#4CAF50",  # 默认辅助色
                "background": "#FFFFFF",
                "text": "#333333",
                "accent1": "#8BC34A",
                "accent2": "#F1F8E9"
            },
            "fonts": {
                "heading": "PingFang SC, Microsoft YaHei, sans-serif",
                "body": "PingFang SC, Microsoft YaHei, sans-serif"
            }
        }
        
        # 尝试从模板提取颜色
        try:
            if hasattr(prs, 'theme') and prs.theme:
                color_scheme = prs.theme.theme_elements.clrScheme
                
                # 提取主要颜色
                if len(color_scheme.srgbClr) > 0:
                    theme["colors"]["primary"] = "#" + color_scheme.srgbClr[0].val
                    
                # 提取更多颜色
                if len(color_scheme.srgbClr) > 1:
                    theme["colors"]["secondary"] = "#" + color_scheme.srgbClr[1].val
                    
                # 提取背景色
                if len(color_scheme.srgbClr) > 2:
                    theme["colors"]["background"] = "#" + color_scheme.srgbClr[2].val
        except Exception as e:
            logger.warning(f"提取主题颜色失败: {str(e)}")
            
        # 尝试提取字体
        try:
            if hasattr(prs, 'theme') and prs.theme:
                font_scheme = prs.theme.theme_elements.fontScheme
                
                # 提取标题字体
                if hasattr(font_scheme, 'majorFont') and font_scheme.majorFont:
                    latin_font = font_scheme.majorFont.latin
                    if hasattr(latin_font, 'typeface'):
                        theme["fonts"]["heading"] = latin_font.typeface
                        
                # 提取正文字体
                if hasattr(font_scheme, 'minorFont') and font_scheme.minorFont:
                    latin_font = font_scheme.minorFont.latin
                    if hasattr(latin_font, 'typeface'):
                        theme["fonts"]["body"] = latin_font.typeface
        except Exception as e:
            logger.warning(f"提取主题字体失败: {str(e)}")
            
        return theme
        
    def _create_html_templates(self, output_dir, layouts, theme):
        """创建HTML模板文件"""
        # 复制基础样式文件
        self._copy_base_resources(output_dir)
        
        # 为每种布局类型创建HTML模板
        layout_types = set(layout["type"] for layout in layouts.values())
        
        for layout_type in layout_types:
            self._create_layout_template(output_dir, layout_type, theme)
            
        # 创建默认模板
        self._create_default_template(output_dir, theme)
        
        logger.info(f"HTML模板创建完成，保存在: {output_dir}")
        
    def _copy_base_resources(self, output_dir):
        """复制基础资源文件"""
        # 基础CSS
        base_css = """
        body {
            margin: 0;
            padding: 0;
            font-family: var(--body-font);
        }
        
        .slide {
            width: 1280px;
            height: 720px;
            position: relative;
            overflow: hidden;
            background-color: var(--background-color);
        }
        
        .slide-header {
            padding: 20px 60px;
        }
        
        .slide-title {
            font-family: var(--heading-font);
            font-size: 36px;
            font-weight: 700;
            color: var(--primary-color);
            margin: 0;
        }
        
        .slide-content {
            padding: 20px 60px;
        }
        
        .slide-footer {
            position: absolute;
            bottom: 0;
            left: 0;
            right: 0;
            padding: 16px 60px;
            display: flex;
            justify-content: flex-end;
            color: #78909C;
            font-size: 18px;
        }
        """
        
        with open(os.path.join(output_dir, "base.css"), 'w', encoding='utf-8') as f:
            f.write(base_css)
            
    def _create_layout_template(self, output_dir, layout_type, theme):
        """创建特定布局类型的HTML模板"""
        # 根据布局类型创建不同的模板
        html_content = ""
        css_content = ""
        
        if layout_type == "title":
            html_content, css_content = self._create_title_layout(theme)
        elif layout_type == "content" or layout_type == "bullet_points":
            html_content, css_content = self._create_content_layout(theme)
        elif layout_type == "image_text":
            html_content, css_content = self._create_image_text_layout(theme)
        elif layout_type == "comparison":
            html_content, css_content = self._create_comparison_layout(theme)
        elif layout_type == "table":
            html_content, css_content = self._create_table_layout(theme)
        elif layout_type == "chart":
            html_content, css_content = self._create_chart_layout(theme)
        else:
            # 默认空白布局
            html_content, css_content = self._create_blank_layout(theme)
            
        # 保存HTML模板
        with open(os.path.join(output_dir, f"{layout_type}.html"), 'w', encoding='utf-8') as f:
            f.write(html_content)
            
        # 保存CSS样式
        with open(os.path.join(output_dir, f"{layout_type}.css"), 'w', encoding='utf-8') as f:
            f.write(css_content)
            
    def _create_default_template(self, output_dir, theme):
        """创建默认模板"""
        # 创建默认HTML
        html_content = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>{{title}}</title>
    <link rel="stylesheet" href="base.css">
    <link rel="stylesheet" href="{{layout_type}}.css">
    <style>
        :root {
            --primary-color: {{primary_color}};
            --secondary-color: {{secondary_color}};
            --background-color: {{background_color}};
            --text-color: {{text_color}};
            --heading-font: {{heading_font}};
            --body-font: {{body_font}};
        }
        {{additional_css}}
    </style>
</head>
<body>
    <div class="slide">
        {{content}}
        <div class="slide-footer">{{page_number}}</div>
    </div>
</body>
</html>"""
        
        with open(os.path.join(output_dir, "default.html"), 'w', encoding='utf-8') as f:
            f.write(html_content)
            
    def _create_title_layout(self, theme):
        """创建标题页布局"""
        html = """<div class="slide-content title-content">
    <h1 class="main-title">{{title}}</h1>
    <h2 class="subtitle">{{subtitle}}</h2>
</div>"""
        
        css = """
.title-content {
    display: flex;
    flex-direction: column;
    justify-content: center;
    align-items: center;
    height: 80%;
    text-align: center;
}

.main-title {
    font-size: 60px;
    font-weight: 700;
    color: var(--primary-color);
    margin-bottom: 30px;
}

.subtitle {
    font-size: 32px;
    font-weight: 500;
    color: var(--secondary-color);
}
"""
        
        return html, css
        
    def _create_content_layout(self, theme):
        """创建内容页布局"""
        html = """<div class="slide-header">
    <h1 class="slide-title">{{title}}</h1>
</div>
<div class="slide-content">
    <div class="bullet-list">
        {{#each bullet_points}}
        <div class="bullet-item">
            <div class="bullet-marker"></div>
            <div class="bullet-text">{{this}}</div>
        </div>
        {{/each}}
    </div>
</div>"""
        
        css = """
.bullet-list {
    margin-top: 20px;
}

.bullet-item {
    display: flex;
    align-items: flex-start;
    margin-bottom: 16px;
}

.bullet-marker {
    width: 12px;
    height: 12px;
    background-color: var(--primary-color);
    border-radius: 50%;
    margin-top: 8px;
    margin-right: 16px;
    flex-shrink: 0;
}

.bullet-text {
    font-size: 24px;
    line-height: 1.5;
    color: var(--text-color);
}
"""
        
        return html, css
        
    def _create_image_text_layout(self, theme):
        """创建图文布局"""
        html = """<div class="slide-header">
    <h1 class="slide-title">{{title}}</h1>
</div>
<div class="slide-content">
    <div class="image-text-container">
        <div class="image-container">
            <img src="{{image_url}}" alt="{{image_alt}}">
        </div>
        <div class="text-container">
            <div class="bullet-list">
                {{#each bullet_points}}
                <div class="bullet-item">
                    <div class="bullet-marker"></div>
                    <div class="bullet-text">{{this}}</div>
                </div>
                {{/each}}
            </div>
        </div>
    </div>
</div>"""
        
        css = """
.image-text-container {
    display: grid;
    grid-template-columns: 1.2fr 1fr;
    gap: 40px;
    height: 100%;
}

.image-container {
    display: flex;
    align-items: center;
    justify-content: center;
}

.image-container img {
    max-width: 100%;
    max-height: 500px;
    object-fit: contain;
}

.text-container {
    padding: 20px 0;
}

.bullet-list {
    margin-top: 0;
}

.bullet-item {
    display: flex;
    align-items: flex-start;
    margin-bottom: 16px;
}

.bullet-marker {
    width: 12px;
    height: 12px;
    background-color: var(--primary-color);
    border-radius: 50%;
    margin-top: 8px;
    margin-right: 16px;
    flex-shrink: 0;
}

.bullet-text {
    font-size: 22px;
    line-height: 1.5;
    color: var(--text-color);
}
"""
        
        return html, css
        
    def _create_comparison_layout(self, theme):
        """创建比较布局"""
        html = """<div class="slide-header">
    <h1 class="slide-title">{{title}}</h1>
</div>
<div class="slide-content">
    <table class="comparison-table">
        <thead>
            <tr>
                <th class="first-column">特征</th>
                {{#each columns}}
                <th>{{this}}</th>
                {{/each}}
            </tr>
        </thead>
        <tbody>
            {{#each rows}}
            <tr>
                <td class="row-title">{{name}}</td>
                {{#each values}}
                <td>{{this}}</td>
                {{/each}}
            </tr>
            {{/each}}
        </tbody>
    </table>
</div>"""
        
        css = """
.comparison-table {
    width: 100%;
    border-collapse: collapse;
    margin-top: 20px;
}

.comparison-table th {
    background-color: var(--primary-color);
    color: white;
    font-weight: 600;
    font-size: 22px;
    text-align: center;
    padding: 12px 20px;
}

.comparison-table .first-column {
    background-color: var(--secondary-color);
}

.comparison-table td {
    padding: 12px 20px;
    border: 1px solid #e0e0e0;
    font-size: 20px;
    text-align: center;
}

.comparison-table .row-title {
    font-weight: 600;
    background-color: rgba(0, 0, 0, 0.03);
    text-align: left;
}

.comparison-table tr:nth-child(even) {
    background-color: rgba(0, 0, 0, 0.02);
}
"""
        
        return html, css
        
    def _create_table_layout(self, theme):
        """创建表格布局"""
        html = """<div class="slide-header">
    <h1 class="slide-title">{{title}}</h1>
</div>
<div class="slide-content">
    <table class="data-table">
        <thead>
            <tr>
                {{#each headers}}
                <th>{{this}}</th>
                {{/each}}
            </tr>
        </thead>
        <tbody>
            {{#each rows}}
            <tr>
                {{#each this}}
                <td>{{this}}</td>
                {{/each}}
            </tr>
            {{/each}}
        </tbody>
    </table>
</div>"""
        
        css = """
.data-table {
    width: 100%;
    border-collapse: collapse;
    margin-top: 20px;
}

.data-table th {
    background-color: var(--primary-color);
    color: white;
    font-weight: 600;
    font-size: 22px;
    text-align: center;
    padding: 12px 20px;
}

.data-table td {
    padding: 12px 20px;
    border: 1px solid #e0e0e0;
    font-size: 20px;
    text-align: center;
}

.data-table tr:nth-child(even) {
    background-color: rgba(0, 0, 0, 0.02);
}
"""
        
        return html, css
        
    def _create_chart_layout(self, theme):
        """创建图表布局"""
        html = """<div class="slide-header">
    <h1 class="slide-title">{{title}}</h1>
</div>
<div class="slide-content">
    <div class="chart-container">
        <img src="{{chart_url}}" alt="{{chart_alt}}">
    </div>
    <div class="chart-description">
        {{description}}
    </div>
</div>"""
        
        css = """
.chart-container {
    display: flex;
    justify-content: center;
    align-items: center;
    height: 70%;
    margin-bottom: 20px;
}

.chart-container img {
    max-width: 90%;
    max-height: 100%;
    object-fit: contain;
}

.chart-description {
    font-size: 20px;
    line-height: 1.5;
    color: var(--text-color);
    text-align: center;
    padding: 0 40px;
}
"""
        
        return html, css
        
    def _create_blank_layout(self, theme):
        """创建空白布局"""
        html = """<div class="slide-header">
    <h1 class="slide-title">{{title}}</h1>
</div>
<div class="slide-content">
    {{content}}
</div>"""
        
        css = """
.slide-content {
    font-size: 24px;
    line-height: 1.6;
    color: var(--text-color);
}
"""
        
        return html, css 
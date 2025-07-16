import os
import sys
import logging
import tempfile
import shutil
import datetime
from pathlib import Path
import time

# 配置日志
logger = logging.getLogger("ppt_engine.renderer")

class Renderer:
    """渲染引擎，负责将HTML转换为PPT"""
    
    def __init__(self, output_dir=None):
        """
        初始化渲染器
        
        Args:
            output_dir: 输出目录
        """
        self.base_path = Path(os.path.dirname(os.path.abspath(__file__)))
        self.output_dir = output_dir or os.path.join(self.base_path.parent, "uploads")
        
        # 确保输出目录存在
        os.makedirs(self.output_dir, exist_ok=True)
        
        # 初始化子系统
        self._init_browser()
        
    def _init_browser(self):
        """初始化浏览器"""
        try:
            from selenium import webdriver
            from selenium.webdriver.chrome.options import Options
            
            options = Options()
            options.add_argument('--headless')
            options.add_argument('--disable-gpu')
            options.add_argument('--no-sandbox')
            options.add_argument('--disable-dev-shm-usage')
            options.add_argument('--window-size=1280,720')
            
            try:
                from selenium.webdriver.chrome.service import Service
                from webdriver_manager.chrome import ChromeDriverManager
                
                service = Service(ChromeDriverManager().install())
                self.browser = webdriver.Chrome(service=service, options=options)
            except:
                # 直接尝试使用Chrome驱动
                self.browser = webdriver.Chrome(options=options)
                
            logger.info("浏览器初始化成功")
            self.has_browser = True
        except Exception as e:
            logger.error(f"初始化浏览器失败: {str(e)}")
            self.has_browser = False
            
    def render(self, slides):
        """
        渲染幻灯片到PPT文件
        
        Args:
            slides: 幻灯片数据列表
            
        Returns:
            output_path: 输出文件路径
        """
        # 1. 生成唯一文件名
        timestamp = datetime.datetime.now().strftime('%Y%m%d%H%M%S')
        output_filename = f"enhanced_ppt_{timestamp}.pptx"
        output_path = os.path.join(self.output_dir, output_filename)
        
        # 2. 创建临时HTML文件
        html_files = []
        temp_dir = tempfile.mkdtemp(prefix="ppt_html_")
        
        try:
            for i, slide in enumerate(slides):
                html_content = self._create_full_html(slide['html'], slide['css'])
                
                # 保存到临时文件
                file_path = os.path.join(temp_dir, f"slide_{i+1}.html")
                with open(file_path, 'w', encoding='utf-8') as f:
                    f.write(html_content)
                    
                html_files.append(file_path)
                
            # 3. 将HTML转换为PPTX
            result = self._convert_html_to_pptx(html_files, output_path)
            
            if not result:
                logger.error("转换HTML到PPTX失败")
                return None
                
            return output_path
            
        except Exception as e:
            logger.error(f"渲染幻灯片失败: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            return None
        finally:
            # 清理临时文件
            try:
                shutil.rmtree(temp_dir)
            except:
                pass
                
    def _create_full_html(self, html_content, css_content):
        """
        创建完整的HTML文档
        
        Args:
            html_content: HTML内容
            css_content: CSS样式
            
        Returns:
            full_html: 完整的HTML文档
        """
        return f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <style>
    /* 基础样式 */
    body {{
        margin: 0;
        padding: 0;
        overflow: hidden;
    }}
    
    /* 引入的CSS */
    {css_content}
    </style>
</head>
<body>
    {html_content}
</body>
</html>"""
        
    def _convert_html_to_pptx(self, html_files, output_path):
        """
        将HTML文件转换为PPTX
        
        Args:
            html_files: HTML文件路径列表
            output_path: 输出文件路径
            
        Returns:
            success: 是否成功
        """
        try:
            # 如果有浏览器，使用截图方法
            if self.has_browser:
                return self._convert_with_browser(html_files, output_path)
            else:
                # 尝试使用现有的转换功能
                try:
                    from ..html_to_ppt import convert_html_files_to_ppt
                    return convert_html_files_to_ppt(html_files, output_path)
                except:
                    logger.error("无法使用现有的HTML转PPT功能")
                    return self._convert_fallback(html_files, output_path)
        except Exception as e:
            logger.error(f"转换HTML到PPTX失败: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            return False
            
    def _convert_with_browser(self, html_files, output_path):
        """
        使用浏览器截图方法转换HTML到PPTX
        
        Args:
            html_files: HTML文件路径列表
            output_path: 输出文件路径
            
        Returns:
            success: 是否成功
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches
            
            # 创建演示文稿
            prs = Presentation()
            
            # 设置幻灯片尺寸为16:9
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            
            # 处理每个HTML文件
            for html_file in html_files:
                # 打开HTML文件
                file_url = f"file:///{os.path.abspath(html_file)}"
                self.browser.get(file_url)
                
                # 等待页面加载
                time.sleep(0.5)
                
                # 截图
                screenshot = self.browser.get_screenshot_as_png()
                
                # 添加幻灯片
                slide_layout = prs.slide_layouts[6]  # 空白布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 保存截图到临时文件
                img_path = tempfile.mktemp(suffix='.png')
                with open(img_path, 'wb') as f:
                    f.write(screenshot)
                    
                try:
                    # 添加图片到幻灯片
                    slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)
                finally:
                    # 删除临时图片
                    if os.path.exists(img_path):
                        os.remove(img_path)
                        
            # 保存演示文稿
            prs.save(output_path)
            logger.info(f"演示文稿已保存到: {output_path}")
            
            return True
        except Exception as e:
            logger.error(f"使用浏览器转换失败: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            return False
            
    def _convert_fallback(self, html_files, output_path):
        """
        备用方法，当其他方法都失败时使用
        
        Args:
            html_files: HTML文件路径列表
            output_path: 输出文件路径
            
        Returns:
            success: 是否成功
        """
        logger.warning("使用备用方法生成PPT")
        
        try:
            from pptx import Presentation
            
            # 创建简单的演示文稿
            prs = Presentation()
            
            for html_file in html_files:
                # 读取HTML内容
                with open(html_file, 'r', encoding='utf-8') as f:
                    html_content = f.read()
                    
                # 提取标题和内容
                import re
                title_match = re.search(r'<h1[^>]*>(.*?)</h1>', html_content, re.DOTALL)
                title = title_match.group(1) if title_match else "Slide"
                
                # 添加幻灯片
                slide_layout = prs.slide_layouts[1]  # 标题和内容布局
                slide = prs.slides.add_slide(slide_layout)
                
                # 设置标题
                slide.shapes.title.text = title
                
                # 添加提示文本
                content = slide.placeholders[1]
                content.text = "HTML内容无法正确渲染，请检查浏览器设置"
                
            # 保存演示文稿
            prs.save(output_path)
            logger.info(f"备用演示文稿已保存到: {output_path}")
            
            return True
        except Exception as e:
            logger.error(f"备用方法失败: {str(e)}")
            return False
            
    def close(self):
        """关闭浏览器"""
        if self.has_browser:
            try:
                self.browser.close()
                self.browser.quit()
            except:
                pass 
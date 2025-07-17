#!/usr/bin/env python
"""
HTML到PPT转换器
将填充好内容的HTML页面转换为PowerPoint演示文稿
"""

import os
import sys
import logging
import tempfile
import time
import base64
import re
import shutil
from io import BytesIO
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt

# 尝试导入Selenium
try:
    from selenium import webdriver
    from selenium.webdriver.chrome.options import Options
    from selenium.webdriver.chrome.service import Service
    HAS_SELENIUM = True
except ImportError:
    HAS_SELENIUM = False
    
# 尝试导入webdriver_manager
try:
    from webdriver_manager.chrome import ChromeDriverManager
    HAS_WEBDRIVER_MANAGER = True
except ImportError:
    HAS_WEBDRIVER_MANAGER = False

# 获取模块日志记录器
logger = logging.getLogger("ppt_engine.html_to_ppt")

class HTMLToPPTConverter:
    """HTML到PPT转换器"""
    
    def __init__(self, headless=True):
        """
        初始化转换器
        
        Args:
            headless: 是否使用无头模式
        """
        self.browser = None
        self.headless = headless
        
        # 检查环境
        if not HAS_SELENIUM:
            logger.error("缺少Selenium库，请安装: pip install selenium")
        
        if not HAS_WEBDRIVER_MANAGER:
            logger.warning("缺少webdriver_manager库，推荐安装: pip install webdriver-manager")
            
        # 初始化浏览器
        self._init_browser()
    
    def _init_browser(self):
        """初始化浏览器"""
        if not HAS_SELENIUM:
            return
            
        try:
            # 配置浏览器选项
            options = Options()
            if self.headless:
                options.add_argument('--headless')
            options.add_argument('--disable-gpu')
            options.add_argument('--no-sandbox')
            options.add_argument('--disable-dev-shm-usage')
            options.add_argument('--window-size=1280,720')
            
            # 初始化浏览器
            if HAS_WEBDRIVER_MANAGER:
                service = Service(ChromeDriverManager().install())
                self.browser = webdriver.Chrome(service=service, options=options)
            else:
                # 尝试直接找到chrome driver
                self.browser = webdriver.Chrome(options=options)
                
            logger.info("浏览器初始化成功")
        except Exception as e:
            logger.error(f"初始化浏览器失败: {str(e)}")
            self.browser = None
    
    def convert_html_files_to_ppt(self, html_files, output_path):
        """
        将多个HTML文件转换为PPT
        
        Args:
            html_files: HTML文件路径列表
            output_path: 输出的PPT文件路径
            
        Returns:
            生成的PPT文件路径
        """
        logger.info(f"开始将{len(html_files)}个HTML文件转换为PPT")
        
        # 检查HTML文件
        valid_files = [f for f in html_files if os.path.exists(f) and f.endswith('.html')]
        if not valid_files:
            logger.error("没有有效的HTML文件")
            return None
            
        # 检查浏览器是否可用
        if not self.browser:
            logger.error("浏览器不可用，无法进行转换")
            return None
            
        try:
            # 创建演示文稿
            prs = Presentation()
            
            # 调整幻灯片大小为16:9
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            
            # 逐个处理HTML文件
            for html_file in valid_files:
                logger.info(f"处理HTML文件: {html_file}")
                
                # 打开HTML文件
                file_url = f"file:///{os.path.abspath(html_file)}"
                self.browser.get(file_url)
                
                # 等待页面加载
                time.sleep(0.5)
                
                # 截图
                screenshot = self.browser.get_screenshot_as_png()
                
                # 添加到PPT
                self._add_image_to_presentation(prs, screenshot)
                
            # 保存PPT
            prs.save(output_path)
            logger.info(f"PPT生成成功: {output_path}")
            
            return output_path
        except Exception as e:
            logger.error(f"转换HTML文件失败: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            return None
    
    def convert_html_contents_to_ppt(self, html_contents, output_path):
        """
        将HTML内容列表转换为PPT
        
        Args:
            html_contents: HTML内容字符串列表
            output_path: 输出的PPT文件路径
            
        Returns:
            生成的PPT文件路径
        """
        logger.info(f"开始将{len(html_contents)}个HTML内容转换为PPT")
        
        # 检查HTML内容
        if not html_contents:
            logger.error("没有HTML内容")
            return None
            
        # 检查浏览器是否可用
        if not self.browser:
            logger.error("浏览器不可用，无法进行转换")
            return None
        
        try:
            # 创建临时目录
            temp_dir = tempfile.mkdtemp(prefix="html_to_ppt_")
            
            # 保存HTML内容到临时文件
            html_files = []
            for i, html in enumerate(html_contents):
                file_path = os.path.join(temp_dir, f"slide_{i+1}.html")
                with open(file_path, 'w', encoding='utf-8') as f:
                    # 处理资源路径
                    html = self._fix_resource_paths(html, temp_dir)
                    f.write(html)
                html_files.append(file_path)
                
            # 转换HTML文件
            result = self.convert_html_files_to_ppt(html_files, output_path)
            
            # 清理临时目录
            try:
                shutil.rmtree(temp_dir)
            except:
                pass
                
            return result
        except Exception as e:
            logger.error(f"转换HTML内容失败: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            return None
    
    def _fix_resource_paths(self, html, base_dir):
        """
        修复HTML中的资源路径，确保可以正确加载样式表和图片
        
        Args:
            html: HTML内容
            base_dir: 基础目录
            
        Returns:
            修复后的HTML内容
        """
        # 修复CSS路径
        def replace_css_link(match):
            css_file = match.group(1)
            # 如果是相对路径且不是http开头，检查文件是否存在
            if not css_file.startswith(('http://', 'https://', '/')):
                css_path = os.path.join(os.path.dirname(base_dir), css_file)
                if os.path.exists(css_path):
                    # 复制到临时目录
                    dest_path = os.path.join(base_dir, os.path.basename(css_file))
                    shutil.copy2(css_path, dest_path)
                    return f'<link rel="stylesheet" href="{os.path.basename(css_file)}">'
            return match.group(0)
            
        html = re.sub(r'<link rel="stylesheet" href="([^"]+)"', replace_css_link, html)
        
        # 修复内联样式
        html = re.sub(r'@import url\([\'"]?([^\'"]+)[\'"]?\);', lambda m: '@import url("' + m.group(1) + '");', html)
        
        # 修复图片路径
        def replace_image_src(match):
            img_src = match.group(1)
            # 处理数据URL和绝对URL
            if img_src.startswith(('data:', 'http://', 'https://', '/')):
                return match.group(0)
                
            # 尝试将相对路径复制到临时目录
            img_path = os.path.join(os.path.dirname(base_dir), img_src)
            if os.path.exists(img_path):
                dest_path = os.path.join(base_dir, os.path.basename(img_src))
                shutil.copy2(img_path, dest_path)
                return f'<img src="{os.path.basename(img_src)}"'
                
            return match.group(0)
            
        html = re.sub(r'<img src="([^"]+)"', replace_image_src, html)
        
        return html
    
    def _add_image_to_presentation(self, presentation, image_data):
        """
        将图像添加到演示文稿
        
        Args:
            presentation: 演示文稿对象
            image_data: 图像数据
            
        Returns:
            添加的幻灯片对象
        """
        try:
            # 添加空白幻灯片
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # 空白布局
            
            # 图片位置和大小
            width = presentation.slide_width
            height = presentation.slide_height
            
            # 处理不同类型的图像数据
            if isinstance(image_data, bytes):
                # 如果是图像数据
                image = Image.open(BytesIO(image_data))
                logger.info("使用二进制图像数据")
            elif isinstance(image_data, str):
                if image_data.startswith('file://'):
                    # 处理file://开头的路径，确保正确提取路径
                    file_path = image_data[7:]
                    # 在Windows上，可能需要额外处理
                    if os.name == 'nt' and file_path.startswith('/'):
                        file_path = file_path[1:]  # 移除开头的/
                    logger.info(f"从file://路径加载图像: {file_path}")
                    
                    # 检查文件是否存在
                    if not os.path.exists(file_path):
                        logger.warning(f"图像文件不存在: {file_path}")
                        # 尝试使用默认图像
                        default_img_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 
                                                      "default_images", "default.jpg")
                        if os.path.exists(default_img_path):
                            file_path = default_img_path
                            logger.info(f"使用默认图像: {file_path}")
                        else:
                            # 创建空白图像
                            image = Image.new('RGB', (800, 600), color=(255, 255, 255))
                            logger.info("创建空白图像替代")
                    else:
                        image = Image.open(file_path)
                elif image_data.startswith('data:image/'):
                    # 处理data URL
                    logger.info("从数据URL加载图像")
                    import base64
                    try:
                        # 提取base64部分
                        base64_data = image_data.split(',')[1]
                        image_bytes = base64.b64decode(base64_data)
                        image = Image.open(BytesIO(image_bytes))
                    except Exception as e:
                        logger.error(f"处理数据URL失败: {str(e)}")
                        # 创建空白图像
                        image = Image.new('RGB', (800, 600), color=(255, 255, 255))
                elif os.path.exists(image_data):
                    # 处理直接的文件路径
                    logger.info(f"从文件路径加载图像: {image_data}")
                    image = Image.open(image_data)
                else:
                    logger.warning(f"无法识别的图像数据: {image_data[:30]}...")
                    # 创建一个空白图像作为占位符
                    image = Image.new('RGB', (800, 600), color=(255, 255, 255))
            else:
                logger.error(f"不支持的图像数据类型: {type(image_data)}")
                # 创建一个空白图像
                image = Image.new('RGB', (800, 600), color=(255, 255, 255))
            
            # 调整图像比例
            img_ratio = image.width / image.height
            slide_ratio = width / height
            
            # 计算合适的图像尺寸和位置
            if img_ratio > slide_ratio:
                # 图像更宽，按宽度缩放
                img_width = width * 0.8
                img_height = img_width / img_ratio
                left = width * 0.1
                top = (height - img_height) / 2
            else:
                # 图像更高，按高度缩放
                img_height = height * 0.8
                img_width = img_height * img_ratio
                top = height * 0.1
                left = (width - img_width) / 2
            
            # 保存为临时文件
            temp_img = BytesIO()
            image.save(temp_img, format='PNG')
            temp_img.seek(0)
            
            # 添加图像到幻灯片，使用计算的位置和尺寸
            slide.shapes.add_picture(temp_img, left, top, width=img_width, height=img_height)
            
            logger.info("成功添加图片到幻灯片")
            
            return slide
        except Exception as e:
            logger.error(f"添加图片到幻灯片失败: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            
            # 尝试添加一个空白幻灯片作为替代
            try:
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                return slide
            except:
                return None
    
    def close(self):
        """关闭浏览器"""
        if self.browser:
            try:
                self.browser.quit()
                logger.info("浏览器已关闭")
            except:
                pass
            self.browser = None

def convert_html_to_ppt(html_contents, output_path):
    """
    便捷函数：将HTML内容转换为PPT
    
    Args:
        html_contents: HTML内容字符串列表
        output_path: 输出的PPT文件路径
        
    Returns:
        生成的PPT文件路径
    """
    converter = HTMLToPPTConverter()
    try:
        return converter.convert_html_contents_to_ppt(html_contents, output_path)
    finally:
        converter.close()

def convert_html_files_to_ppt(html_files, output_path):
    """
    便捷函数：将HTML文件转换为PPT
    
    Args:
        html_files: HTML文件路径列表
        output_path: 输出的PPT文件路径
        
    Returns:
        生成的PPT文件路径
    """
    converter = HTMLToPPTConverter()
    try:
        return converter.convert_html_files_to_ppt(html_files, output_path)
    finally:
        converter.close()

if __name__ == "__main__":
    # 命令行接口
    if len(sys.argv) < 3:
        print("用法: python html_to_ppt.py <HTML文件1,HTML文件2,...> <输出PPT路径>")
        sys.exit(1)
        
    # 解析HTML文件
    html_files_arg = sys.argv[1]
    output_path = sys.argv[2]
    
    # 检查是目录还是文件列表
    if os.path.isdir(html_files_arg):
        # 如果是目录，获取所有HTML文件
        html_dir = html_files_arg
        html_files = sorted([os.path.join(html_dir, f) for f in os.listdir(html_dir) if f.endswith('.html')])
    else:
        # 如果是文件列表，按逗号分隔
        html_files = [f.strip() for f in html_files_arg.split(',') if f.strip()]
        
    # 转换文件
    result = convert_html_files_to_ppt(html_files, output_path)
    
    if result:
        print(f"PPT生成成功: {result}")
    else:
        print("PPT生成失败")
        sys.exit(1) 
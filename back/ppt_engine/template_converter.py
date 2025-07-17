#!/usr/bin/env python
"""
PPT模板转换器
将PowerPoint模板转换为HTML模板格式，便于内容填充
"""

import os
import json
import logging
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 获取模块日志记录器
logger = logging.getLogger("ppt_engine.template_converter")

class PPTTemplateConverter:
    """PPT模板转换器，将PPT模板转换为HTML模板"""
    
    def __init__(self, templates_dir=None, html_templates_dir=None):
        """
        初始化转换器
        
        Args:
            templates_dir: PPT模板目录
            html_templates_dir: HTML模板输出目录
        """
        # 获取模块的基础路径
        if templates_dir is None:
            templates_dir = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "ppt_templates")
        
        if html_templates_dir is None:
            html_templates_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "html_templates")
            
        self.templates_dir = templates_dir
        self.html_templates_dir = html_templates_dir
        
        # 确保目录存在
        os.makedirs(self.templates_dir, exist_ok=True)
        os.makedirs(self.html_templates_dir, exist_ok=True)
    
    def analyze_ppt_template(self, template_path):
        """
        分析PPT模板，提取结构和元素信息
        
        Args:
            template_path: PPT模板文件路径
            
        Returns:
            模板结构信息字典
        """
        logger.info(f"分析PPT模板: {template_path}")
        
        # 加载PPT文件
        prs = Presentation(template_path)
        
        # 提取模板信息
        template_info = {
            "name": os.path.basename(template_path),
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "slides": []
        }
        
        # 分析每个幻灯片
        for i, slide in enumerate(prs.slides):
            slide_data = {
                "index": i,
                "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else f"Layout {i}",
                "elements": [],
                "suitable_for": self._determine_slide_purpose(slide, i, len(prs.slides))
            }
            
            # 分析幻灯片中的所有形状
            for shape in slide.shapes:
                element = self._analyze_shape(shape)
                if element:
                    slide_data["elements"].append(element)
                    
            template_info["slides"].append(slide_data)
        
        return template_info
    
    def _analyze_shape(self, shape):
        """
        分析PPT中的形状元素
        
        Args:
            shape: 形状对象
            
        Returns:
            元素信息字典
        """
        # 基本信息
        element = {
            "id": getattr(shape, "shape_id", 0),
            "name": shape.name,
            "type": str(shape.shape_type),
            "left": shape.left,
            "top": shape.top,
            "width": shape.width,
            "height": shape.height
        }
        
        # 分析文本框
        if shape.has_text_frame:
            element["content_type"] = "text"
            element["text"] = shape.text
            
            # 检查是否为占位符
            if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                placeholder_type = shape.placeholder_format.type
                if placeholder_type == 1:  # 标题
                    element["role"] = "title"
                elif placeholder_type == 2:  # 正文
                    element["role"] = "content"
                else:
                    element["role"] = f"placeholder_{placeholder_type}"
            
            # 提取文本样式
            element["text_style"] = self._extract_text_style(shape.text_frame)
        
        # 分析图片
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element["content_type"] = "image"
        
        # 分析表格
        elif hasattr(shape, "has_table") and shape.has_table:
            element["content_type"] = "table"
            element["rows"] = shape.table.rows._length
            element["columns"] = shape.table.columns._length
        
        # 其他形状类型
        else:
            element["content_type"] = "shape"
        
        return element
    
    def _extract_text_style(self, text_frame):
        """
        提取文本框的样式信息
        
        Args:
            text_frame: 文本框对象
            
        Returns:
            样式信息字典
        """
        style = {}
        
        # 检查是否有段落
        if not text_frame.paragraphs:
            return style
            
        # 从第一段文本获取样式信息
        p = text_frame.paragraphs[0]
        
        style["alignment"] = str(p.alignment) if hasattr(p, "alignment") else None
        style["level"] = p.level
        
        # 检查是否有文本运行
        if not p.runs:
            return style
            
        # 从第一个运行获取字体信息
        r = p.runs[0]
        if hasattr(r, "font"):
            font = r.font
            style["font_name"] = font.name
            style["font_size"] = font.size.pt if hasattr(font, "size") and font.size else None
            style["bold"] = font.bold
            style["italic"] = font.italic
            
            # 获取颜色信息
            if hasattr(font, "color") and hasattr(font.color, "rgb"):
                rgb = font.color.rgb
                if rgb:
                    style["color"] = f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"
        
        return style
    
    def _determine_slide_purpose(self, slide, index, total_slides):
        """
        确定幻灯片的用途
        
        Args:
            slide: 幻灯片对象
            index: 幻灯片索引
            total_slides: 幻灯片总数
            
        Returns:
            用途列表
        """
        purposes = []
        
        # 根据位置判断
        if index == 0:
            purposes.append("cover")
        elif index == total_slides - 1:
            purposes.append("conclusion")
        else:
            # 检查是否包含特定元素
            has_title = False
            has_content = False
            has_image = False
            has_table = False
            
            # 分析形状
            for shape in slide.shapes:
                # 检查是否为标题
                if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                    if shape.placeholder_format.type == 1:
                        has_title = True
                    elif shape.placeholder_format.type == 2:
                        has_content = True
                
                # 检查是否为图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    has_image = True
                
                # 检查是否为表格
                if hasattr(shape, "has_table") and shape.has_table:
                    has_table = True
            
            # 根据元素组合确定用途
            if has_title and has_content:
                purposes.append("content")
            if has_image:
                if has_title or has_content:
                    purposes.append("image_content")
                else:
                    purposes.append("image")
            if has_table:
                purposes.append("table")
            
            # 如果没有明确用途，标记为通用内容
            if not purposes:
                purposes.append("general")
        
        return purposes
    
    def convert_to_html_template(self, template_path, output_dir=None):
        """
        将PPT模板转换为HTML模板
        
        Args:
            template_path: PPT模板文件路径
            output_dir: 输出目录路径，如果为None则使用默认目录
            
        Returns:
            HTML模板目录路径，转换失败则返回None
        """
        try:
            # 检查模板文件是否存在
            if not os.path.exists(template_path):
                logger.error(f"模板文件不存在: {template_path}")
                return None
                
            # 获取模板名称
            template_name = os.path.splitext(os.path.basename(template_path))[0]
            
            # 确定输出目录
            if output_dir is None:
                output_dir = os.path.join(self.html_templates_dir, template_name)
            
            logger.info(f"开始转换模板: {template_path} -> {output_dir}")
            
            # 确保输出目录存在
            os.makedirs(output_dir, exist_ok=True)
            
            # 分析模板
            template_info = self.analyze_ppt_template(template_path)
            if not template_info:
                logger.error("分析模板失败")
                return None
                
            # 保存模板信息
            info_path = os.path.join(output_dir, "template_info.json")
            with open(info_path, 'w', encoding='utf-8') as f:
                json.dump(template_info, f, ensure_ascii=False, indent=2)
                
            # 为每个幻灯片创建HTML模板
            for slide_info in template_info["slides"]:
                html_content = self._create_html_template(slide_info, template_info)
                if html_content:
                    # 保存HTML文件
                    slide_index = slide_info["index"]
                    purpose = "_".join(slide_info["suitable_for"])
                    html_filename = f"slide_{slide_index + 1}_{purpose}.html"
                    html_path = os.path.join(output_dir, html_filename)
                    
                    with open(html_path, 'w', encoding='utf-8') as f:
                        f.write(html_content)
                        
            # 生成CSS文件
            css_content = self._generate_css(template_info)
            css_path = os.path.join(output_dir, "styles.css")
            with open(css_path, 'w', encoding='utf-8') as f:
                f.write(css_content)
                
            # 创建通用模板文件
            self._create_common_templates(output_dir, template_info)
            
            logger.info(f"模板转换成功: {template_path} -> {output_dir}")
            logger.info(f"生成了 {len(template_info['slides'])} 个HTML模板文件")
            
            # 返回HTML模板目录路径
            return output_dir
        except Exception as e:
            logger.error(f"转换模板失败: {str(e)}")
            import traceback
            logger.error(traceback.format_exc())
            return None
            
    def _create_html_template(self, slide_info, template_info):
        """
        为幻灯片创建HTML模板
        
        Args:
            slide_info: 幻灯片信息
            template_info: 模板信息
            
        Returns:
            HTML模板内容
        """
        template_name = template_info["name"]
        slide_width = template_info["slide_width"] / 914400  # 转换为像素
        slide_height = template_info["slide_height"] / 914400
        
        # 创建HTML头部
        html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
  <meta charset="UTF-8">
  <title>{slide_info["layout_name"]}</title>
  <meta name="viewport" content="width=device-width, initial-scale=1.0">
  <link rel="stylesheet" href="style.css">
</head>
<body>
  <div class="slide-container" data-purpose="{','.join(slide_info["suitable_for"])}">
"""
        
        # 添加元素
        for element in slide_info["elements"]:
            element_html = self._create_element_html(element)
            html += f"    {element_html}\n"
            
        # 添加HTML尾部
        html += """  </div>
</body>
</html>"""
        
        return html
    
    def _create_element_html(self, element):
        """
        为元素创建HTML代码
        
        Args:
            element: 元素信息
            
        Returns:
            HTML代码
        """
        element_id = f"element_{element['id']}"
        content_type = element.get("content_type", "shape")
        role = element.get("role", "")
        
        # 计算位置和大小（转换为像素）
        left = element["left"] / 914400
        top = element["top"] / 914400
        width = element["width"] / 914400
        height = element["height"] / 914400
        
        # 基本样式
        style = f"style=\"position:absolute; left:{left}px; top:{top}px; width:{width}px; height:{height}px;\""
        
        # 根据元素类型创建HTML
        if content_type == "text":
            classes = f"element text {role}"
            
            if role == "title":
                return f'<h1 id="{element_id}" class="{classes}" {style}>{{{{title}}}}</h1>'
            elif role == "content":
                return f'<div id="{element_id}" class="{classes}" {style}>{{{{content}}}}</div>'
            elif "placeholder" in role:
                return f'<div id="{element_id}" class="{classes}" {style}>{{{{placeholder_{element_id}}}}}</div>'
            else:
                return f'<div id="{element_id}" class="{classes}" {style}>{element.get("text", "")}</div>'
                
        elif content_type == "image":
            classes = f"element image {role}"
            return f'<div id="{element_id}" class="{classes}" {style}><img src="{{{{image_url}}}}" alt="幻灯片图片" style="width:100%; height:100%; object-fit:contain;"></div>'
            
        elif content_type == "table":
            classes = f"element table {role}"
            return f'<div id="{element_id}" class="{classes}" {style}>{{{{table_html}}}}</div>'
            
        else:
            classes = f"element {content_type} {role}"
            return f'<div id="{element_id}" class="{classes}" {style}></div>'
    
    def _generate_css(self, template_info):
        """
        生成CSS样式表
        
        Args:
            template_info: 模板信息
            
        Returns:
            CSS内容
        """
        # 提取幻灯片尺寸
        width = template_info.get("slide_width", 960)
        height = template_info.get("slide_height", 540)
        
        # 计算宽高比
        aspect_ratio = height / width
        
        # 提取主题颜色
        primary_color = "#2E7D32"  # 默认绿色
        secondary_color = "#4CAF50"
        text_color = "#333333"
        background_color = "#FFFFFF"
        
        # 尝试从模板中提取颜色
        for slide in template_info.get("slides", []):
            for element in slide.get("elements", []):
                if element.get("content_type") == "text" and element.get("text_style", {}).get("color"):
                    text_color = element.get("text_style", {}).get("color")
                    break
        
        # 生成CSS
        css = f"""
/* 全局样式 */
* {{
    margin: 0;
    padding: 0;
    box-sizing: border-box;
}}

body {{
    font-family: 'Microsoft YaHei', Arial, sans-serif;
    background-color: #f0f0f0;
    margin: 0;
    padding: 20px;
}}

/* 幻灯片容器 */
.slide {{
    width: 100%;
    max-width: 1200px;
    margin: 0 auto;
    background-color: {background_color};
    box-shadow: 0 4px 8px rgba(0, 0, 0, 0.1);
    position: relative;
    overflow: hidden;
    padding: 40px;
    aspect-ratio: {1/aspect_ratio};
}}

/* 标题样式 */
.title {{
    font-size: 36px;
    font-weight: bold;
    color: {primary_color};
    margin-bottom: 20px;
    text-align: left;
}}

/* 内容样式 */
.content {{
    font-size: 24px;
    color: {text_color};
    line-height: 1.5;
    margin-bottom: 20px;
}}

/* 子标题样式 */
.subtitle {{
    font-size: 28px;
    color: {secondary_color};
    margin-bottom: 15px;
}}

/* 项目符号列表 */
.bullet-points {{
    list-style-type: disc;
    margin-left: 30px;
    font-size: 24px;
    color: {text_color};
    line-height: 1.5;
}}

.bullet-points li {{
    margin-bottom: 10px;
    position: relative;
}}

/* 图片容器 */
.image-container {{
    text-align: center;
    margin: 20px 0;
    max-width: 100%;
    height: auto;
}}

.image-container img {{
    max-width: 100%;
    max-height: 400px;
    object-fit: contain;
}}

/* 表格样式 */
table {{
    width: 100%;
    border-collapse: collapse;
    margin: 20px 0;
}}

th, td {{
    border: 1px solid #ddd;
    padding: 12px;
    text-align: left;
}}

th {{
    background-color: {primary_color};
    color: white;
}}

tr:nth-child(even) {{
    background-color: #f2f2f2;
}}

/* 特殊幻灯片样式 */
.cover-slide {{
    text-align: center;
    display: flex;
    flex-direction: column;
    justify-content: center;
    align-items: center;
    padding: 60px;
}}

.cover-slide .title {{
    font-size: 48px;
    margin-bottom: 30px;
    text-align: center;
}}

.cover-slide .subtitle {{
    font-size: 32px;
    margin-bottom: 40px;
    text-align: center;
}}

.conclusion-slide {{
    background-color: #f9f9f9;
}}

.conclusion-slide .title {{
    color: {primary_color};
}}

/* 响应式设计 */
@media (max-width: 768px) {{
    .slide {{
        padding: 20px;
    }}
    
    .title {{
        font-size: 28px;
    }}
    
    .content, .bullet-points {{
        font-size: 18px;
    }}
    
    .cover-slide .title {{
        font-size: 36px;
    }}
    
    .cover-slide .subtitle {{
        font-size: 24px;
    }}
}}

/* 导航栏样式 - 模拟PPT左侧导航 */
.slide-nav {{
    position: absolute;
    left: 0;
    top: 0;
    width: 150px;
    height: 100%;
    background-color: #f0f0f0;
    border-right: 1px solid #ddd;
    padding: 20px 0;
}}

.slide-nav-item {{
    padding: 8px 15px;
    cursor: pointer;
    font-size: 14px;
    color: #666;
}}

.slide-nav-item.active {{
    background-color: {primary_color};
    color: white;
}}

/* 带导航的幻灯片内容区域 */
.with-nav .slide-content {{
    margin-left: 150px;
}}
"""
        return css

    def _create_common_templates(self, output_dir, template_info):
        """
        创建通用HTML模板文件（封面、内容、结论等）
        
        Args:
            output_dir: 输出目录
            template_info: 模板信息
        """
        # 创建封面模板
        cover_html = self._create_cover_template(template_info)
        with open(os.path.join(output_dir, "cover.html"), 'w', encoding='utf-8') as f:
            f.write(cover_html)
            
        # 创建内容模板
        content_html = self._create_content_template(template_info)
        with open(os.path.join(output_dir, "content.html"), 'w', encoding='utf-8') as f:
            f.write(content_html)
            
        # 创建图文模板
        image_html = self._create_image_template(template_info)
        with open(os.path.join(output_dir, "image.html"), 'w', encoding='utf-8') as f:
            f.write(image_html)
            
        # 创建表格模板
        table_html = self._create_table_template(template_info)
        with open(os.path.join(output_dir, "table.html"), 'w', encoding='utf-8') as f:
            f.write(table_html)
            
        # 创建结论模板
        conclusion_html = self._create_conclusion_template(template_info)
        with open(os.path.join(output_dir, "conclusion.html"), 'w', encoding='utf-8') as f:
            f.write(conclusion_html)
            
        logger.info(f"已创建通用模板文件: cover.html, content.html, image.html, table.html, conclusion.html")
        
    def _create_cover_template(self, template_info):
        """创建封面模板"""
        html = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>封面</title>
    <link rel="stylesheet" href="styles.css">
    <style>
        .cover-slide {
            display: flex;
            flex-direction: column;
            justify-content: center;
            align-items: center;
            text-align: center;
            height: 100%;
            padding: 40px;
        }
        .title {
            font-size: 48px;
            font-weight: bold;
            margin-bottom: 30px;
            color: #2E7D32;
        }
        .subtitle {
            font-size: 24px;
            margin-bottom: 20px;
            color: #4CAF50;
        }
        .author {
            font-size: 18px;
            margin-top: 40px;
            color: #666;
        }
        .image-container {
            margin: 30px 0;
            max-width: 80%;
        }
        .image-container img {
            max-width: 100%;
            max-height: 300px;
            object-fit: contain;
        }
    </style>
</head>
<body>
    <div class="slide cover-slide">
        <h1 class="title">{{ title }}</h1>
        <div class="subtitle">{{ content }}</div>
        {% if image_url %}
        <div class="image-container">
            <img src="{{ image_url }}" alt="{{ title }}">
        </div>
        {% endif %}
        {% if slide_data.author %}
        <div class="author">{{ slide_data.author }}</div>
        {% endif %}
    </div>
</body>
</html>"""
        return html
        
    def _create_content_template(self, template_info):
        """创建内容模板"""
        html = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>内容</title>
    <link rel="stylesheet" href="styles.css">
    <style>
        .content-slide {
            display: flex;
            flex-direction: column;
            padding: 40px;
            height: 100%;
        }
        .title {
            font-size: 36px;
            font-weight: bold;
            margin-bottom: 30px;
            color: #2E7D32;
        }
        .content {
            font-size: 24px;
            line-height: 1.5;
            margin-bottom: 20px;
            flex: 1;
        }
        .bullet-points {
            list-style-type: disc;
            margin-left: 30px;
            font-size: 24px;
            line-height: 1.5;
        }
        .bullet-points li {
            margin-bottom: 15px;
        }
        .slide-nav {
            position: absolute;
            left: 0;
            top: 0;
            width: 150px;
            height: 100%;
            background-color: #f0f0f0;
            border-right: 1px solid #ddd;
            padding: 20px 0;
        }
        .slide-content {
            margin-left: 150px;
        }
    </style>
</head>
<body>
    <div class="slide content-slide">
        <div class="slide-nav">
            <!-- 导航占位符 -->
        </div>
        <div class="slide-content">
            <h1 class="title">{{ title }}</h1>
            <div class="content">{{ content }}</div>
            {% if bullet_points and bullet_points|length > 0 %}
            <ul class="bullet-points">
                {% for point in bullet_points %}
                <li>{{ point }}</li>
                {% endfor %}
            </ul>
            {% endif %}
        </div>
    </div>
</body>
</html>"""
        return html
        
    def _create_image_template(self, template_info):
        """创建图文模板"""
        html = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>图文</title>
    <link rel="stylesheet" href="styles.css">
    <style>
        .image-slide {
            display: flex;
            flex-direction: column;
            padding: 40px;
            height: 100%;
        }
        .title {
            font-size: 36px;
            font-weight: bold;
            margin-bottom: 20px;
            color: #2E7D32;
        }
        .content-container {
            display: flex;
            flex: 1;
        }
        .text-content {
            flex: 1;
            padding-right: 20px;
        }
        .content {
            font-size: 24px;
            line-height: 1.5;
            margin-bottom: 20px;
        }
        .image-container {
            flex: 1;
            display: flex;
            justify-content: center;
            align-items: center;
        }
        .image-container img {
            max-width: 100%;
            max-height: 400px;
            object-fit: contain;
        }
        .slide-nav {
            position: absolute;
            left: 0;
            top: 0;
            width: 150px;
            height: 100%;
            background-color: #f0f0f0;
            border-right: 1px solid #ddd;
            padding: 20px 0;
        }
        .slide-content {
            margin-left: 150px;
        }
    </style>
</head>
<body>
    <div class="slide image-slide">
        <div class="slide-nav">
            <!-- 导航占位符 -->
        </div>
        <div class="slide-content">
            <h1 class="title">{{ title }}</h1>
            <div class="content-container">
                <div class="text-content">
                    <div class="content">{{ content }}</div>
                </div>
                {% if image_url %}
                <div class="image-container">
                    <img src="{{ image_url }}" alt="{{ title }}">
                </div>
                {% endif %}
            </div>
        </div>
    </div>
</body>
</html>"""
        return html
        
    def _create_conclusion_template(self, template_info):
        """创建结论模板"""
        html = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>结论</title>
    <link rel="stylesheet" href="styles.css">
    <style>
        .conclusion-slide {
            display: flex;
            flex-direction: column;
            padding: 40px;
            height: 100%;
            background-color: #f9f9f9;
        }
        .title {
            font-size: 36px;
            font-weight: bold;
            margin-bottom: 30px;
            color: #2E7D32;
            text-align: center;
        }
        .content {
            font-size: 24px;
            line-height: 1.5;
            margin: 20px auto;
            max-width: 80%;
            text-align: center;
        }
        .image-container {
            margin: 20px auto;
            text-align: center;
            max-width: 60%;
        }
        .image-container img {
            max-width: 100%;
            max-height: 300px;
            object-fit: contain;
        }
        .slide-nav {
            position: absolute;
            left: 0;
            top: 0;
            width: 150px;
            height: 100%;
            background-color: #f0f0f0;
            border-right: 1px solid #ddd;
            padding: 20px 0;
        }
        .slide-content {
            margin-left: 150px;
        }
    </style>
</head>
<body>
    <div class="slide conclusion-slide">
        <div class="slide-nav">
            <!-- 导航占位符 -->
        </div>
        <div class="slide-content">
            <h1 class="title">{{ title }}</h1>
            <div class="content">{{ content }}</div>
            {% if image_url %}
            <div class="image-container">
                <img src="{{ image_url }}" alt="{{ title }}">
            </div>
            {% endif %}
        </div>
    </div>
</body>
</html>"""
        return html

    def _create_table_template(self, template_info):
        """创建表格模板"""
        html = """<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>表格</title>
    <link rel="stylesheet" href="styles.css">
    <style>
        .table-slide {
            display: flex;
            flex-direction: column;
            padding: 40px;
            height: 100%;
        }
        .title {
            font-size: 36px;
            font-weight: bold;
            margin-bottom: 20px;
            color: #2E7D32;
        }
        .content {
            font-size: 18px;
            line-height: 1.5;
            margin-bottom: 20px;
        }
        .table-container {
            flex: 1;
            overflow: auto;
        }
        table {
            width: 100%;
            border-collapse: collapse;
            margin: 20px 0;
        }
        th {
            background-color: #2E7D32;
            color: white;
            font-weight: bold;
            text-align: center;
            padding: 12px;
            border: 1px solid #ddd;
        }
        td {
            padding: 10px;
            border: 1px solid #ddd;
            text-align: left;
        }
        tr:nth-child(even) {
            background-color: #f2f2f2;
        }
        .slide-nav {
            position: absolute;
            left: 0;
            top: 0;
            width: 150px;
            height: 100%;
            background-color: #f0f0f0;
            border-right: 1px solid #ddd;
            padding: 20px 0;
        }
        .slide-content {
            margin-left: 150px;
        }
    </style>
</head>
<body>
    <div class="slide table-slide">
        <div class="slide-nav">
            <!-- 导航占位符 -->
        </div>
        <div class="slide-content">
            <h1 class="title">{{ title }}</h1>
            <div class="content">{{ content }}</div>
            <div class="table-container">
                {% if table_html %}
                {{ table_html|safe }}
                {% else %}
                <table>
                    <thead>
                        <tr>
                            {% for header in table[0] %}
                            <th>{{ header }}</th>
                            {% endfor %}
                        </tr>
                    </thead>
                    <tbody>
                        {% for row in table[1:] %}
                        <tr>
                            {% for cell in row %}
                            <td>{{ cell }}</td>
                            {% endfor %}
                        </tr>
                        {% endfor %}
                    </tbody>
                </table>
                {% endif %}
            </div>
        </div>
    </div>
</body>
</html>"""
        return html

def convert_template(template_path, output_dir=None):
    """
    转换PPT模板为HTML模板的便捷函数
    
    Args:
        template_path: PPT模板文件路径
        output_dir: 输出目录路径
        
    Returns:
        HTML模板目录路径
    """
    converter = PPTTemplateConverter()
    return converter.convert_to_html_template(template_path, output_dir)

if __name__ == "__main__":
    # 命令行接口
    if len(sys.argv) < 2:
        print("用法: python template_converter.py <模板路径> [输出目录]")
        sys.exit(1)
        
    template_path = sys.argv[1]
    output_dir = sys.argv[2] if len(sys.argv) > 2 else None
    
    convert_template(template_path, output_dir) 
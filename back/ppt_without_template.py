#!/usr/bin/env python
"""
无模板PPT生成工具
直接生成基本的PPT，不使用模板
"""

import os
import sys
import json
import logging
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import requests
from image_service import get_image_for_slide
import tempfile
import re
import base64

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("ppt_generation.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("ppt_without_template")

# 安全的字体大小设置函数
def safe_font_size(size):
    """确保字体大小在有效范围内 (100-400000)"""
    if size <= 0:
        return Pt(12)  # 默认值
    if size < 100:
        return Pt(size)  # 正常情况，使用传入的值
    return Pt(72)  # 如果值过大，使用最大安全值

def create_title_slide(prs, slide_data):
    """创建标题页"""
    slide_layout = prs.slide_layouts[0]  # 使用标题布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title = slide.shapes.title
    title.text = slide_data.get('title', '演示文稿')
    
    # 设置标题格式
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = safe_font_size(44)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 112, 192)
    
    # 添加副标题
    if slide.placeholders[1].has_text_frame:
        subtitle = slide.placeholders[1]
        subtitle.text = slide_data.get('content', '')
        
        # 设置副标题格式
        for paragraph in subtitle.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = safe_font_size(24)
                run.font.color.rgb = RGBColor(89, 89, 89)
    
    return slide

def create_content_slide(prs, slide_data):
    """创建内容页"""
    slide_layout = prs.slide_layouts[1]  # 使用标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title = slide.shapes.title
    title.text = slide_data.get('title', '内容页')
    
    # 设置标题格式
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = safe_font_size(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 112, 192)
    
    # 添加内容
    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2:  # 2 = BODY
            content_placeholder = shape
            break
    
    if content_placeholder and content_placeholder.has_text_frame:
        tf = content_placeholder.text_frame
        tf.text = slide_data.get('content', '')
        
        # 设置内容格式
        for paragraph in tf.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = safe_font_size(24)
                run.font.color.rgb = RGBColor(0, 0, 0)
    
    return slide

def create_bullet_slide(prs, slide_data):
    """创建要点页"""
    slide_layout = prs.slide_layouts[1]  # 使用标题和内容布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title = slide.shapes.title
    title.text = slide_data.get('title', '要点页')
    
    # 设置标题格式
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = safe_font_size(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 112, 192)
    
    # 添加要点
    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.type == 2:  # 2 = BODY
            content_placeholder = shape
            break
    
    if content_placeholder and content_placeholder.has_text_frame:
        tf = content_placeholder.text_frame
        tf.clear()
        
        # 添加内容作为第一段
        if slide_data.get('content'):
            p = tf.paragraphs[0]
            p.text = slide_data.get('content')
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = safe_font_size(24)
                run.font.color.rgb = RGBColor(0, 0, 0)
            
            # 添加空行
            tf.add_paragraph()
        
        # 添加要点
        keypoints = slide_data.get('keypoints', [])
        for i, point in enumerate(keypoints):
            p = tf.add_paragraph()
            
            # 处理不同类型的要点
            if isinstance(point, dict):
                p.text = point.get('text', '')
                p.level = point.get('level', 0)
            elif isinstance(point, list):
                p.text = ' - '.join([str(item) for item in point])
            else:
                p.text = str(point)
            
            p.bullet = True
            p.alignment = PP_ALIGN.LEFT
            
            for run in p.runs:
                run.font.size = safe_font_size(20)
                run.font.color.rgb = RGBColor(0, 0, 0)
    
    return slide

def create_image_slide(prs, slide_data):
    """创建图片页"""
    slide_layout = prs.slide_layouts[5]  # 使用只有标题的布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title = slide.shapes.title
    title.text = slide_data.get('title', '图片页')
    
    # 设置标题格式
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = safe_font_size(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 112, 192)
    
    # 添加内容文本框
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = slide_data.get('content', '')
    
    # 设置内容格式
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = safe_font_size(20)
            run.font.color.rgb = RGBColor(0, 0, 0)
    
    # 添加图片
    try:
        # 获取图片
        image_source = None
        
        # 检查是否已有image_url
        if slide_data.get('image_url'):
            image_source = slide_data.get('image_url')
        # 然后检查image字段
        elif slide_data.get('image'):
            image_source = slide_data.get('image')
        
        # 如果没有找到图片源，尝试调用服务生成
        if not image_source:
            image_source = get_image_for_slide(slide_data)
        
        # 添加图片到幻灯片
        if image_source:
            left = Inches(2)
            top = Inches(2.5)
            width = Inches(6)
            height = Inches(4)
            
            # 处理不同类型的图片源
            if isinstance(image_source, bytes):
                # 如果是二进制数据，直接使用
                slide.shapes.add_picture(BytesIO(image_source), left, top, width, height)
                
            elif isinstance(image_source, str):
                # 处理字符串类型的图片源
                if image_source.startswith('file://'):
                    # 本地文件路径，去掉file://前缀
                    file_path = image_source[7:]
                    if os.path.exists(file_path):
                        try:
                            # 使用Pillow打开图片，然后保存为内存缓冲区
                            from PIL import Image
                            img = Image.open(file_path)
                            
                            # 保存为内存缓冲区
                            img_io = BytesIO()
                            img_format = img.format if img.format else 'PNG'
                            img.save(img_io, format=img_format)
                            img_io.seek(0)
                            
                            # 添加到幻灯片
                            slide.shapes.add_picture(img_io, left, top, width, height)
                        except Exception as e:
                            logger.error(f"处理本地图片失败: {str(e)}")
                            raise e
                    else:
                        logger.error(f"图片文件不存在: {file_path}")
                        raise FileNotFoundError(f"图片文件不存在: {file_path}")
                        
                elif image_source.startswith(('http://', 'https://')):
                    # 网络URL，下载图片
                    response = requests.get(image_source, timeout=10)
                    response.raise_for_status()
                    slide.shapes.add_picture(BytesIO(response.content), left, top, width, height)
                    
                elif os.path.exists(image_source):
                    # 直接使用本地文件路径
                    slide.shapes.add_picture(image_source, left, top, width, height)
                    
                else:
                    # 未知格式，可能是图片描述，尝试生成图片
                    image_data = get_image_for_slide(slide_data)
                    if isinstance(image_data, bytes):
                        slide.shapes.add_picture(BytesIO(image_data), left, top, width, height)
                    elif isinstance(image_data, str) and os.path.exists(image_data.replace('file://', '')):
                        file_path = image_data.replace('file://', '')
                        slide.shapes.add_picture(file_path, left, top, width, height)
                    else:
                        raise ValueError(f"无法处理的图片格式: {type(image_data)}")
            else:
                logger.error(f"无法处理的图片源类型: {type(image_source)}")
                raise TypeError(f"无法处理的图片源类型: {type(image_source)}")
                
        else:
            logger.warning("未找到图片源")
            # 添加提示文本
            left = Inches(2)
            top = Inches(3.5)
            width = Inches(6)
            height = Inches(1)
            error_box = slide.shapes.add_textbox(left, top, width, height)
            error_box.text_frame.text = "未提供图片"
            for paragraph in error_box.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = safe_font_size(20)
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(192, 0, 0)
    except Exception as e:
        logger.error(f"添加图片失败: {str(e)}")
        # 添加错误提示
        left = Inches(2)
        top = Inches(3.5)
        width = Inches(6)
        height = Inches(1)
        error_box = slide.shapes.add_textbox(left, top, width, height)
        error_box.text_frame.text = f"图片加载失败: {str(e)}"
        for paragraph in error_box.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = safe_font_size(20)
                run.font.italic = True
                run.font.color.rgb = RGBColor(192, 0, 0)
    
    return slide

def create_table_slide(prs, slide_data):
    """创建表格页"""
    slide_layout = prs.slide_layouts[5]  # 使用只有标题的布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title = slide.shapes.title
    title.text = slide_data.get('title', '表格页')
    
    # 设置标题格式
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = safe_font_size(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 112, 192)
    
    # 添加内容文本框
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(1)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = slide_data.get('content', '')
    
    # 设置内容格式
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.LEFT
        for run in paragraph.runs:
            run.font.size = safe_font_size(20)
            run.font.color.rgb = RGBColor(0, 0, 0)
    
    # 添加表格
    table_data = slide_data.get('table')
    if table_data:
        try:
            # 确定表格行列数
            if isinstance(table_data, list):
                rows = len(table_data)
                cols = max(len(row) if isinstance(row, list) else 1 for row in table_data)
                
                if rows > 0 and cols > 0:
                    left = Inches(1)
                    top = Inches(2.5)
                    width = Inches(8)
                    height = Inches(4)
                    
                    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
                    
                    # 填充表格数据
                    for i, row_data in enumerate(table_data):
                        if isinstance(row_data, list):
                            for j, cell_data in enumerate(row_data):
                                if j < cols:
                                    cell = table.cell(i, j)
                                    cell.text = str(cell_data)
                                    
                                    # 设置单元格格式
                                    for paragraph in cell.text_frame.paragraphs:
                                        paragraph.alignment = PP_ALIGN.CENTER
                                        for run in paragraph.runs:
                                            run.font.size = safe_font_size(16)
                        else:
                            cell = table.cell(i, 0)
                            cell.text = str(row_data)
                            
                            # 设置单元格格式
                            for paragraph in cell.text_frame.paragraphs:
                                paragraph.alignment = PP_ALIGN.CENTER
                                for run in paragraph.runs:
                                    run.font.size = safe_font_size(16)
        except Exception as e:
            logger.error(f"添加表格失败: {str(e)}")
            # 添加错误提示
            left = Inches(2)
            top = Inches(3.5)
            width = Inches(6)
            height = Inches(1)
            error_box = slide.shapes.add_textbox(left, top, width, height)
            error_box.text_frame.text = "表格创建失败"
            for paragraph in error_box.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = safe_font_size(20)
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(192, 0, 0)
    
    return slide

def create_summary_slide(prs, slide_data):
    """创建总结页"""
    slide_layout = prs.slide_layouts[5]  # 使用只有标题的布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title = slide.shapes.title
    title.text = slide_data.get('title', '总结')
    
    # 设置标题格式
    for paragraph in title.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = safe_font_size(44)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 112, 192)
    
    # 添加内容文本框
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = slide_data.get('content', '')
    
    # 设置内容格式
    for paragraph in tf.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.size = safe_font_size(32)
            run.font.color.rgb = RGBColor(0, 0, 0)
    
    # 添加要点
    keypoints = slide_data.get('keypoints', [])
    if keypoints:
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(3)
        
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        
        for i, point in enumerate(keypoints):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            # 处理不同类型的要点
            if isinstance(point, dict):
                p.text = point.get('text', '')
                p.level = point.get('level', 0)
            elif isinstance(point, list):
                p.text = ' - '.join([str(item) for item in point])
            else:
                p.text = str(point)
            
            p.bullet = True
            p.alignment = PP_ALIGN.LEFT
            
            for run in p.runs:
                run.font.size = safe_font_size(24)
                run.font.color.rgb = RGBColor(0, 0, 0)
    
    return slide

def generate_ppt_without_template(slides_data, output_path):
    """
    不使用模板直接生成PPT
    
    Args:
        slides_data: 幻灯片数据列表
        output_path: 输出文件路径
        
    Returns:
        bool: 成功返回True，失败返回False
    """
    try:
        logger.info(f"开始创建无模板PPT: {output_path}")
        logger.info(f"幻灯片数量: {len(slides_data)}")
        
        # 创建演示文稿
        prs = Presentation()
        
        # 为每个幻灯片数据创建一页
        for i, slide_data in enumerate(slides_data):
            logger.info(f"处理第 {i+1} 张幻灯片")
            
            # 获取幻灯片类型和布局
            slide_type = slide_data.get('type', '').lower()
            slide_layout = slide_data.get('layout', '').lower()
            
            # 更智能地判断幻灯片类型
            # 首页默认为封面
            if i == 0 or slide_type == 'cover' or slide_layout == 'cover' or "封面" in slide_data.get('title', ''):
                logger.info(f"创建封面页: {slide_data.get('title', '未命名')}")
                create_title_slide(prs, slide_data)
                
            # 最后一页默认为总结
            elif i == len(slides_data) - 1 or slide_type == 'summary' or slide_type == 'conclusion' or slide_layout == 'summary' or "总结" in slide_data.get('title', '') or "结论" in slide_data.get('title', ''):
                logger.info(f"创建总结页: {slide_data.get('title', '未命名')}")
                create_summary_slide(prs, slide_data)
                
            # 如果有表格数据，创建表格页
            elif 'table' in slide_data and slide_data['table']:
                logger.info(f"创建表格页: {slide_data.get('title', '未命名')}")
                create_table_slide(prs, slide_data)
                
            # 如果有图片数据或图片关键词，创建图片页
            elif ('image' in slide_data and slide_data['image']) or slide_data.get('layout', '').lower() == 'image' or any(keyword in slide_data.get('title', '').lower() for keyword in ['图片', '图示', '示意图', 'image', 'picture', 'figure']):
                logger.info(f"创建图片页: {slide_data.get('title', '未命名')}")
                create_image_slide(prs, slide_data)
                
            # 如果有要点数据或要点关键词，创建要点页
            elif ('keypoints' in slide_data and slide_data['keypoints']) or any(keyword in slide_data.get('title', '').lower() for keyword in ['要点', '关键点', '重点', 'key points', 'bullet points']):
                logger.info(f"创建要点页: {slide_data.get('title', '未命名')}")
                create_bullet_slide(prs, slide_data)
                
            # 如果内容较长，创建内容页
            elif slide_data.get('content') and len(slide_data.get('content', '')) > 100:
                logger.info(f"创建内容页(长文本): {slide_data.get('title', '未命名')}")
                create_content_slide(prs, slide_data)
                
            # 默认创建内容页
            else:
                logger.info(f"创建内容页(默认): {slide_data.get('title', '未命名')}")
                create_content_slide(prs, slide_data)
        
        # 保存演示文稿
        logger.info(f"保存PPT到: {output_path}")
        prs.save(output_path)
        logger.info("PPT保存成功")
        
        return True
    except Exception as e:
        logger.error(f"创建PPT失败: {str(e)}")
        logger.error(traceback.format_exc())
        return False

if __name__ == "__main__":
    # 命令行使用方式
    if len(sys.argv) < 3:
        print("用法: python ppt_without_template.py <slides_data.json> <output.pptx>")
        sys.exit(1)
        
    slides_data_path = sys.argv[1]
    output_path = sys.argv[2]
    
    try:
        # 加载幻灯片数据
        with open(slides_data_path, 'r', encoding='utf-8') as f:
                slides_data = json.load(f)
        
        # 生成PPT
        success = generate_ppt_without_template(slides_data, output_path)
        
        if success:
            print(f"PPT生成成功: {output_path}")
        else:
            print("PPT生成失败")
            sys.exit(1)
    except Exception as e:
        print(f"错误: {e}")
        sys.exit(1) 
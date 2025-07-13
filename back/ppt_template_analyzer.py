#!/usr/bin/env python
"""
PPT模板分析工具
分析PPT模板并生成对应的JSON元数据文件，便于AI进行模板填充
"""

import os
import sys
import json
import logging
import traceback
import glob
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("ppt_template_analyzer.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("ppt_template_analyzer")

def analyze_template(template_path, output_path=None):
    """
    分析PPT模板并生成JSON元数据文件
    
    Args:
        template_path: 模板文件路径
        output_path: 输出JSON文件路径，如果为None则使用模板名称
        
    Returns:
        生成的JSON文件路径
    """
    logger.info(f"开始分析模板: {template_path}")
    
    # 确定输出路径
    if output_path is None:
        template_name = os.path.basename(template_path)
        base_name, _ = os.path.splitext(template_name)
        template_dir = os.path.dirname(template_path)
        output_path = os.path.join(template_dir, f"{base_name}.json")
    
    try:
        # 加载演示文稿
        prs = Presentation(template_path)
        
        # 创建模板元数据
        template_data = {
            "name": os.path.basename(template_path),
            "slide_count": len(prs.slides),
            "slide_width": prs.slide_width,
            "slide_height": prs.slide_height,
            "slides": []
        }
        
        # 分析每个幻灯片
        for i, slide in enumerate(prs.slides):
            slide_data = {
                "index": i,
                "layout_name": slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else f"Layout {i}",
                "shapes": [],
                "placeholders": [],
                "has_title": False,
                "has_content": False,
                "has_image": False,
                "has_table": False,
                "has_chart": False,
                "suitable_for": []  # 适合的内容类型
            }
            
            # 分析幻灯片中的形状
            for shape in slide.shapes:
                shape_data = {
                    "type": str(shape.shape_type),
                    "name": shape.name,
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height
                }
                
                # 分析文本框
                if shape.has_text_frame:
                    shape_data["has_text"] = True
                    shape_data["text"] = shape.text
                    
                    # 检查是否为标题
                    if shape.is_placeholder and shape.placeholder_format.type == 1:  # 1 = TITLE
                        slide_data["has_title"] = True
                        shape_data["is_title"] = True
                    
                    # 检查是否为内容
                    if shape.is_placeholder and shape.placeholder_format.type == 2:  # 2 = BODY
                        slide_data["has_content"] = True
                        shape_data["is_content"] = True
                
                # 分析表格
                if shape.has_table:
                    slide_data["has_table"] = True
                    shape_data["is_table"] = True
                    shape_data["rows"] = shape.table.rows._length
                    shape_data["columns"] = shape.table.columns._length
                
                # 分析图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_data["has_image"] = True
                    shape_data["is_image"] = True
                
                # 分析图表
                if shape.has_chart:
                    slide_data["has_chart"] = True
                    shape_data["is_chart"] = True
                    shape_data["chart_type"] = str(shape.chart.chart_type)
                
                # 添加形状数据
                slide_data["shapes"].append(shape_data)
            
            # 分析占位符
            for placeholder in slide.placeholders:
                placeholder_data = {
                    "idx": placeholder.placeholder_format.idx,
                    "type": placeholder.placeholder_format.type,
                    "name": placeholder.name,
                    "left": placeholder.left,
                    "top": placeholder.top,
                    "width": placeholder.width,
                    "height": placeholder.height,
                    "has_text": placeholder.has_text_frame
                }
                
                if placeholder.has_text_frame:
                    placeholder_data["text"] = placeholder.text
                
                slide_data["placeholders"].append(placeholder_data)
            
            # 确定幻灯片适合的内容类型
            if i == 0:
                slide_data["suitable_for"].append("cover")
            elif i == len(prs.slides) - 1:
                slide_data["suitable_for"].append("conclusion")
                slide_data["suitable_for"].append("summary")
            else:
                if slide_data["has_title"] and slide_data["has_content"]:
                    slide_data["suitable_for"].append("content")
                    slide_data["suitable_for"].append("keypoints")
                if slide_data["has_image"] or any(p["type"] == 18 for p in slide_data["placeholders"]):  # 18 = PICTURE
                    slide_data["suitable_for"].append("image")
                if slide_data["has_table"] or any(p["type"] == 12 for p in slide_data["placeholders"]):  # 12 = TABLE
                    slide_data["suitable_for"].append("table")
                if slide_data["has_chart"] or any(p["type"] == 13 for p in slide_data["placeholders"]):  # 13 = CHART
                    slide_data["suitable_for"].append("chart")
            
            # 添加幻灯片数据
            template_data["slides"].append(slide_data)
        
        # 添加模板适用性信息
        template_data["suitable_for"] = {
            "cover": any(s["suitable_for"] and "cover" in s["suitable_for"] for s in template_data["slides"]),
            "content": any(s["suitable_for"] and "content" in s["suitable_for"] for s in template_data["slides"]),
            "keypoints": any(s["suitable_for"] and "keypoints" in s["suitable_for"] for s in template_data["slides"]),
            "image": any(s["suitable_for"] and "image" in s["suitable_for"] for s in template_data["slides"]),
            "table": any(s["suitable_for"] and "table" in s["suitable_for"] for s in template_data["slides"]),
            "chart": any(s["suitable_for"] and "chart" in s["suitable_for"] for s in template_data["slides"]),
            "summary": any(s["suitable_for"] and "summary" in s["suitable_for"] for s in template_data["slides"])
        }
        
        # 保存JSON文件
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(template_data, f, ensure_ascii=False, indent=2)
        
        logger.info(f"模板分析完成，元数据已保存到: {output_path}")
        return output_path
    
    except Exception as e:
        logger.error(f"分析模板失败: {str(e)}")
        logger.error(traceback.format_exc())
        raise

def analyze_all_templates(template_dir):
    """
    分析指定目录中的所有PPT模板
        
        Args:
        template_dir: 模板目录路径
    """
    logger.info(f"开始分析目录中的所有模板: {template_dir}")
    
    # 确保模板目录存在
    if not os.path.isdir(template_dir):
        logger.error(f"模板目录不存在: {template_dir}")
        return
    
    # 查找所有PPTX文件
    pptx_files = glob.glob(os.path.join(template_dir, "*.pptx"))
    
    if not pptx_files:
        logger.warning(f"未找到PPTX文件: {template_dir}")
        return
    
    logger.info(f"找到{len(pptx_files)}个PPTX文件")
    
    # 分析每个模板
    success_count = 0
    for pptx_file in pptx_files:
        try:
            # 跳过临时文件
            if os.path.basename(pptx_file).startswith("~$"):
                continue
                
            logger.info(f"处理模板: {pptx_file}")
            json_path = analyze_template(pptx_file)
            logger.info(f"生成JSON元数据: {json_path}")
            success_count += 1
        except Exception as e:
            logger.error(f"处理模板失败: {pptx_file}, 错误: {str(e)}")
    
    logger.info(f"模板分析完成: 成功 {success_count}/{len(pptx_files)}")

def main():
    """命令行入口"""
    if len(sys.argv) < 2:
        # 默认模板目录
        template_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "ppt_templates")
        logger.info(f"未指定模板路径，使用默认目录: {template_dir}")
        analyze_all_templates(template_dir)
    else:
        path = sys.argv[1]
        if os.path.isdir(path):
            # 分析目录中的所有模板
            analyze_all_templates(path)
        elif os.path.isfile(path) and path.lower().endswith(".pptx"):
            # 分析单个模板
            output_path = sys.argv[2] if len(sys.argv) > 2 else None
            analyze_template(path, output_path)
        else:
            logger.error(f"无效的路径: {path}")
            print("用法: python ppt_template_analyzer.py <template.pptx|template_dir> [output.json]")
        sys.exit(1)
        
if __name__ == "__main__":
    main() 
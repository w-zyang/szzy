#!/usr/bin/env python
"""
改进版PPT生成器
使用插件式架构和更灵活的组件化设计，便于AI修改和扩展
"""

import os
import sys
import json
import logging
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from io import BytesIO
import requests
from abc import ABC, abstractmethod
import importlib
import inspect

# 尝试导入其他可能有用的库
try:
    from PIL import Image, ImageEnhance, ImageFilter
    HAS_PIL = True
except ImportError:
    HAS_PIL = False

try:
    import numpy as np
    HAS_NUMPY = True
except ImportError:
    HAS_NUMPY = False

try:
    import matplotlib.pyplot as plt
    HAS_MATPLOTLIB = True
except ImportError:
    HAS_MATPLOTLIB = False

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("ppt_generation.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("improved_ppt_generator")

# 配置常量
DEFAULT_IMAGES_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'default_images')

class SlideComponentBase(ABC):
    """幻灯片组件的基类，定义组件接口"""
    
    @abstractmethod
    def apply(self, slide, data, context=None):
        """
        应用组件到幻灯片
        
        Args:
            slide: 幻灯片对象
            data: 组件数据
            context: 上下文信息
            
        Returns:
            处理后的幻灯片
        """
        pass

class TitleComponent(SlideComponentBase):
    """标题组件"""
    
    def apply(self, slide, data, context=None):
        """应用标题组件"""
        if not data.get('title'):
            return slide
            
        # 查找标题占位符
        title = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 1:  # 1 = TITLE
                title = shape
                break
                
        # 如果找不到标题占位符，创建一个文本框作为标题
        if not title:
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1.2)
            title = slide.shapes.add_textbox(left, top, width, height)
            
        # 设置标题文本
        title.text_frame.text = data['title']
        
        # 设置标题格式
        for paragraph in title.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(context.get('title_font_size', 36))
                run.font.bold = True
                
                # 使用主题颜色或默认颜色
                if context and context.get('theme_color'):
                    run.font.color.rgb = context['theme_color']
                else:
                    run.font.color.rgb = RGBColor(0, 112, 192)
                    
        return slide

class ContentComponent(SlideComponentBase):
    """内容组件"""
    
    def apply(self, slide, data, context=None):
        """应用内容组件"""
        if not data.get('content'):
            return slide
            
        # 查找内容占位符
        content_shape = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 2:  # 2 = BODY
                content_shape = shape
                break
                
        # 如果找不到内容占位符，创建一个文本框作为内容
        if not content_shape:
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(9)
            height = Inches(5)
            content_shape = slide.shapes.add_textbox(left, top, width, height)
            
        # 设置内容文本
        tf = content_shape.text_frame
        tf.text = data['content']
        
        # 自动调整文本框大小
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        # 设置内容格式
        for paragraph in tf.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(context.get('content_font_size', 18))
                
                # 使用主题颜色或默认颜色
                if context and context.get('text_color'):
                    run.font.color.rgb = context['text_color']
                else:
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    
        return slide

class BulletPointsComponent(SlideComponentBase):
    """项目符号组件"""
    
    def apply(self, slide, data, context=None):
        """应用项目符号组件"""
        if not data.get('keypoints') and not data.get('bullet_points'):
            return slide
            
        # 使用keypoints或bullet_points
        bullet_points = data.get('keypoints', data.get('bullet_points', []))
        if not bullet_points:
            return slide
            
        # 查找内容占位符
        content_shape = None
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type == 2:  # 2 = BODY
                content_shape = shape
                break
                
        # 如果找不到内容占位符，创建一个文本框
        if not content_shape:
            left = Inches(0.5)
            top = Inches(2)
            width = Inches(9)
            height = Inches(5)
            content_shape = slide.shapes.add_textbox(left, top, width, height)
            
        # 获取文本框
        tf = content_shape.text_frame
        tf.clear()
        
        # 添加项目符号
        for i, point in enumerate(bullet_points):
            # 添加段落
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            
            # 处理不同类型的要点
            if isinstance(point, dict):
                p.text = point.get('text', '')
                p.level = point.get('level', 0)
            elif isinstance(point, list):
                p.text = ' - '.join([str(item) for item in point])
            else:
                p.text = str(point)
                
            # 设置项目符号格式
            p.bullet = True
            p.alignment = PP_ALIGN.LEFT
            
            for run in p.runs:
                run.font.size = Pt(context.get('bullet_font_size', 18))
                if context and context.get('text_color'):
                    run.font.color.rgb = context['text_color']
                else:
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    
        return slide

class ImageComponent(SlideComponentBase):
    """图片组件"""
    
    def apply(self, slide, data, context=None):
        """应用图片组件"""
        logger.info(f"处理图片组件，数据: {data.get('title', '')}")
            
        try:
            # 查找图片占位符
            image_placeholder = None
            for shape in slide.shapes:
                if shape.is_placeholder and shape.placeholder_format.type == 18:  # 18 = PICTURE
                    image_placeholder = shape
                    break
                    
            # 确定图片数据
            image_data = None
            image_path = None
            image_url = None
            
            # 如果提供了图片URL或数据
            if data.get('image'):
                image_url = data['image']
                logger.info(f"幻灯片提供了图片信息: {image_url[:100] if isinstance(image_url, str) else type(image_url)}")
                
                # 处理不同类型的图片数据
                if isinstance(image_url, str):
                    if image_url.startswith('http://') or image_url.startswith('https://'):
                    # 从URL下载图片
                        try:
                            logger.info(f"尝试下载图片URL: {image_url}")
                            response = requests.get(image_url, stream=True, timeout=10)
                    response.raise_for_status()
                    image_data = response.content
                            logger.info(f"成功下载图片，大小: {len(image_data)} 字节")
                        except Exception as e:
                            logger.warning(f"下载图片失败: {str(e)}")
                    elif image_url.startswith('/'):
                        # 本地路径，构建完整路径
                        local_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 
                                                 image_url.lstrip('/'))
                        logger.info(f"处理本地图片路径: {local_path}")
                        if os.path.exists(local_path):
                            try:
                                with open(local_path, 'rb') as f:
                        image_data = f.read()
                                logger.info(f"成功加载本地图片，大小: {len(image_data)} 字节")
                            except Exception as e:
                                logger.warning(f"读取本地图片失败: {str(e)}")
                        else:
                            logger.warning(f"本地图片路径不存在: {local_path}")
                    else:
                        # 可能是图片描述而非URL，尝试使用图片服务生成
                        logger.info(f"检测到图片描述，尝试生成图片: {image_url}")
            
            # 如果没有有效的图片数据，尝试使用图片服务
            if not image_data and context and context.get('image_service'):
                image_service = context['image_service']
                logger.info("使用图片服务生成图片")
                
                # 构建图片描述
                if isinstance(data.get('image'), str) and not data['image'].startswith(('http://', 'https://', '/')):
                    # 如果已有图片描述，直接使用
                    image_prompt = data['image']
                else:
                    # 否则从幻灯片内容构建描述
                    image_prompt = f"{data.get('title', '')}"
                    if data.get('content'):
                        image_prompt += f" - {data.get('content', '')[:100]}"
                    if data.get('keypoints'):
                        points_str = ", ".join([str(p) for p in data.get('keypoints', [])[:3]])
                        image_prompt += f" - {points_str}"
                
                logger.info(f"图片生成提示词: {image_prompt}")
                
                # 调用图片服务获取图片
                try:
                    # 直接使用图片服务的图片生成函数
                    image_url = image_service.generate_image(image_prompt)
                    logger.info(f"图片生成成功，URL: {image_url}")
                    
                    # 从URL加载图片数据
                    if image_url:
                        if image_url.startswith('http'):
                            # 网络URL
                            response = requests.get(image_url, timeout=10)
                            response.raise_for_status()
                            image_data = response.content
                        elif image_url.startswith('/'):
                            # 本地路径
                            local_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 
                                                     image_url.lstrip('/'))
                            if os.path.exists(local_path):
                                with open(local_path, 'rb') as f:
                                    image_data = f.read()
                except Exception as e:
                    logger.error(f"图片生成/获取失败: {str(e)}")
            
            # 如果还是没有图片数据，尝试使用默认图片
            if not image_data:
                logger.warning("无法获取或生成图片，使用默认图片")
                default_image_path = self._get_default_image(data)
                if default_image_path and os.path.exists(default_image_path):
                    with open(default_image_path, 'rb') as f:
                        image_data = f.read()
                else:
                    logger.error(f"默认图片不存在: {default_image_path}")
            
            # 如果有图片数据，添加到幻灯片
            if image_data:
                # 确定图片位置和大小
                if image_placeholder:
                    left = image_placeholder.left
                    top = image_placeholder.top
                    width = image_placeholder.width
                    height = image_placeholder.height
                    # 删除占位符
                    sp = image_placeholder.element
                    sp.getparent().remove(sp)
                else:
                    # 默认位置和大小
                    left = Inches(1)
                    top = Inches(2.5)
                    width = Inches(8)
                    height = Inches(4.5)
                
                # 添加图片
                slide.shapes.add_picture(BytesIO(image_data), left, top, width, height)
                logger.info("图片已成功添加到幻灯片")
            else:
                logger.error("无法获取有效的图片数据")
            
        except Exception as e:
            logger.error(f"添加图片失败: {str(e)}")
            logger.error(traceback.format_exc())
            
            # 添加错误提示
            left = Inches(2)
            top = Inches(3.5)
            width = Inches(6)
            height = Inches(1)
            error_box = slide.shapes.add_textbox(left, top, width, height)
            error_box.text_frame.text = "图片加载失败"
            for paragraph in error_box.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(20)
                    run.font.italic = True
                    run.font.color.rgb = RGBColor(192, 0, 0)
                    
        return slide
        
    def _get_default_image(self, slide_data):
        """获取适合幻灯片内容的默认图片"""
        # 默认图片目录
        default_images_dir = DEFAULT_IMAGES_DIR
        if not os.path.exists(default_images_dir):
            return None
        
        # 根据幻灯片类型选择默认图片
        slide_type = slide_data.get('type', '').lower()
        
        # 映射幻灯片类型到默认图片
        type_to_image = {
            'cover': 'cover.jpg',
            'title': 'cover.jpg',
            'conclusion': 'conclusion.jpg',
            'summary': 'summary.jpg'
        }
        
        # 检查标题关键词
        title = slide_data.get('title', '').lower()
        keywords_map = {
            'biology': ['生物', '细胞', '植物', '动物'],
            'cell': ['细胞', '生物结构'],
            'plant_cell': ['植物细胞', '植物'],
            'animal_cell': ['动物细胞', '动物']
        }
        
        # 根据标题关键词选择图片
        selected_image = None
        for image_key, keywords in keywords_map.items():
            for keyword in keywords:
                if keyword in title:
                    selected_image = f"{image_key}.jpg"
                    break
            if selected_image:
                break
        
        # 如果根据关键词没找到，使用类型映射
        if not selected_image and slide_type in type_to_image:
            selected_image = type_to_image[slide_type]
        
        # 如果都没有匹配，使用通用默认图片
        if not selected_image:
            selected_image = 'default.jpg'
        
        # 构建完整路径
        image_path = os.path.join(default_images_dir, selected_image)
        if os.path.exists(image_path):
            return image_path
            
        # 如果指定图片不存在，使用任何可用的图片
        images = [f for f in os.listdir(default_images_dir) 
                 if f.lower().endswith(('.jpg', '.jpeg', '.png'))]
        if images:
            return os.path.join(default_images_dir, images[0])
                
        return None

class TableComponent(SlideComponentBase):
    """表格组件"""
    
    def apply(self, slide, data, context=None):
        """应用表格组件"""
        if not data.get('table'):
            return slide
            
        table_data = data['table']
        
        # 确保table_data是一个二维数组
        if not isinstance(table_data, list):
            return slide
            
        # 计算表格维度
        rows = len(table_data)
        cols = max([len(row) if isinstance(row, list) else 1 for row in table_data])
        
        if rows == 0 or cols == 0:
            return slide
            
        # 设置表格位置和大小
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        
        # 添加表格
        table = slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # 填充表格数据
        for i, row_data in enumerate(table_data):
            if not isinstance(row_data, list):
                row_data = [row_data]
                
            for j, cell_data in enumerate(row_data):
                if j < cols:
                    cell = table.cell(i, j)
                    cell.text = str(cell_data)
                    
                    # 设置单元格格式
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.size = Pt(context.get('table_font_size', 14))
        
        # 设置表头样式（如果有）
        if rows > 0:
            for j in range(cols):
                header_cell = table.cell(0, j)
                for paragraph in header_cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        if context and context.get('theme_color'):
                            run.font.color.rgb = context['theme_color']
                            
        return slide

class ChartComponent(SlideComponentBase):
    """图表组件"""
    
    def apply(self, slide, data, context=None):
        """应用图表组件"""
        # 需要matplotlib支持
        if not data.get('chart') or not HAS_MATPLOTLIB:
            return slide
            
        try:
            chart_data = data['chart']
            chart_type = chart_data.get('type', 'bar')
            
            # 生成图表图像
            plt.figure(figsize=(10, 6))
            
            if chart_type == 'bar':
                self._create_bar_chart(chart_data)
            elif chart_type == 'line':
                self._create_line_chart(chart_data)
            elif chart_type == 'pie':
                self._create_pie_chart(chart_data)
            else:
                # 默认为柱状图
                self._create_bar_chart(chart_data)
                
            # 保存图表到内存
            img_data = BytesIO()
            plt.savefig(img_data, format='png')
            img_data.seek(0)
            
            # 设置图表位置
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4.5)
            
            # 添加图表图像
            slide.shapes.add_picture(img_data, left, top, width, height)
            
            # 关闭matplotlib图表
            plt.close()
            
        except Exception as e:
            logger.error(f"添加图表失败: {str(e)}")
            logger.error(traceback.format_exc())
            
        return slide
    
    def _create_bar_chart(self, chart_data):
        """创建柱状图"""
        labels = chart_data.get('labels', [])
        values = chart_data.get('values', [])
        title = chart_data.get('title', '柱状图')
        
        plt.bar(labels, values)
        plt.title(title)
        plt.tight_layout()
    
    def _create_line_chart(self, chart_data):
        """创建折线图"""
        labels = chart_data.get('labels', [])
        values = chart_data.get('values', [])
        title = chart_data.get('title', '折线图')
        
        plt.plot(labels, values)
        plt.title(title)
        plt.tight_layout()
    
    def _create_pie_chart(self, chart_data):
        """创建饼图"""
        labels = chart_data.get('labels', [])
        values = chart_data.get('values', [])
        title = chart_data.get('title', '饼图')
        
        plt.pie(values, labels=labels, autopct='%1.1f%%')
        plt.title(title)
        plt.axis('equal')
        plt.tight_layout()

class PPTGenerator:
    """改进的PPT生成器，使用组件化架构"""
    
    def __init__(self, template_path=None):
        """
        初始化PPT生成器
        
        Args:
            template_path: 模板文件路径，如果不提供则创建空白演示文稿
        """
        self.template_path = template_path
        self.prs = None
        self.slides_data = []
        self.template_layouts = {}  # 保存模板布局信息
        self.theme = {}  # 演示文稿主题配置
        self.components = {}  # 组件注册表
        
        # 初始化演示文稿
        self.init_presentation()
        
        # 注册默认组件
        self.register_component('title', TitleComponent())
        self.register_component('content', ContentComponent())
        self.register_component('bullet_points', BulletPointsComponent())
        self.register_component('image', ImageComponent())
        self.register_component('table', TableComponent())
        self.register_component('chart', ChartComponent())
        
    def init_presentation(self):
        """初始化演示文稿"""
        try:
            if self.template_path and os.path.exists(self.template_path):
                logger.info(f"使用模板创建演示文稿: {self.template_path}")
                self.prs = Presentation(self.template_path)
            else:
                logger.info("创建空白演示文稿")
                self.prs = Presentation()
                
            # 分析模板布局
            self._analyze_template_layouts()
            
            # 初始化默认主题
            self._init_default_theme()
                
        except Exception as e:
            logger.error(f"初始化演示文稿失败: {str(e)}")
            logger.error(traceback.format_exc())
            raise
            
    def _analyze_template_layouts(self):
        """分析模板布局"""
        if not self.prs:
            return
            
        # 分析幻灯片布局
        for i, layout in enumerate(self.prs.slide_layouts):
            layout_name = layout.name if hasattr(layout, 'name') else f"Layout {i}"
            
            # 分析布局的占位符
            placeholders = []
            for ph in layout.placeholders:
                placeholders.append({
                    'idx': ph.placeholder_format.idx,
                    'type': ph.placeholder_format.type,
                    'name': ph.name
                })
                
            # 存储布局信息
            self.template_layouts[layout_name] = {
                'index': i,
                'placeholders': placeholders
            }
            
        logger.info(f"分析了 {len(self.template_layouts)} 个幻灯片布局")
            
    def _init_default_theme(self):
        """初始化默认主题"""
        self.theme = {
            'title_font_size': 36,
            'subtitle_font_size': 24,
            'content_font_size': 18,
            'bullet_font_size': 18,
            'table_font_size': 14,
            'theme_color': RGBColor(0, 112, 192),  # 蓝色
            'accent_color': RGBColor(0, 176, 80),  # 绿色
            'text_color': RGBColor(0, 0, 0),       # 黑色
            'background': None  # 默认背景
        }
        
    def register_component(self, name, component):
        """
        注册幻灯片组件
        
        Args:
            name: 组件名称
            component: 组件实例
        """
        if isinstance(component, SlideComponentBase):
            self.components[name] = component
            logger.info(f"注册组件: {name}")
        else:
            logger.error(f"无法注册组件: {name}，组件必须继承自SlideComponentBase")
            
    def set_theme(self, theme_config):
        """
        设置演示文稿主题
        
        Args:
            theme_config: 主题配置字典
        """
        self.theme.update(theme_config)
        logger.info(f"更新演示文稿主题: {theme_config.keys()}")
        
    def create_slide(self, slide_data, context=None):
        """
        创建单个幻灯片
        
        Args:
            slide_data: 幻灯片数据
            context: 上下文信息
            
        Returns:
            创建的幻灯片对象
        """
        try:
            # 合并上下文和主题
            ctx = dict(self.theme)
            if context:
                ctx.update(context)
                
            # 确定幻灯片布局
            layout_type = slide_data.get('layout', '').lower()
            slide_type = slide_data.get('type', '').lower()
            
            # 选择合适的布局
            slide_layout = self._get_slide_layout(slide_type, layout_type)
            
            # 添加幻灯片
            slide = self.prs.slides.add_slide(slide_layout)
            
            # 应用各组件
            for component_name, component in self.components.items():
                slide = component.apply(slide, slide_data, ctx)
                
            return slide
            
        except Exception as e:
            logger.error(f"创建幻灯片失败: {str(e)}")
            logger.error(traceback.format_exc())
            return None
            
    def _get_slide_layout(self, slide_type, layout_type):
        """
        获取合适的幻灯片布局
        
        Args:
            slide_type: 幻灯片类型
            layout_type: 布局类型
            
        Returns:
            幻灯片布局对象
        """
        # 默认布局索引
        default_layout_idx = 1  # 标题和内容布局
        
        # 根据幻灯片类型选择布局
        if slide_type == 'cover' or slide_type == 'title':
            return self.prs.slide_layouts[0]  # 标题幻灯片
        elif slide_type == 'section':
            return self.prs.slide_layouts[2]  # 章节标题
        elif slide_type == 'two_content':
            return self.prs.slide_layouts[3]  # 两栏内容
        elif slide_type == 'image' or layout_type == 'image':
            # 尝试找到适合图片的布局
            for name, layout in self.template_layouts.items():
                if 'picture' in name.lower() or 'image' in name.lower():
                    return self.prs.slide_layouts[layout['index']]
            # 默认使用标题和内容布局
            return self.prs.slide_layouts[1]
            
        # 默认布局
        return self.prs.slide_layouts[default_layout_idx]
            
    def create_ppt(self, slides_data, output_path, context=None):
        """
        创建PPT演示文稿
        
        Args:
            slides_data: 幻灯片数据列表
            output_path: 输出文件路径
            context: 上下文信息
            
        Returns:
            bool: 成功返回True，失败返回False
        """
        try:
            logger.info(f"开始创建PPT: {output_path}")
            logger.info(f"幻灯片数量: {len(slides_data)}")
            
            # 存储幻灯片数据
            self.slides_data = slides_data
            
            # 创建所有幻灯片
            for i, slide_data in enumerate(slides_data):
                logger.info(f"创建第 {i+1}/{len(slides_data)} 张幻灯片")
                self.create_slide(slide_data, context)
                
            # 保存演示文稿
            self.prs.save(output_path)
            logger.info(f"PPT已保存到: {output_path}")
            
            return True
            
        except Exception as e:
            logger.error(f"创建PPT失败: {str(e)}")
            logger.error(traceback.format_exc())
            return False
            
    def load_plugins(self, plugin_dir=None):
        """
        从指定目录加载插件
        
        Args:
            plugin_dir: 插件目录，如果为None则使用默认目录
        """
        if plugin_dir is None:
            plugin_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'ppt_plugins')
            
        if not os.path.exists(plugin_dir):
            logger.warning(f"插件目录不存在: {plugin_dir}")
            return
            
        logger.info(f"从目录加载插件: {plugin_dir}")
        
        # 将插件目录添加到Python路径
        sys.path.append(plugin_dir)
        
        # 查找所有Python文件
        for file_name in os.listdir(plugin_dir):
            if file_name.endswith('.py') and not file_name.startswith('__'):
                module_name = file_name[:-3]  # 去除.py扩展名
                try:
                    # 导入模块
                    module = importlib.import_module(module_name)
                    
                    # 查找所有继承自SlideComponentBase的类
                    for name, obj in inspect.getmembers(module):
                        if inspect.isclass(obj) and issubclass(obj, SlideComponentBase) and obj != SlideComponentBase:
                            # 实例化并注册组件
                            component = obj()
                            component_name = name.lower().replace('component', '')
                            self.register_component(component_name, component)
                            logger.info(f"从插件加载组件: {component_name}")
                            
                except Exception as e:
                    logger.error(f"加载插件失败: {module_name}, 错误: {str(e)}")
                    logger.error(traceback.format_exc())
                    
# 便捷函数
def generate_ppt(slides_data, output_path, template_path=None, theme=None, image_service=None):
    """
    生成PPT的便捷函数
    
    Args:
        slides_data: 幻灯片数据列表
        output_path: 输出文件路径
        template_path: 模板文件路径，如果不提供则创建空白演示文稿
        theme: 主题配置
        image_service: 图片服务对象
        
    Returns:
        bool: 成功返回True，失败返回False
    """
    try:
        logger.info("开始生成PPT，共 %d 张幻灯片", len(slides_data))
        
        # 创建PPT生成器
        generator = PPTGenerator(template_path)
        
        # 设置主题
        if theme:
            generator.set_theme(theme)
            
        # 加载插件
        generator.load_plugins()
        
        # 创建上下文
        context = {}
        # 确保图片服务传入上下文
        if image_service:
            logger.info("已提供图片服务，将用于生成图片")
            context['image_service'] = image_service
        else:
            logger.warning("未提供图片服务，将使用默认图片")
            
        # 生成PPT
        success = generator.create_ppt(slides_data, output_path, context)
        
        if success:
            logger.info("PPT生成成功: %s", output_path)
        else:
            logger.error("PPT生成失败")
        
        return success
        
    except Exception as e:
        logger.error(f"生成PPT失败: {str(e)}")
        logger.error(traceback.format_exc())
        return False

# 测试代码
if __name__ == "__main__":
    # 测试数据
    test_slides = [
        {
            "type": "cover",
            "title": "演示文稿标题",
            "content": "副标题内容"
        },
        {
            "type": "content",
            "title": "内容页标题",
            "content": "这是内容页的文本内容，可以包含各种信息。"
        },
        {
            "type": "bullet",
            "title": "项目符号页",
            "bullet_points": [
                "要点1",
                "要点2",
                "要点3",
                {"text": "子要点", "level": 1}
            ]
        }
    ]
    
    # 输出路径
    output_path = "test_output.pptx"
    
    # 生成PPT
    success = generate_ppt(test_slides, output_path)
    print(f"PPT生成{'成功' if success else '失败'}") 
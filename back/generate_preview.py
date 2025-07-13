#!/usr/bin/env python
"""
模板预览图生成工具
此脚本用于生成PPT模板的预览图，直接提取第一页并保存为PNG图片
支持多种方法来确保能够生成高质量的预览图
"""

import os
import sys
import logging
import traceback
import tempfile
import shutil
import subprocess
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from io import BytesIO

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
)
logger = logging.getLogger("preview_generator")

def generate_preview(template_path, output_path=None):
    """
    生成PPT模板的预览图
    
    Args:
        template_path: 模板文件路径
        output_path: 输出文件路径，如果为None则使用模板名称
        
    Returns:
        生成的预览图路径
    """
    logger.info(f"正在为模板生成预览图: {template_path}")
    
    # 确定输出路径
    if output_path is None:
        template_name = os.path.basename(template_path)
        base_name, _ = os.path.splitext(template_name)
        template_dir = os.path.dirname(template_path)
        preview_dir = os.path.join(template_dir, 'previews')
        os.makedirs(preview_dir, exist_ok=True)
        output_path = os.path.join(preview_dir, f"{base_name}_preview.png")
    
    # 创建临时目录
    temp_dir = tempfile.mkdtemp()
    temp_pptx = os.path.join(temp_dir, "temp.pptx")
    temp_png = os.path.join(temp_dir, "slide1.png")
    
    try:
        # 复制模板文件到临时位置
        shutil.copy2(template_path, temp_pptx)
        
        # 尝试使用不同方法生成预览图
        preview_generated = False
        
        # 方法1: 使用PowerPoint COM对象 (Windows专用)
        if not preview_generated and os.name == 'nt':
            try:
                import win32com.client
                import pythoncom
                
                logger.info("尝试使用PowerPoint COM对象生成预览图")
                
                # 初始化COM组件
                pythoncom.CoInitialize()
                
                # 使用PowerPoint COM对象导出第一页为图片
                powerpoint = win32com.client.Dispatch("PowerPoint.Application")
                powerpoint.Visible = False
                
                # 打开演示文稿
                presentation = powerpoint.Presentations.Open(os.path.abspath(temp_pptx))
                
                # 确保有幻灯片
                if presentation.Slides.Count > 0:
                    # 导出第一页为PNG
                    presentation.Slides(1).Export(os.path.abspath(temp_png), "PNG", 1024, 768)
                    logger.info(f"成功导出第一页幻灯片到: {temp_png}")
                    
                    # 检查文件是否生成
                    if os.path.exists(temp_png):
                        # 复制到输出路径
                        shutil.copy2(temp_png, output_path)
                        preview_generated = True
                
                # 关闭PowerPoint
                presentation.Close()
                powerpoint.Quit()
                
                # 释放COM组件
                pythoncom.CoUninitialize()
                
            except Exception as e:
                logger.warning(f"使用PowerPoint COM对象生成预览图失败: {str(e)}")
        
        # 方法2: 使用LibreOffice (跨平台)
        if not preview_generated:
            try:
                logger.info("尝试使用LibreOffice生成预览图")
                
                # 检查LibreOffice可执行文件
                libreoffice_paths = [
                    "libreoffice", "soffice",
                    r"C:\Program Files\LibreOffice\program\soffice.exe",
                    r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
                    "/usr/bin/libreoffice",
                    "/usr/bin/soffice",
                    "/Applications/LibreOffice.app/Contents/MacOS/soffice"
                ]
                
                libreoffice_path = None
                for path in libreoffice_paths:
                    try:
                        # 检查命令是否存在
                        if os.path.isfile(path) or shutil.which(path):
                            libreoffice_path = path
                            break
                    except:
                        continue
                
                if libreoffice_path:
                    # 导出为PNG
                    export_dir = os.path.join(temp_dir, "export")
                    os.makedirs(export_dir, exist_ok=True)
                    
                    # 构建命令
                    cmd = [
                        libreoffice_path,
                        "--headless",
                        "--convert-to", "png",
                        "--outdir", export_dir,
                        temp_pptx
                    ]
                    
                    # 执行命令
                    process = subprocess.run(cmd, capture_output=True, text=True)
                    
                    if process.returncode == 0:
                        # 查找生成的PNG文件
                        png_files = [f for f in os.listdir(export_dir) if f.endswith('.png')]
                        if png_files:
                            # 使用第一个PNG文件
                            png_path = os.path.join(export_dir, png_files[0])
                            
                            # 处理图像
                            with Image.open(png_path) as img:
                                # 调整大小为1024x768，保持纵横比
                                img.thumbnail((1024, 768), Image.Resampling.LANCZOS)
                                img.save(output_path)
                                preview_generated = True
                                logger.info(f"使用LibreOffice成功生成预览图: {output_path}")
            except Exception as e:
                logger.warning(f"使用LibreOffice生成预览图失败: {str(e)}")
        
        # 方法3: 使用python-pptx提取信息并创建预览图
        if not preview_generated:
            try:
                from pptx import Presentation
                
                logger.info("尝试使用python-pptx生成预览图")
                
                # 加载演示文稿
                prs = Presentation(template_path)
                
                # 确保有幻灯片
                if len(prs.slides) > 0:
                    slide = prs.slides[0]
                    
                    # 创建一个图像
                    img = Image.new('RGB', (1024, 768), color=(240, 240, 240))
                    draw = ImageDraw.Draw(img)
                    
                    # 尝试获取背景颜色
                    try:
                        if hasattr(slide.background, 'fill') and slide.background.fill.type:
                            if hasattr(slide.background.fill, 'solid'):
                                if hasattr(slide.background.fill.solid_fill, 'fore_color'):
                                    color = slide.background.fill.solid_fill.fore_color
                                    if hasattr(color, 'rgb'):
                                        r, g, b = color.rgb
                                        img = Image.new('RGB', (1024, 768), color=(r, g, b))
                                        draw = ImageDraw.Draw(img)
                    except Exception as e:
                        logger.warning(f"无法获取背景颜色: {str(e)}")
                    
                    # 尝试绘制幻灯片内容
                    try:
                        # 绘制形状
                        for shape in slide.shapes:
                            # 绘制文本框
                            if hasattr(shape, 'text_frame') and shape.text_frame.text:
                                text = shape.text_frame.text
                                x, y = shape.left / 914400 * 1024, shape.top / 914400 * 768
                                width, height = shape.width / 914400 * 1024, shape.height / 914400 * 768
                                
                                # 绘制文本框背景
                                draw.rectangle([x, y, x + width, y + height], outline=(200, 200, 200), width=1)
                                
                                # 尝试加载字体
                                try:
                                    # 尝试使用系统字体
                                    font_path = None
                                    for font_name in ["arial.ttf", "simhei.ttf", "simsun.ttc", "msyh.ttc"]:
                                        try:
                                            # 检查Windows系统字体目录
                                            windows_font_path = os.path.join("C:\\Windows\\Fonts", font_name)
                                            if os.path.exists(windows_font_path):
                                                font_path = windows_font_path
                                                break
                                        except:
                                            pass
                                    
                                    if font_path:
                                        font = ImageFont.truetype(font_path, 14)
                                    else:
                                        font = ImageFont.load_default()
                                    
                                    # 绘制文本
                                    draw.text((x + 5, y + 5), text[:100], fill=(0, 0, 0), font=font)
                                except Exception as font_error:
                                    logger.warning(f"字体加载失败: {str(font_error)}")
                                    draw.text((x + 5, y + 5), text[:100], fill=(0, 0, 0))
                            
                            # 绘制形状
                            elif hasattr(shape, 'shape_type'):
                                x, y = shape.left / 914400 * 1024, shape.top / 914400 * 768
                                width, height = shape.width / 914400 * 1024, shape.height / 914400 * 768
                                
                                # 绘制矩形
                                draw.rectangle([x, y, x + width, y + height], outline=(150, 150, 150), width=1)
                    except Exception as shape_error:
                        logger.warning(f"绘制形状失败: {str(shape_error)}")
                    
                    # 绘制模板名称
                    template_name_display = os.path.splitext(os.path.basename(template_path))[0]
                    try:
                        # 尝试使用系统字体
                        font_path = None
                        for font_name in ["arial.ttf", "simhei.ttf", "simsun.ttc", "msyh.ttc"]:
                            try:
                                # 检查Windows系统字体目录
                                windows_font_path = os.path.join("C:\\Windows\\Fonts", font_name)
                                if os.path.exists(windows_font_path):
                                    font_path = windows_font_path
                                    break
                            except:
                                pass
                        
                        if font_path:
                            font = ImageFont.truetype(font_path, 24)
                        else:
                            font = ImageFont.load_default()
                        
                        # 绘制模板名称
                        text_width = draw.textlength(template_name_display, font=font)
                        draw.text((512 - text_width/2, 700), template_name_display, fill=(0, 0, 0), font=font)
                    except Exception as font_error:
                        logger.warning(f"字体加载失败: {str(font_error)}")
                        draw.text((512, 700), template_name_display, fill=(0, 0, 0))
                    
                    # 保存预览图
                    img.save(output_path)
                    preview_generated = True
                    logger.info(f"使用python-pptx成功生成预览图: {output_path}")
            except Exception as e:
                logger.warning(f"使用python-pptx生成预览图失败: {str(e)}")
        
        # 如果所有方法都失败，创建一个基本预览图
        if not preview_generated:
            logger.warning("所有方法都失败，创建基本预览图")
            
            # 创建一个基本预览图
            img = Image.new('RGB', (1024, 768), color=(240, 240, 240))
            draw = ImageDraw.Draw(img)
            
            # 绘制边框
            draw.rectangle([10, 10, 1014, 758], outline=(200, 200, 200), width=2)
            
            # 获取模板名称
            template_name_display = os.path.splitext(os.path.basename(template_path))[0]
            
            # 尝试加载字体
            try:
                # 尝试使用系统字体
                font_path = None
                for font_name in ["arial.ttf", "simhei.ttf", "simsun.ttc", "msyh.ttc"]:
                    try:
                        # 检查Windows系统字体目录
                        windows_font_path = os.path.join("C:\\Windows\\Fonts", font_name)
                        if os.path.exists(windows_font_path):
                            font_path = windows_font_path
                            break
                    except:
                        pass
                
                if font_path:
                    font_title = ImageFont.truetype(font_path, 36)
                    font_subtitle = ImageFont.truetype(font_path, 24)
                else:
                    font_title = ImageFont.load_default()
                    font_subtitle = ImageFont.load_default()
                
                # 绘制标题
                title_width = draw.textlength(template_name_display, font=font_title)
                draw.text((512 - title_width/2, 300), template_name_display, fill=(0, 0, 0), font=font_title)
                
                # 绘制提示
                subtitle = "模板预览"
                subtitle_width = draw.textlength(subtitle, font=font_subtitle)
                draw.text((512 - subtitle_width/2, 400), subtitle, fill=(100, 100, 100), font=font_subtitle)
            except Exception as font_error:
                logger.warning(f"字体加载失败: {str(font_error)}")
                # 使用简单绘制
                draw.text((512, 300), template_name_display, fill=(0, 0, 0))
                draw.text((512, 400), "模板预览", fill=(100, 100, 100))
            
            # 保存预览图
            img.save(output_path)
            logger.info(f"创建基本预览图: {output_path}")
        
        return output_path
    except Exception as e:
        logger.error(f"生成预览图失败: {str(e)}")
        logger.error(traceback.format_exc())
        raise
    finally:
        # 清理临时文件
        try:
            shutil.rmtree(temp_dir, ignore_errors=True)
        except Exception as cleanup_error:
            logger.warning(f"清理临时文件失败: {str(cleanup_error)}")

def main():
    """命令行入口"""
    if len(sys.argv) < 2:
        print("用法: python generate_preview.py <template.pptx> [output.png]")
        sys.exit(1)
    
    template_path = sys.argv[1]
    output_path = sys.argv[2] if len(sys.argv) > 2 else None
    
    try:
        preview_path = generate_preview(template_path, output_path)
        print(f"预览图已生成: {preview_path}")
    except Exception as e:
        print(f"生成预览图失败: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main() 
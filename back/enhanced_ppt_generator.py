import os
import logging
import traceback
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from io import BytesIO

# 配置日志
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("ppt_generation.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("enhanced_ppt_generator")

class PPTGenerator:
    """增强型PPT生成器，提供更可靠的模板填充功能"""
    
    def __init__(self, template_path=None):
        """
        初始化PPT生成器
        
        Args:
            template_path: 模板文件路径，如果不提供则创建空白演示文稿
        """
        self.template_path = template_path
        self.prs = None
        self.slides_data = []
        self.template_layouts = {}  # 保存模板布局信息
        self.init_presentation()
    
    def init_presentation(self):
        """初始化演示文稿对象"""
        try:
            if self.template_path and os.path.exists(self.template_path):
                logger.info(f"加载PPT模板: {self.template_path}")
                self.prs = Presentation(self.template_path)
                logger.info(f"模板加载成功，包含 {len(self.prs.slides)} 张幻灯片和 {len(self.prs.slide_layouts)} 种布局")
                
                # 记录模板布局信息
                for i, layout in enumerate(self.prs.slide_layouts):
                    placeholders = [p for p in layout.placeholders]
                    logger.debug(f"布局 {i}: 包含 {len(placeholders)} 个占位符")
                    self.template_layouts[i] = {
                        "name": getattr(layout, "name", f"Layout {i}"),
                        "placeholders": len(placeholders)
                    }
                    
                # 删除模板中的示例幻灯片
                self._clear_template_slides()
            else:
                logger.info("创建空白演示文稿")
                self.prs = Presentation()
                logger.info(f"空白演示文稿创建成功，包含 {len(self.prs.slide_layouts)} 种布局")
        except Exception as e:
            logger.error(f"初始化演示文稿失败: {str(e)}")
            logger.error(traceback.format_exc())
            raise RuntimeError(f"初始化演示文稿失败: {str(e)}")
    
    def _clear_template_slides(self):
        """清除模板中的示例幻灯片"""
        try:
            while len(self.prs.slides) > 0:
                rId = self.prs.slides._sldIdLst[0].rId
                self.prs.part.drop_rel(rId)
                del self.prs.slides._sldIdLst[0]
            logger.info("已清除模板中的示例幻灯片")
        except Exception as e:
            logger.warning(f"清除模板幻灯片失败: {str(e)}")
    
    def add_slide(self, slide_data, layout_idx=None):
        """
        添加幻灯片
        
        Args:
            slide_data: 幻灯片数据
            layout_idx: 布局索引，如果不提供则自动选择
            
        Returns:
            添加的幻灯片对象
        """
        # 确定布局索引
        if layout_idx is None:
            layout_idx = self._choose_layout_for_slide(slide_data)
        
        # 确保布局索引有效
        if layout_idx >= len(self.prs.slide_layouts):
            logger.warning(f"布局索引 {layout_idx} 超出范围，使用默认布局 0")
            layout_idx = 0
        
        # 添加幻灯片
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_idx])
        logger.info(f"添加幻灯片，使用布局 {layout_idx}: {getattr(self.prs.slide_layouts[layout_idx], 'name', f'Layout {layout_idx}')}")
        
        # 填充幻灯片内容
        self._fill_slide_content(slide, slide_data)
        
        return slide
    
    def _choose_layout_for_slide(self, slide_data):
        """
        为幻灯片选择合适的布局
        
        Args:
            slide_data: 幻灯片数据
            
        Returns:
            布局索引
        """
        slide_type = slide_data.get('type', '').lower()
        layout = slide_data.get('layout', '').lower()
        
        # 类型/布局到索引的映射
        layout_map = {
            'cover': 0,        # 封面
            'title': 0,        # 标题
            'keypoints': 1,    # 要点
            'content': 1,      # 内容
            'compare': 2,      # 对比
            'image': 6,        # 图文
            'summary': 5,      # 总结
            'conclusion': 5,   # 结论
            'flow': 7,         # 流程
            'chart': 8,        # 图表
            'table': 3,        # 表格
            'quote': 4,        # 引用
        }
        
        # 首先尝试使用layout字段
        if layout in layout_map:
            layout_idx = layout_map[layout]
            if layout_idx < len(self.prs.slide_layouts):
                return layout_idx
        
        # 其次尝试使用type字段
        if slide_type in layout_map:
            layout_idx = layout_map[slide_type]
            if layout_idx < len(self.prs.slide_layouts):
                return layout_idx
        
        # 默认使用内容布局
        return min(1, len(self.prs.slide_layouts) - 1)
    
    def _fill_slide_content(self, slide, slide_data):
        """
        填充幻灯片内容
        
        Args:
            slide: 幻灯片对象
            slide_data: 幻灯片数据
        """
        logger.info(f"填充幻灯片内容: {slide_data.get('title', '无标题')}")
        
        # 填充标题
        title = slide_data.get('title', '')
        if title:
            self._set_slide_title(slide, title)
        
        # 填充内容
        content = slide_data.get('content', '')
        keypoints = slide_data.get('keypoints', [])
        
        if keypoints and isinstance(keypoints, list) and len(keypoints) > 0:
            # 优先使用要点
            self._set_slide_keypoints(slide, keypoints)
        elif content:
            # 其次使用内容
            self._set_slide_content(slide, content)
        
        # 添加图片
        image_url = slide_data.get('image')
        if image_url:
            self._add_image_to_slide(slide, image_url)
        
        # 添加表格
        table_data = slide_data.get('table')
        if table_data and isinstance(table_data, list):
            self._add_table_to_slide(slide, table_data)
    
    def _set_slide_title(self, slide, title):
        """设置幻灯片标题"""
        # 查找标题占位符
        title_shape = None
        
        # 尝试使用内置标题
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            title_shape = slide.shapes.title
        
        # 如果没有内置标题，查找标题占位符
        if not title_shape:
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 1:  # 1=标题占位符
                    title_shape = shape
                    break
        
        # 如果仍然没有找到，创建文本框作为标题
        if not title_shape:
            left = Inches(0.5)
            top = Inches(0.5)
            width = Inches(9)
            height = Inches(1.2)
            title_shape = slide.shapes.add_textbox(left, top, width, height)
        
        # 设置标题文本
        title_shape.text = title
        
        # 应用标题格式
        for paragraph in title_shape.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(36)
                run.font.bold = True
                run.font.name = "微软雅黑"
    
    def _set_slide_content(self, slide, content):
        """设置幻灯片内容"""
        # 查找内容占位符
        content_shape = None
        
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:  # 2=内容占位符
                content_shape = shape
                break
        
        # 如果没有找到内容占位符，创建文本框
        if not content_shape:
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4.5)
            content_shape = slide.shapes.add_textbox(left, top, width, height)
        
        # 设置内容
        content_shape.text = content
        
        # 应用内容格式
        for paragraph in content_shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(24)
                run.font.name = "微软雅黑"
    
    def _set_slide_keypoints(self, slide, keypoints):
        """设置幻灯片要点"""
        # 查找内容占位符
        content_shape = None
        
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 2:  # 2=内容占位符
                content_shape = shape
                break
        
        # 如果没有找到内容占位符，创建文本框
        if not content_shape:
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(4.5)
            content_shape = slide.shapes.add_textbox(left, top, width, height)
        
        # 清除现有内容
        content_shape.text = ""
        
        # 添加要点
        text_frame = content_shape.text_frame
        
        for i, point in enumerate(keypoints):
            # 添加段落
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            # 处理不同类型的要点
            if isinstance(point, list):
                if len(point) == 2:  # 键值对 [key, value]
                    p.text = f"{point[0]}: {point[1]}"
                else:
                    p.text = " - ".join([str(item) for item in point])
            else:
                p.text = str(point)
            
            # 设置段落格式
            p.level = 0
            try:
                p.bullet = True
            except:
                logger.warning("不支持项目符号，跳过")
            
            # 设置字体
            for run in p.runs:
                run.font.size = Pt(24)
                run.font.name = "微软雅黑"
    
    def _add_image_to_slide(self, slide, image_url):
        """添加图片到幻灯片"""
        logger.info(f"添加图片: {image_url}")
        
        try:
            # 获取图片数据
            img_data = None
            
            if image_url.startswith(('http://', 'https://')):
                response = requests.get(image_url, timeout=10)
                if response.status_code == 200:
                    img_data = BytesIO(response.content)
                else:
                    logger.warning(f"获取图片失败，状态码: {response.status_code}")
                    return
            elif image_url.startswith('file://'):
                file_path = image_url[7:]
                if os.path.exists(file_path):
                    with open(file_path, 'rb') as f:
                        img_data = BytesIO(f.read())
                else:
                    logger.warning(f"本地图片不存在: {file_path}")
                    return
            
            if not img_data:
                logger.warning("无法获取图片数据")
                return
            
            # 查找图片占位符
            image_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 18:  # 18=图片占位符
                    image_placeholder = shape
                    break
            
            # 如果找到图片占位符，使用它
            if image_placeholder:
                left = image_placeholder.left
                top = image_placeholder.top
                width = image_placeholder.width
                height = image_placeholder.height
                
                # 删除占位符
                sp = image_placeholder._element
                sp.getparent().remove(sp)
                
                # 添加图片
                slide.shapes.add_picture(img_data, left, top, width, height)
                logger.info("图片添加到占位符位置")
            else:
                # 没有占位符，在幻灯片右侧添加图片
                slide_width = slide.slide_width
                slide_height = slide.slide_height
                
                # 计算图片位置
                left = slide_width * 0.6
                top = slide_height * 0.25
                width = slide_width * 0.35
                height = slide_height * 0.5
                
                slide.shapes.add_picture(img_data, left, top, width, height)
                logger.info("图片添加到幻灯片右侧")
        except Exception as e:
            logger.error(f"添加图片失败: {str(e)}")
            logger.error(traceback.format_exc())
    
    def _add_table_to_slide(self, slide, table_data):
        """添加表格到幻灯片"""
        logger.info("添加表格到幻灯片")
        
        try:
            # 如果表格数据是一维数组，转换为二维数组
            if not isinstance(table_data[0], list):
                cols = 3  # 默认3列
                rows = (len(table_data) + cols - 1) // cols
                
                new_table_data = []
                for i in range(rows):
                    row_data = table_data[i*cols:(i+1)*cols]
                    while len(row_data) < cols:
                        row_data.append("")
                    new_table_data.append(row_data)
                
                table_data = new_table_data
            
            # 获取行数和列数
            rows = len(table_data)
            cols = len(table_data[0]) if table_data[0] else 0
            
            if rows == 0 or cols == 0:
                logger.warning("表格数据为空")
                return
            
            # 查找表格占位符
            table_placeholder = None
            for shape in slide.placeholders:
                if shape.placeholder_format.type == 8:  # 8=表格占位符
                    table_placeholder = shape
                    break
            
            # 创建表格
            if table_placeholder:
                # 使用占位符位置
                left = table_placeholder.left
                top = table_placeholder.top
                width = table_placeholder.width
                height = table_placeholder.height
                
                # 删除占位符
                sp = table_placeholder._element
                sp.getparent().remove(sp)
            else:
                # 默认位置
                left = Inches(1)
                top = Inches(2.5)
                width = Inches(8)
                height = rows * Inches(0.5)  # 每行高度
            
            # 创建表格
            table = slide.shapes.add_table(rows, cols, left, top, width, height).table
            
            # 填充表格内容
            for i in range(rows):
                for j in range(cols):
                    cell = table.cell(i, j)
                    
                    # 设置单元格内容
                    if i < len(table_data) and j < len(table_data[i]):
                        cell_data = table_data[i][j]
                        
                        # 处理复杂的单元格数据（如字典）
                        if isinstance(cell_data, dict):
                            for key in ['value', 'text', 'content', 'name', 'title']:
                                if key in cell_data:
                                    cell_data = cell_data[key]
                                    break
                            else:
                                if cell_data:
                                    cell_data = next(iter(cell_data.values()))
                        
                        cell.text = str(cell_data)
                    else:
                        cell.text = ""
                    
                    # 设置表头格式
                    if i == 0:
                        cell.text_frame.paragraphs[0].font.bold = True
            
            logger.info(f"表格添加成功: {rows}行 {cols}列")
        except Exception as e:
            logger.error(f"添加表格失败: {str(e)}")
            logger.error(traceback.format_exc())
    
    def generate_slides_from_data(self, slides_data):
        """
        根据幻灯片数据生成幻灯片
        
        Args:
            slides_data: 幻灯片数据列表
        """
        logger.info(f"生成 {len(slides_data)} 张幻灯片")
        self.slides_data = slides_data
        
        # 添加幻灯片
        for i, slide_data in enumerate(slides_data):
            logger.info(f"处理第 {i+1} 张幻灯片，类型: {slide_data.get('type', '未指定')}，标题: {slide_data.get('title', '无标题')}")
            
            # 为第一页和最后一页特殊处理
            if i == 0:
                # 第一页通常是封面
                layout_idx = min(0, len(self.prs.slide_layouts) - 1)  # 封面布局
                slide_data['type'] = slide_data.get('type', 'cover')
                self.add_slide(slide_data, layout_idx)
            elif i == len(slides_data) - 1:
                # 最后一页通常是总结
                layout_idx = min(5, len(self.prs.slide_layouts) - 1)  # 总结布局
                slide_data['type'] = slide_data.get('type', 'summary')
                self.add_slide(slide_data, layout_idx)
            else:
                # 中间页面按类型选择布局
                self.add_slide(slide_data)
    
    def save(self, output_path):
        """
        保存演示文稿
        
        Args:
            output_path: 保存路径
        """
        logger.info(f"保存PPT到: {output_path}")
        
        try:
            # 确保输出目录存在
            output_dir = os.path.dirname(output_path)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir, exist_ok=True)
            
            # 保存PPT
            self.prs.save(output_path)
            logger.info("PPT保存成功")
            return True
        except Exception as e:
            logger.error(f"保存PPT失败: {str(e)}")
            logger.error(traceback.format_exc())
            return False

if __name__ == "__main__":
    # 测试代码
    test_data = [
        {"title": "演示文稿标题", "content": "演示文稿副标题", "type": "cover"},
        {"title": "内容页面", "content": "这是一个内容页面", "keypoints": ["要点1", "要点2", "要点3"]},
        {"title": "表格页面", "table": [["姓名", "年龄", "职业"], ["张三", "25", "工程师"], ["李四", "30", "设计师"]]},
        {"title": "总结页面", "content": "总结内容", "type": "summary"}
    ]
    
    try:
        # 测试使用模板
        template_path = "ppt_templates/绿色圆点.pptx"
        if os.path.exists(template_path):
            generator = PPTGenerator(template_path)
            generator.generate_slides_from_data(test_data)
            generator.save("test_template_output.pptx")
            print("使用模板生成PPT成功")
        
        # 测试不使用模板
        generator = PPTGenerator()
        generator.generate_slides_from_data(test_data)
        generator.save("test_no_template_output.pptx")
        print("不使用模板生成PPT成功")
    except Exception as e:
        print(f"测试失败: {str(e)}")
        traceback.print_exc() 
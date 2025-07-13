#!/usr/bin/env python
"""
测试PPT生成功能
"""

import os
import sys
import json
import time
import logging
import traceback
from datetime import datetime
from pptx import Presentation

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("ppt_test.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("ppt_test")

# 测试数据
TEST_OUTLINE = [
    {
        "type": "cover",
        "title": "测试演示文稿",
        "content": "这是一个测试演示文稿",
        "layout": "cover"
    },
    {
        "type": "content",
        "title": "内容页",
        "content": "这是一个内容页，用于测试文本内容的显示效果。这是一个内容页，用于测试文本内容的显示效果。这是一个内容页，用于测试文本内容的显示效果。",
        "layout": "content"
    },
    {
        "type": "keypoints",
        "title": "要点页",
        "content": "这是一个要点页",
        "keypoints": ["要点1", "要点2", "要点3"],
        "layout": "keypoints"
    },
    {
        "type": "image",
        "title": "图片页",
        "content": "这是一个图片页",
        "image": "https://example.com/image.jpg",  # 这个URL不存在，会触发错误处理
        "layout": "image"
    },
    {
        "type": "table",
        "title": "表格页",
        "content": "这是一个表格页",
        "table": [
            ["表头1", "表头2", "表头3"],
            ["数据1", "数据2", "数据3"],
            ["数据4", "数据5", "数据6"]
        ],
        "layout": "table"
    },
    {
        "type": "summary",
        "title": "总结页",
        "content": "这是一个总结页",
        "keypoints": ["总结1", "总结2", "总结3"],
        "layout": "summary"
    }
]

def test_ppt_without_template():
    """测试无模板PPT生成"""
    logger.info("=== 测试无模板PPT生成 ===")
    
    try:
        # 导入无模板PPT生成模块
        from ppt_without_template import generate_ppt_without_template
        
        # 生成测试文件名
        timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
        output_path = f"test_no_template_{timestamp}.pptx"
        
        # 生成PPT
        logger.info(f"开始生成无模板PPT: {output_path}")
        success = generate_ppt_without_template(TEST_OUTLINE, output_path)
        
        if success and os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            logger.info(f"无模板PPT生成成功: {output_path}, 大小: {file_size} 字节")
            
            # 尝试打开文件验证
            try:
                prs = Presentation(output_path)
                slide_count = len(prs.slides)
                logger.info(f"PPT幻灯片数量: {slide_count}")
                
                if slide_count == len(TEST_OUTLINE):
                    logger.info("幻灯片数量正确")
                else:
                    logger.warning(f"幻灯片数量不匹配: 期望 {len(TEST_OUTLINE)}, 实际 {slide_count}")
                
                return True
            except Exception as e:
                logger.error(f"PPT文件验证失败: {str(e)}")
                return False
        else:
            logger.error(f"无模板PPT生成失败")
            return False
    except Exception as e:
        logger.error(f"测试无模板PPT生成时发生异常: {str(e)}")
        logger.error(traceback.format_exc())
        return False

def test_ppt_with_template():
    """测试模板PPT生成"""
    logger.info("=== 测试模板PPT生成 ===")
    
    try:
        # 导入模板PPT生成模块
        from ppt_fill_template import fill_ppt_template
        
        # 查找模板文件
        template_dir = "ppt_templates"
        template_files = [f for f in os.listdir(template_dir) if f.endswith('.pptx') and not f.startswith('~$')]
        
        if not template_files:
            logger.warning("未找到模板文件，跳过模板测试")
            return False
        
        template_path = os.path.join(template_dir, template_files[0])
        logger.info(f"使用模板: {template_path}")
        
        # 查找元数据文件
        template_name_base = os.path.splitext(template_files[0])[0]
        metadata_path = os.path.join(template_dir, f"{template_name_base}.json")
        
        # 生成测试文件名
        timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
        output_path = f"test_with_template_{timestamp}.pptx"
        
        # 生成PPT
        logger.info(f"开始生成模板PPT: {output_path}")
        success = fill_ppt_template(template_path, TEST_OUTLINE, output_path, metadata_path)
        
        if success and os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            logger.info(f"模板PPT生成成功: {output_path}, 大小: {file_size} 字节")
            
            # 尝试打开文件验证
            try:
                prs = Presentation(output_path)
                slide_count = len(prs.slides)
                logger.info(f"PPT幻灯片数量: {slide_count}")
                
                return True
            except Exception as e:
                logger.error(f"PPT文件验证失败: {str(e)}")
                return False
        else:
            logger.error(f"模板PPT生成失败")
            return False
    except Exception as e:
        logger.error(f"测试模板PPT生成时发生异常: {str(e)}")
        logger.error(traceback.format_exc())
        return False

def test_app_ppt_generation():
    """测试通过app.py生成PPT"""
    logger.info("=== 测试通过app.py生成PPT ===")
    
    try:
        # 导入app模块
        sys.path.append(os.path.dirname(os.path.abspath(__file__)))
        from app import preprocess_outline_data, UPLOAD_FOLDER
        
        # 预处理大纲数据
        processed_outline = preprocess_outline_data(TEST_OUTLINE)
        
        # 生成测试文件名
        timestamp = datetime.now().strftime('%Y%m%d%H%M%S')
        filename = f"test_app_{timestamp}.pptx"
        output_path = os.path.join(UPLOAD_FOLDER, filename)
        
        # 导入无模板PPT生成模块
        from ppt_without_template import generate_ppt_without_template
        
        # 生成PPT
        logger.info(f"开始通过app.py生成PPT: {output_path}")
        success = generate_ppt_without_template(processed_outline, output_path)
        
        if success and os.path.exists(output_path):
            file_size = os.path.getsize(output_path)
            logger.info(f"通过app.py生成PPT成功: {output_path}, 大小: {file_size} 字节")
            
            # 尝试打开文件验证
            try:
                prs = Presentation(output_path)
                slide_count = len(prs.slides)
                logger.info(f"PPT幻灯片数量: {slide_count}")
                
                return True
            except Exception as e:
                logger.error(f"PPT文件验证失败: {str(e)}")
                return False
        else:
            logger.error(f"通过app.py生成PPT失败")
            return False
    except Exception as e:
        logger.error(f"测试通过app.py生成PPT时发生异常: {str(e)}")
        logger.error(traceback.format_exc())
        return False

def main():
    """主函数"""
    logger.info("开始PPT生成测试")
    
    # 测试无模板PPT生成
    no_template_result = test_ppt_without_template()
    logger.info(f"无模板PPT生成测试结果: {'成功' if no_template_result else '失败'}")
    
    # 测试模板PPT生成
    template_result = test_ppt_with_template()
    logger.info(f"模板PPT生成测试结果: {'成功' if template_result else '失败'}")
    
    # 测试通过app.py生成PPT
    app_result = test_app_ppt_generation()
    logger.info(f"通过app.py生成PPT测试结果: {'成功' if app_result else '失败'}")
    
    # 总结测试结果
    if no_template_result and template_result and app_result:
        logger.info("所有测试通过，PPT生成系统工作正常")
        return 0
    else:
        logger.error("部分测试失败，PPT生成系统可能存在问题")
        return 1

if __name__ == "__main__":
    sys.exit(main()) 